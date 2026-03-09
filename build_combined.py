"""Build editable PPTX: Claude AI text + NotebookLM matched images per slide.

Pipeline:
1. Claude AI generates slide content per topic (parallel, with web research)
2. NotebookLM generates visual slides per topic using Claude's content as source
3. Extract images from NotebookLM PDFs, skip cover + summary pages
4. Pair images 1:1 with Claude content slides (topic-level matching)
5. Build PPTX: editable text + matched image per slide + AI shape diagrams

If NotebookLM fails for a topic, that topic gets full-width text + shape diagrams only.
"""
import os, json, asyncio, tempfile, shutil, logging, time
from pathlib import Path

logging.basicConfig(level=logging.INFO, format='%(asctime)s %(message)s')
log = logging.getLogger(__name__)

CONTEXT_PATH = Path.home() / 'AppData' / 'Local' / 'Temp' / 'context.json'
NBLM_CACHE_DIR = Path.home() / 'AppData' / 'Local' / 'Temp' / 'nblm_slides_cache'
SLIDES_PER_TOPIC = 7
NBLM_TIMEOUT_PER_TOPIC = 300  # 5 minutes per topic


# ---------------------------------------------------------------------------
# NotebookLM: generate visual slides for one topic
# ---------------------------------------------------------------------------

async def generate_nblm_topic(topic_text: str, topic_title: str, topic_key: str,
                              cache_dir: Path) -> list:
    """Generate NotebookLM slides for one topic and return list of image paths.

    Args:
        topic_text: The Claude-generated slide content formatted as text.
        topic_title: Human-readable topic title.
        topic_key: e.g. "LU1_T1" for caching.
        cache_dir: Directory to store extracted images.

    Returns:
        List of image file paths (cover + summary already filtered out).
        Empty list if generation fails.
    """
    from notebooklm import NotebookLMClient

    topic_cache = cache_dir / topic_key
    topic_cache.mkdir(parents=True, exist_ok=True)

    # Check cache first
    cached_images = sorted(topic_cache.glob('slide_*.png'))
    if cached_images:
        print(f'  [{topic_key}] Using {len(cached_images)} cached images', flush=True)
        return [str(p) for p in cached_images]

    try:
        client = await NotebookLMClient.from_storage()
        async with client:
            # 1. Create notebook
            nb_title = f"Slides: {topic_title[:60]}"
            notebook = await client.notebooks.create(nb_title)
            nb_id = notebook.id
            print(f'  [{topic_key}] Created notebook: {nb_id[:12]}...', flush=True)

            # 2. Add Claude's content as the source
            source = await client.sources.add_text(
                nb_id,
                f"Training Content: {topic_title}",
                topic_text[:50000]
            )
            src_id = source.id

            # 3. Wait for source processing
            try:
                await client.sources.wait_for_sources(nb_id, [src_id], timeout=30.0)
            except Exception as e:
                print(f'  [{topic_key}] Source wait warning: {e}', flush=True)
                await asyncio.sleep(5)

            # 4. Generate slides with instructions to match content
            instructions = (
                f"Create a visual presentation for: {topic_title}. "
                f"Each slide should have relevant diagrams, flowcharts, icons, or illustrations. "
                f"Follow the exact content structure provided in the source material. "
                f"Make each slide visually rich with professional graphics. "
                f"Do NOT add a cover slide or summary slide. "
                f"Focus on one concept per slide with supporting visuals."
            )

            gen_result = await client.artifacts.generate_slide_deck(
                nb_id,
                source_ids=[src_id],
                instructions=instructions[:3000],
            )
            task_id = gen_result.task_id if hasattr(gen_result, 'task_id') else str(gen_result)
            print(f'  [{topic_key}] Generation started, polling...', flush=True)

            # 5. Poll for completion
            completed = False
            start_time = time.time()
            while time.time() - start_time < NBLM_TIMEOUT_PER_TOPIC:
                try:
                    decks = await client.artifacts.list_slide_decks(nb_id)
                    for deck in decks:
                        status = getattr(deck, 'status', None)
                        if getattr(deck, 'is_completed', False) or status == 3 or str(status) == '3':
                            completed = True
                            break
                        if getattr(deck, 'is_failed', False) or status == 4 or str(status) == '4':
                            print(f'  [{topic_key}] Generation FAILED', flush=True)
                            return []
                    if completed:
                        break
                except Exception:
                    pass
                await asyncio.sleep(3)
                elapsed = int(time.time() - start_time)
                if elapsed % 30 == 0 and elapsed > 0:
                    print(f'  [{topic_key}] Still generating... ({elapsed}s)', flush=True)

            if not completed:
                print(f'  [{topic_key}] TIMEOUT after {NBLM_TIMEOUT_PER_TOPIC}s', flush=True)
                return []

            # 6. Download as PDF
            pdf_path = str(topic_cache / f'{topic_key}.pdf')
            await client.artifacts.download_slide_deck(nb_id, pdf_path)
            print(f'  [{topic_key}] Downloaded slides', flush=True)

            # 7. Extract images from PDF
            images = _extract_pdf_images(pdf_path, topic_cache, topic_key)
            print(f'  [{topic_key}] Extracted {len(images)} images (cover/summary filtered)', flush=True)
            return images

    except Exception as e:
        print(f'  [{topic_key}] NotebookLM error: {e}', flush=True)
        return []


def _extract_pdf_images(pdf_path: str, output_dir: Path, prefix: str) -> list:
    """Extract page images from PDF, skip first (cover) and last (summary) pages.

    Returns list of image file paths.
    """
    import fitz

    doc = fitz.open(pdf_path)
    total_pages = len(doc)
    if total_pages == 0:
        return []

    # Skip first page (cover) and last page (summary/outro)
    # If only 1-2 pages, keep all (no cover/summary to skip)
    if total_pages <= 2:
        start_page = 0
        end_page = total_pages
    else:
        start_page = 1           # Skip cover (page 0)
        end_page = total_pages - 1  # Skip summary (last page)

    images = []
    for page_idx in range(start_page, end_page):
        page = doc[page_idx]
        # Render at 2x resolution for high quality
        mat = fitz.Matrix(2.0, 2.0)
        pix = page.get_pixmap(matrix=mat)
        img_path = str(output_dir / f'slide_{page_idx:03d}.png')
        pix.save(img_path)
        images.append(img_path)

    doc.close()
    return images


# ---------------------------------------------------------------------------
# Format Claude content as NotebookLM source text
# ---------------------------------------------------------------------------

def format_topic_as_source(topic_data: dict, course_title: str, lu_title: str) -> str:
    """Format Claude-generated topic slides as a text document for NotebookLM.

    This ensures NotebookLM generates images matching our exact slide content.
    """
    title = topic_data.get('title', 'Topic')
    slides = topic_data.get('slides', [])
    activity = topic_data.get('activity', [])

    lines = [
        f"Course: {course_title}",
        f"Learning Unit: {lu_title}",
        f"Topic: {title}",
        "",
        "=" * 60,
        ""
    ]

    for i, slide in enumerate(slides):
        s_title = slide.get('title', '')
        s_bullets = slide.get('bullets', [])
        lines.append(f"--- Slide {i + 1}: {s_title} ---")
        lines.append("")
        for bullet in s_bullets:
            lines.append(f"  - {bullet}")
        lines.append("")

    if activity:
        lines.append("--- Activity ---")
        for step in activity:
            lines.append(f"  {step}")
        lines.append("")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Main build pipeline
# ---------------------------------------------------------------------------

async def build():
    from courseware_agents.slides.slides_agent import generate_slide_content
    from generate_slides.build_pptx import build_lu_deck

    ctx = json.loads(CONTEXT_PATH.read_text())
    lus = ctx.get('Learning_Units', [])
    course_title = ctx.get('Course_Title', 'Course')

    print(f'Course: {course_title}', flush=True)
    print(f'Learning Units: {len(lus)}', flush=True)

    # Prepare cache directory
    NBLM_CACHE_DIR.mkdir(parents=True, exist_ok=True)

    all_pptx = []
    for lu_idx, lu in enumerate(lus):
        lu_num = lu.get('LU_Number') or f'LU{lu_idx+1}'
        lu_title = lu.get('LU_Title', 'Learning Unit')
        topics = lu.get('Topics', [])
        is_first = (lu_idx == 0)
        is_last = (lu_idx == len(lus) - 1)

        # ------------------------------------------------------------------
        # STEP 1: Claude AI generates slide content (all topics in parallel)
        # ------------------------------------------------------------------
        print(f'\n{"="*60}', flush=True)
        print(f'[{lu_num}] STEP 1: Claude AI generating content for {len(topics)} topics...', flush=True)
        slides_data = await generate_slide_content(
            ctx, lu, lu_idx, len(lus), slides_per_topic=SLIDES_PER_TOPIC
        )
        topic_list = slides_data.get("topics", [])
        total_content_slides = sum(len(t.get('slides', [])) for t in topic_list)
        print(f'[{lu_num}] Claude generated {len(topic_list)} topics, {total_content_slides} content slides', flush=True)

        # ------------------------------------------------------------------
        # STEP 2: NotebookLM generates visual slides per topic (sequential)
        # ------------------------------------------------------------------
        print(f'\n[{lu_num}] STEP 2: NotebookLM generating images per topic...', flush=True)
        images_per_topic = {}  # {topic_idx: [img_path, ...]}

        for ti, topic_data in enumerate(topic_list):
            topic_key = f'{lu_num}_T{ti+1}'
            t_title = topic_data.get('title', f'Topic {ti+1}')
            n_slides = len(topic_data.get('slides', []))

            # Format Claude content as source for NotebookLM
            source_text = format_topic_as_source(topic_data, course_title, lu_title)

            print(f'\n  [{topic_key}] "{t_title}" ({n_slides} slides)', flush=True)

            # Generate NotebookLM slides for this topic
            nblm_images = await generate_nblm_topic(
                source_text, t_title, topic_key, NBLM_CACHE_DIR
            )

            if nblm_images:
                # Pair 1:1: take min(nblm_images, n_content_slides) images
                paired = nblm_images[:n_slides]
                images_per_topic[ti] = paired
                print(f'  [{topic_key}] Paired {len(paired)}/{n_slides} slides with images', flush=True)
            else:
                print(f'  [{topic_key}] No NotebookLM images — using shape diagrams only', flush=True)

            # Delay between topics to avoid rate limiting
            if ti < len(topic_list) - 1:
                await asyncio.sleep(3)

        # ------------------------------------------------------------------
        # STEP 3: Build PPTX with text + matched images + shape diagrams
        # ------------------------------------------------------------------
        print(f'\n[{lu_num}] STEP 3: Building PPTX...', flush=True)
        n_topics_with_images = sum(1 for v in images_per_topic.values() if v)
        n_total_images = sum(len(v) for v in images_per_topic.values())
        print(f'  {n_topics_with_images}/{len(topic_list)} topics have NotebookLM images ({n_total_images} total)', flush=True)

        pptx_path, slide_count = build_lu_deck(
            ctx, lu_idx, slides_data,
            is_first=is_first, is_last=is_last,
            images_per_topic=images_per_topic if images_per_topic else None,
        )
        print(f'[{lu_num}] Built {slide_count} total slides', flush=True)
        all_pptx.append(pptx_path)

    # ------------------------------------------------------------------
    # STEP 4: Merge LUs and save
    # ------------------------------------------------------------------
    if len(all_pptx) > 1:
        from generate_slides.slides_generation import _merge_pptx_to_single
        merged = tempfile.mktemp(suffix='_ALL.pptx')
        merged, total = _merge_pptx_to_single(all_pptx, merged)
    else:
        from pptx import Presentation
        merged = all_pptx[0]
        total = len(Presentation(merged).slides)

    safe_title = course_title.replace(':', '').replace('/', '-').replace(' ', '_')[:40]
    dl_path = str(Path.home() / 'Downloads' / f'{safe_title}_Slides.pptx')
    shutil.copy2(merged, dl_path)
    print(f'\n{"="*60}', flush=True)
    print(f'DONE! {total} slides -> {dl_path}', flush=True)
    print(f'{"="*60}', flush=True)


if __name__ == '__main__':
    asyncio.run(build())

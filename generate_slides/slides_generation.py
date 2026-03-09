"""
Slides Generation Module

Generates presentation slides from extracted course info using NotebookLM.
No file upload needed — uses course context from Extract Course Info page.

Flow:
1. User extracts course info on the Extract Course Info page
2. App formats the structured context as source text
3. App calls NotebookLM directly:
   - create notebook
   - add course content as source text
   - research internet for key topics (optional)
   - import research sources into notebook
   - generate slide deck with all sources
4. User gets slides in NotebookLM Studio

Dependencies:
    - streamlit
    - notebooklm-py[browser] (pip install notebooklm-py[browser])
"""

import streamlit as st
import asyncio
import logging
import os
import re
import shutil
import time
import urllib.parse
from typing import Optional, Dict, Any, List

logger = logging.getLogger(__name__)


def _auto_match_company(org_name: str) -> dict:
    """Match organisation name from CP against company database.

    Returns the matching company dict, or empty dict if no match.
    """
    if not org_name:
        return {}
    try:
        from generate_ap_fg_lg.utils.organizations import get_organizations
        organizations = get_organizations()
        if not organizations:
            return {}
        def _norm(s: str) -> str:
            return s.lower().strip().replace(".", "").replace(",", "")
        target = _norm(org_name)
        for org in organizations:
            if _norm(org.get("name", "")) == target:
                return org
        # Fallback: substring match
        for org in organizations:
            org_norm = _norm(org.get("name", ""))
            if target in org_norm or org_norm in target:
                return org
    except Exception:
        pass
    return {}


# =============================================================================
# Helpers for building lesson plan and activity source text
# =============================================================================

def _build_lesson_plan_text(context: dict) -> list:
    """Build Lesson Plan slide content from course context.

    Creates a day-by-day schedule matching supervisor's reference PPTX:
    - Digital Attendance (AM) in each day
    - Topics distributed across days
    - Lunch Break
    - Digital Attendance (PM)
    - Assessment on last day
    """
    lines = []
    lus = context.get('Learning_Units', [])
    total_hours = context.get('Total_Course_Duration_Hours', '')
    try:
        hours = float(str(total_hours).replace('hours', '').replace('hrs', '').replace('h', '').strip() or '0')
    except (ValueError, TypeError):
        hours = 8
    days = max(1, round(hours / 8))

    # Gather all topics in order with global numbering
    all_topics = []
    global_topic_num = 0
    for lu_idx, lu in enumerate(lus):
        for t_idx, topic in enumerate(lu.get('Topics', [])):
            global_topic_num += 1
            topic_title = topic.get('Topic_Title', f'Topic {global_topic_num}')
            all_topics.append((global_topic_num, topic_title))

    # Distribute topics across days (excluding assessment time on last day)
    topics_per_day = max(1, len(all_topics) // days)
    topic_idx = 0

    for day in range(1, days + 1):
        lines.append(f"### Day {day}")

        lines.append("- **Digital Attendance (AM)**")
        if day == 1:
            lines.append("- Trainer and Learners Introduction")
            lines.append("- Learning Outcomes")
            lines.append("- Course Outline")

        # Assign topics to this day
        day_topics = []
        end_idx = min(topic_idx + topics_per_day, len(all_topics))
        if day == days:
            end_idx = len(all_topics)
        for ti in range(topic_idx, end_idx):
            day_topics.append(all_topics[ti])
        topic_idx = end_idx

        # First half topics (before lunch)
        mid = max(1, len(day_topics) // 2)
        for num, title in day_topics[:mid]:
            lines.append(f"- Topic {num}: {title}")

        lines.append("- Lunch Break")
        lines.append("- **Digital Attendance (PM)**")

        # Second half topics (after lunch)
        for num, title in day_topics[mid:]:
            lines.append(f"- Topic {num}: {title}")

        if day == days:
            lines.append("- **Course Feedback and TRAQOM Survey**")
            lines.append("- **Digital Attendance (Assessment)**")
            lines.append("- Final Assessment & End of Class")

        lines.append(f"- End of Day {day}")
        lines.append("")

    return lines


def _build_activity_text(topic_title: str, topic_bullets: list,
                          assessment_methods: list) -> list:
    """Build activity/lab slide content for a topic.

    Following the supervisor's reference PPTX pattern:
    - "Activity: [Topic Title]" or "Lab - [Topic Title]"
    - Hands-on activity instructions with clear steps
    - Link to assessment method from CP
    """
    lines = []
    lines.append(f"### Activity: {topic_title}")
    lines.append("")
    lines.append(f"**Objective:** Apply and demonstrate understanding of {topic_title} concepts.")
    lines.append("")

    # Generate specific activity steps based on topic content
    lines.append("**Instructions:**")
    if topic_bullets and len(topic_bullets) >= 2:
        lines.append(f"1. Review the key concepts covered: {', '.join(topic_bullets[:3])}")
        lines.append(f"2. In groups of 3-4, analyze the scenario/case study provided by the trainer")
        lines.append(f"3. Apply the {topic_title} framework/process to the given scenario")
        lines.append(f"4. Identify key findings, risks, and recommended actions")
        lines.append(f"5. Prepare a brief presentation (5 minutes) summarizing your analysis")
        lines.append(f"6. Present your findings to the class for discussion and feedback")
    else:
        lines.append(f"1. Review the key concepts covered in {topic_title}")
        lines.append(f"2. Work in groups to analyze the provided case study")
        lines.append(f"3. Apply the concepts to solve the problem scenario")
        lines.append(f"4. Document your approach and findings")
        lines.append(f"5. Present your solution to the class")
    lines.append("")

    lines.append(f"**Expected Output:** Written analysis document with recommendations")
    lines.append(f"**Duration:** 20-30 minutes (including group discussion)")
    lines.append("")

    # Link to assessment methods from CP
    if assessment_methods:
        method_names = [am.get('Assessment_Method', '') for am in assessment_methods if am.get('Assessment_Method')]
        if method_names:
            lines.append(f"**Assessment Link:** This activity prepares you for the {', '.join(method_names)} assessment.")
            lines.append("Submit your work as part of the final assessment.")
            lines.append("")

    return lines


# =============================================================================
# Logo stamping on PDF first page
# =============================================================================

def _stamp_logos_on_first_page(pdf_path: str, company_name: str = None) -> str:
    """Overlay WSQ and company logos on the cover slide (first page) of a PDF.

    Logos are scaled proportionally to the page size so they're always visible.
    WSQ logo top-left, company logo top-right. Copyright footer on ALL pages.

    Args:
        pdf_path: Path to the PDF file to stamp.
        company_name: Optional company name to look up company-specific logo.

    Returns:
        Path to the stamped PDF (same as input, overwritten).
    """
    import os
    from pathlib import Path

    try:
        from reportlab.pdfgen import canvas as rl_canvas
        from reportlab.lib.colors import Color
        from pypdf import PdfReader, PdfWriter
        import tempfile
    except ImportError as e:
        logger.warning(f"Logo stamping skipped — missing dependency: {e}")
        return pdf_path

    # Locate logo files
    project_root = Path(__file__).resolve().parent.parent
    logo_dir = project_root / "assets" / "slide_logos"
    wsq_logo = logo_dir / "wsq_logo.png"
    company_logo = logo_dir / "tertiary_infotech_logo.png"

    # Try company-specific logo from company/logo/ directory
    if company_name:
        safe = company_name.lower().replace(" ", "_").replace(".", "")
        for ext in (".png", ".jpg", ".jpeg"):
            candidate = project_root / "company" / "logo" / f"{safe}{ext}"
            if candidate.exists():
                company_logo = candidate
                break

    if not wsq_logo.exists() and not company_logo.exists():
        logger.warning("No logo files found for stamping")
        return pdf_path

    try:
        reader = PdfReader(pdf_path)
        if not reader.pages:
            return pdf_path

        first_page = reader.pages[0]
        page_width = float(first_page.mediabox.width)
        page_height = float(first_page.mediabox.height)

        # Scale logos proportionally to page size (reference: 1376x768 widescreen)
        scale = page_width / 1376.0
        margin_x = int(15 * scale)
        logo_y = page_height - int(110 * scale)  # Top area, slightly lower for bigger logos

        # WSQ logo: top-left, BIG and prominent
        wsq_w = int(220 * scale)
        wsq_h = int(150 * scale)

        # Company logo: top-right, BIG and prominent
        comp_w = int(400 * scale)
        comp_h = int(110 * scale)

        # --- Cover slide overlay (logos) ---
        overlay_tmp = tempfile.mktemp(suffix="_logo_overlay.pdf")
        c = rl_canvas.Canvas(overlay_tmp, pagesize=(page_width, page_height))

        if wsq_logo.exists():
            c.drawImage(str(wsq_logo), margin_x, logo_y,
                        width=wsq_w, height=wsq_h, preserveAspectRatio=True, mask='auto')

        if company_logo.exists():
            c.drawImage(str(company_logo), page_width - comp_w - margin_x, logo_y,
                        width=comp_w, height=comp_h, preserveAspectRatio=True, mask='auto')

        c.save()

        # --- Copyright footer overlay (all pages) — white background strip ---
        footer_tmp = tempfile.mktemp(suffix="_footer_overlay.pdf")
        fc = rl_canvas.Canvas(footer_tmp, pagesize=(page_width, page_height))
        footer_text = "This material belongs to Tertiary Infotech Academy Pte Ltd (UEN: 20120096W). All Rights Reserved"
        font_size = max(10, int(12 * scale))
        fc.setFont("Helvetica", font_size)
        text_width = fc.stringWidth(footer_text, "Helvetica", font_size)
        strip_height = font_size + int(16 * scale)  # Tall white strip
        footer_y = int(6 * scale)  # Text vertically centered in strip
        # White background strip — covers any slide content at bottom
        fc.setFillColor(Color(1, 1, 1))  # White
        fc.rect(0, 0, page_width, strip_height, fill=True, stroke=False)
        # Black text on white background
        fc.setFillColor(Color(0, 0, 0))  # Black
        fc.drawString((page_width - text_width) / 2, footer_y, footer_text)
        fc.save()

        # Build output PDF
        overlay_reader = PdfReader(overlay_tmp)
        overlay_page = overlay_reader.pages[0]
        footer_reader = PdfReader(footer_tmp)
        footer_page = footer_reader.pages[0]

        writer = PdfWriter()

        # First page: logos + footer
        first_page.merge_page(overlay_page)
        first_page.merge_page(footer_page)
        writer.add_page(first_page)

        # Remaining pages: footer only
        for page in reader.pages[1:]:
            page.merge_page(footer_page)
            writer.add_page(page)

        with open(pdf_path, "wb") as f:
            writer.write(f)

        # Cleanup
        for tmp in [overlay_tmp, footer_tmp]:
            try:
                os.unlink(tmp)
            except OSError:
                pass

        logger.info(f"Stamped logos + footer on {len(reader.pages)} pages of {pdf_path}")
        return pdf_path

    except Exception as e:
        logger.warning(f"Logo stamping failed: {e}")
        return pdf_path


# =============================================================================
# Replace certificate page with branded template image
# =============================================================================

def _replace_certificate_page(pdf_path: str) -> str:
    """Find the 'Certificate of Accomplishment' page in a PDF and replace it
    with the company's branded certificate template image.

    Scans each page for text containing 'certificate' AND 'accomplishment'.
    When found, replaces that page with a full-page rendering of
    assets/slide_logos/certificate_template.png.

    Args:
        pdf_path: Path to the PDF file to process.

    Returns:
        Path to the processed PDF (same as input, overwritten).
    """
    import os
    from pathlib import Path

    try:
        from reportlab.pdfgen import canvas as rl_canvas
        from reportlab.lib.utils import ImageReader
        from pypdf import PdfReader, PdfWriter
        import tempfile
    except ImportError as e:
        logger.warning(f"Certificate replacement skipped — missing dependency: {e}")
        return pdf_path

    # Locate certificate template image
    project_root = Path(__file__).resolve().parent.parent
    cert_img_path = project_root / "assets" / "slide_logos" / "certificate_template.png"

    if not cert_img_path.exists():
        logger.warning("Certificate template image not found at assets/slide_logos/certificate_template.png")
        return pdf_path

    try:
        reader = PdfReader(pdf_path)
        if not reader.pages:
            return pdf_path

        # Find pages containing "certificate" and "accomplishment"
        cert_page_indices = []
        for i, page in enumerate(reader.pages):
            try:
                text = (page.extract_text() or "").lower()
                if "certificate" in text and "accomplishment" in text:
                    cert_page_indices.append(i)
            except Exception:
                continue

        if not cert_page_indices:
            logger.info("No certificate page found in PDF — skipping replacement")
            return pdf_path

        logger.info(f"Found certificate page(s) at index(es): {cert_page_indices}")

        # Get page dimensions from the first certificate page
        ref_page = reader.pages[cert_page_indices[0]]
        page_width = float(ref_page.mediabox.width)
        page_height = float(ref_page.mediabox.height)

        # Create a single-page PDF with the certificate image filling the whole page
        cert_pdf_tmp = tempfile.mktemp(suffix="_cert_page.pdf")
        c = rl_canvas.Canvas(cert_pdf_tmp, pagesize=(page_width, page_height))

        # Draw white background first
        c.setFillColorRGB(1, 1, 1)
        c.rect(0, 0, page_width, page_height, fill=1, stroke=0)

        # Draw certificate image scaled to fill the page while preserving aspect ratio
        img = ImageReader(str(cert_img_path))
        img_w, img_h = img.getSize()
        scale_w = page_width / img_w
        scale_h = page_height / img_h
        scale = min(scale_w, scale_h)  # Fit within page

        draw_w = img_w * scale
        draw_h = img_h * scale
        # Center on page
        x_offset = (page_width - draw_w) / 2
        y_offset = (page_height - draw_h) / 2

        c.drawImage(str(cert_img_path), x_offset, y_offset,
                    width=draw_w, height=draw_h,
                    preserveAspectRatio=True, mask='auto')
        c.save()

        # Read the certificate page PDF
        cert_reader = PdfReader(cert_pdf_tmp)
        cert_replacement = cert_reader.pages[0]

        # Build new PDF, replacing certificate page(s) with the branded image
        writer = PdfWriter()
        for i, page in enumerate(reader.pages):
            if i in cert_page_indices:
                writer.add_page(cert_replacement)
            else:
                writer.add_page(page)

        with open(pdf_path, "wb") as f:
            writer.write(f)

        # Cleanup
        try:
            os.unlink(cert_pdf_tmp)
        except OSError:
            pass

        logger.info(f"Replaced {len(cert_page_indices)} certificate page(s) in {pdf_path}")
        return pdf_path

    except Exception as e:
        logger.warning(f"Certificate page replacement failed: {e}")
        return pdf_path


# =============================================================================
# PPTX: Download, logo stamping, and certificate replacement
# =============================================================================

async def _download_pptx_direct(client, notebook_id: str, output_path: str) -> bool:
    """Try to download PPTX directly from NotebookLM raw artifact data.

    NotebookLM web UI added PPTX export (Feb 2026). The PPTX URL may exist
    at a position beyond the known PDF URL at art[16][3]. This function probes
    the raw artifact metadata for any .pptx URL.

    Returns True if direct PPTX download succeeded, False if not available.
    """
    try:
        raw = await client.artifacts._list_raw(notebook_id)
        for art in raw:
            if not isinstance(art, list) or len(art) <= 4:
                continue
            # ArtifactTypeCode.SLIDE_DECK = 10
            if art[2] != 10:
                continue
            # Status COMPLETED = 3
            if art[4] not in (3, 'ARTIFACT_STATUS_COMPLETED'):
                continue

            # Search the entire artifact for a .pptx URL
            art_json = str(art)
            if '.pptx' not in art_json.lower() and 'export=pptx' not in art_json.lower():
                # No PPTX URL in this artifact
                continue

            # Deep search for PPTX URL in nested structure
            def _find_pptx_url(obj):
                if isinstance(obj, str):
                    if obj.startswith('http') and ('pptx' in obj.lower() or 'export=pptx' in obj.lower() or 'format=pptx' in obj.lower()):
                        return obj
                elif isinstance(obj, list):
                    for item in obj:
                        result = _find_pptx_url(item)
                        if result:
                            return result
                return None

            pptx_url = _find_pptx_url(art)
            if pptx_url:
                logger.info(f"Found direct PPTX URL in artifact data!")
                # Download using the client's internal download method
                await client.artifacts._download_url(pptx_url, output_path)
                if os.path.exists(output_path) and os.path.getsize(output_path) > 1000:
                    logger.info(f"Direct PPTX download success: {output_path}")
                    return True

            # Also try known positions: art[16][4], art[16][5], art[17], art[18]
            if len(art) > 16 and isinstance(art[16], list):
                metadata = art[16]
                for pos in range(4, min(len(metadata), 10)):
                    val = metadata[pos]
                    if isinstance(val, str) and val.startswith('http'):
                        try:
                            await client.artifacts._download_url(val, output_path)
                            if os.path.exists(output_path) and os.path.getsize(output_path) > 1000:
                                # Check if it's actually a PPTX (zip magic bytes PK)
                                with open(output_path, 'rb') as f:
                                    header = f.read(4)
                                if header[:2] == b'PK':
                                    logger.info(f"Direct PPTX from art[16][{pos}]: {output_path}")
                                    return True
                                else:
                                    # Not a PPTX, remove it
                                    os.unlink(output_path)
                        except Exception:
                            if os.path.exists(output_path):
                                os.unlink(output_path)

            # Check higher indices too
            for idx in range(17, min(len(art), 22)):
                if isinstance(art[idx], str) and art[idx].startswith('http'):
                    try:
                        await client.artifacts._download_url(art[idx], output_path)
                        if os.path.exists(output_path) and os.path.getsize(output_path) > 1000:
                            with open(output_path, 'rb') as f:
                                header = f.read(4)
                            if header[:2] == b'PK':
                                logger.info(f"Direct PPTX from art[{idx}]: {output_path}")
                                return True
                            else:
                                os.unlink(output_path)
                    except Exception:
                        if os.path.exists(output_path):
                            os.unlink(output_path)

    except Exception as e:
        logger.debug(f"Direct PPTX probe failed (will fall back to PDF): {e}")

    return False


def _pdf_to_pptx(pdf_path: str, pptx_path: str) -> str:
    """Convert a PDF file to PPTX by rendering each page as a high-res image slide.

    Each PDF page becomes a full-slide image in the PPTX, matching what
    NotebookLM's own PPTX export does (image layers).

    Args:
        pdf_path: Path to the input PDF file.
        pptx_path: Path to save the output PPTX file.

    Returns:
        Path to the created PPTX file.
    """
    import tempfile
    from pathlib import Path

    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
        import fitz  # pymupdf — already installed
    except ImportError as e:
        logger.warning(f"PDF to PPTX conversion skipped — missing dependency: {e}")
        return pdf_path

    try:
        doc = fitz.open(pdf_path)
        prs = Presentation()

        # Set slide dimensions to widescreen 16:9 (matching NotebookLM)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # Blank layout

        for page_num in range(len(doc)):
            page = doc[page_num]
            # Render at 2x resolution for high quality
            mat = fitz.Matrix(2.0, 2.0)
            pix = page.get_pixmap(matrix=mat, alpha=False)

            # Save page as temp PNG
            img_tmp = tempfile.mktemp(suffix=f"_page{page_num}.png")
            pix.save(img_tmp)

            # Add slide with full-page image
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                img_tmp, Emu(0), Emu(0),
                prs.slide_width, prs.slide_height
            )

            # Clean up temp image
            try:
                import os
                os.unlink(img_tmp)
            except Exception:
                pass

        doc.close()
        prs.save(pptx_path)
        logger.info(f"Converted PDF to PPTX: {pptx_path} ({len(doc)} slides)")
        return pptx_path

    except Exception as e:
        logger.warning(f"PDF to PPTX conversion failed: {e}")
        return pdf_path


def _stamp_logos_on_pptx(pptx_path: str, company_name: str = None) -> str:
    """Stamp WSQ and company logos on the cover slide of a PPTX file.

    Also adds copyright footer text to ALL slides.

    Args:
        pptx_path: Path to the PPTX file.
        company_name: Optional company name for company-specific logo.

    Returns:
        Path to the stamped PPTX (same as input, overwritten).
    """
    from pathlib import Path

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError as e:
        logger.warning(f"PPTX logo stamping skipped — missing dependency: {e}")
        return pptx_path

    project_root = Path(__file__).resolve().parent.parent
    wsq_logo = project_root / "assets" / "slide_logos" / "wsq_logo.png"
    company_logo = project_root / "assets" / "slide_logos" / "tertiary_infotech_logo.png"

    if not wsq_logo.exists() and not company_logo.exists():
        logger.warning("No logo files found — skipping PPTX logo stamping")
        return pptx_path

    try:
        prs = Presentation(pptx_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Stamp logos on FIRST slide (cover)
        if prs.slides and len(prs.slides) > 0:
            cover = prs.slides[0]

            # WSQ logo — top-left
            if wsq_logo.exists():
                logo_w = Inches(2.2)
                logo_h = Inches(1.5)
                cover.shapes.add_picture(
                    str(wsq_logo),
                    Inches(0.3), Inches(0.2),
                    logo_w, logo_h
                )

            # Company logo — top-right
            if company_logo.exists():
                logo_w = Inches(3.5)
                logo_h = Inches(1.0)
                left = slide_width - logo_w - Inches(0.3)
                cover.shapes.add_picture(
                    str(company_logo),
                    left, Inches(0.3),
                    logo_w, logo_h
                )

        # Add copyright footer to ALL slides
        footer_text = "This material belongs to Tertiary Infotech Academy Pte Ltd (UEN: 20120096W). All Rights Reserved"
        for slide in prs.slides:
            txBox = slide.shapes.add_textbox(
                Inches(0.5),
                slide_height - Inches(0.45),
                slide_width - Inches(1.0),
                Inches(0.35)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = footer_text
            p.font.size = Pt(7)
            p.font.color.rgb = RGBColor(100, 100, 100)
            p.alignment = PP_ALIGN.CENTER

        prs.save(pptx_path)
        logger.info(f"Stamped logos on PPTX: {pptx_path}")
        return pptx_path

    except Exception as e:
        logger.warning(f"PPTX logo stamping failed: {e}")
        return pptx_path


def _replace_certificate_in_pptx(pptx_path: str) -> str:
    """Replace the certificate slide in a PPTX with the company template.

    Scans slides for text containing 'certificate' and 'accomplishment',
    then replaces matching slides with the branded certificate image.

    Args:
        pptx_path: Path to the PPTX file.

    Returns:
        Path to the processed PPTX (same as input, overwritten).
    """
    from pathlib import Path

    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError as e:
        logger.warning(f"PPTX certificate replacement skipped — missing dependency: {e}")
        return pptx_path

    project_root = Path(__file__).resolve().parent.parent
    cert_img = project_root / "assets" / "slide_logos" / "certificate_template.png"

    if not cert_img.exists():
        logger.warning("Certificate template not found — skipping PPTX certificate replacement")
        return pptx_path

    try:
        prs = Presentation(pptx_path)

        for slide in prs.slides:
            # Extract text from all shapes
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text += shape.text_frame.text.lower() + " "

            if "certificate" in slide_text and "accomplishment" in slide_text:
                # Clear all shapes from this slide
                for shape in list(slide.shapes):
                    sp = shape._element
                    sp.getparent().remove(sp)

                # Add certificate image filling the entire slide
                slide.shapes.add_picture(
                    str(cert_img),
                    Emu(0), Emu(0),
                    prs.slide_width, prs.slide_height
                )
                logger.info("Replaced certificate slide in PPTX")

        prs.save(pptx_path)
        return pptx_path

    except Exception as e:
        logger.warning(f"PPTX certificate replacement failed: {e}")
        return pptx_path


# =============================================================================
# Format course info as text for NotebookLM
# =============================================================================

def _format_course_info_as_text(context: dict) -> str:
    """
    Convert extracted course info dict to comprehensive structured text for NotebookLM.
    Follows the professional WSQ slide template standard with detailed content.

    Args:
        context: Course context dict from Extract Course Info page.

    Returns:
        Formatted text document suitable as a NotebookLM source.
    """
    lines = []

    course_title = context.get('Course_Title', 'Course')
    tgs_ref = context.get('TGS_Ref_No', '')
    tsc_title = context.get('TSC_Title', '')
    tsc_code = context.get('TSC_Code', '')
    duration = context.get('Total_Course_Duration_Hours', '')
    training_hours = context.get('Total_Training_Hours', '')

    # ── SLIDE 1: COVER ──
    lines.append(f"# {course_title}")
    lines.append("")
    if tgs_ref:
        lines.append(f"Course Code: {tgs_ref}")
    if tsc_code:
        lines.append(f"TSC Code: {tsc_code}")
    lines.append("")

    # ── SLIDE 2: DIGITAL ATTENDANCE ──
    lines.append("## Digital Attendance")
    lines.append("- It is mandatory for you to take both AM, PM and Assessment digital attendance for WSQ funded courses.")
    lines.append("- The trainer or administrator will show you the digital attendance QR code generated from SSG portal.")
    lines.append("- Please scan the QR code from your mobile phone camera and submit your attendance.")
    lines.append("")

    # ── SLIDE 3: ABOUT THE TRAINER ──
    lines.append("## About the Trainer")
    lines.append("- Just a professional trainer illustration image only")
    lines.append("- Do NOT write any text — no names, no credentials, no experience, no specialization")
    lines.append("- Keep this slide completely plain with only the title and an illustration")
    lines.append("")

    # ── SLIDE 4: LET'S KNOW EACH OTHER ──
    lines.append("## Let's Know Each Other")
    lines.append("- Name")
    lines.append("- Current role / company")
    lines.append("- Experience in the field")
    lines.append("- What do you hope to learn from this course?")
    lines.append("")

    # ── SLIDE 5: GROUND RULES ──
    lines.append("## Ground Rules")
    lines.append("- Set your mobile phone to silent mode")
    lines.append("- Participate actively in the class. No question is stupid.")
    lines.append("- Mutual respect. Agree to disagree.")
    lines.append("- One conversation at one time.")
    lines.append("- Be punctual. Back from breaks on time.")
    lines.append("- Exit the class silently if you need to step out for phone call, toilet break etc.")
    lines.append("- 75% attendance is required for WSQ funding eligibility.")
    lines.append("")

    # ── SLIDE 6-7: LESSON PLAN (before Skills Framework) ──
    lines.append("## Lesson Plan")
    lines.extend(_build_lesson_plan_text(context))

    # ── SLIDE 8: SKILLS FRAMEWORK ──
    lines.append("## Skills Framework")
    if tsc_title:
        lines.append(f"**TSC Title:** {tsc_title}")
    if tsc_code:
        lines.append(f"**TSC Code:** {tsc_code}")
    lines.append("")

    # TSC Description
    description = context.get('TSC_Description') or context.get('Proficiency_Description')
    if description:
        lines.append("### TSC Description")
        lines.append(description)
        lines.append("")

    # ── SLIDE 9: TSC K&A STATEMENTS ──
    all_k_statements = []
    all_a_statements = []
    for lu in context.get('Learning_Units', []):
        for k in lu.get('K_numbering_description', []):
            k_entry = f"**{k.get('K_number', '')}:** {k.get('Description', '')}"
            if k_entry not in all_k_statements:
                all_k_statements.append(k_entry)
        for a in lu.get('A_numbering_description', []):
            a_entry = f"**{a.get('A_number', '')}:** {a.get('Description', '')}"
            if a_entry not in all_a_statements:
                all_a_statements.append(a_entry)

    if all_k_statements:
        lines.append("### TSC Knowledge Statements")
        for k in all_k_statements:
            lines.append(f"- {k}")
        lines.append("")

    if all_a_statements:
        lines.append("### TSC Ability Statements")
        for a in all_a_statements:
            lines.append(f"- {a}")
        lines.append("")

    # ── SLIDE 10: LEARNING OUTCOMES ──
    lines.append("## Learning Outcomes")
    for lu_idx, lu in enumerate(context.get('Learning_Units', []), 1):
        lo_num = lu.get('LO_Number', f'LO{lu_idx}')
        if lu.get('LO'):
            lines.append(f"- **{lo_num}:** {lu['LO']}")
    lines.append("")

    # ── SLIDE 11-12: COURSE OUTLINE ──
    lines.append("## Course Outline")
    for lu_idx, lu in enumerate(context.get('Learning_Units', []), 1):
        lu_title = lu.get('LU_Title', 'Learning Unit')
        lu_num = lu.get('LU_Number', f'LU{lu_idx}')

        # Collect K&A numbers for this LU
        k_nums = [k.get('K_number', '') for k in lu.get('K_numbering_description', [])]
        a_nums = [a.get('A_number', '') for a in lu.get('A_numbering_description', [])]
        ka_refs = ', '.join(k_nums + a_nums)
        ka_str = f" ({ka_refs})" if ka_refs else ""

        lines.append(f"### {lu_num}: {lu_title}{ka_str}")
        for t_idx, topic in enumerate(lu.get('Topics', []), 1):
            topic_title = topic.get('Topic_Title', 'Topic')
            lines.append(f"- T{t_idx}: {topic_title}")
        lines.append("")

    # ── SECTION 5: ASSESSMENT INFORMATION ──
    assessment_details = context.get('Assessment_Methods_Details', [])
    if assessment_details:
        lines.append("## Final Assessment")
        for am in assessment_details:
            method = am.get('Assessment_Method', '')
            abbr = am.get('Method_Abbreviation', '')
            duration_am = am.get('Total_Delivery_Hours', '')
            label = f"{method} ({abbr})" if abbr else method
            lines.append(f"- {label} - {duration_am}")
        lines.append("- Assessment format: Open Book")
        lines.append("- Open book assessment ONLY includes Slides, Learner Guide or any approved materials.")
        lines.append("- Appeal process")
        lines.append("")

        lines.append("## Briefing for Assessment")
        lines.append("- Place phones and other materials under the table or on the floor")
        lines.append("- No photos or recording of assessment scripts")
        lines.append("- No discussion during assessment")
        lines.append("- Use black/blue pen for assessment (hard copies)")
        lines.append("- No usage of liquid paper or correction tape")
        lines.append("- Assessment scripts will be collected when time is up")
        lines.append("")

    # ── CRITERIA FOR FUNDING ──
    lines.append("## Criteria for Funding")
    lines.append("- Minimum attendance rate of 75% based on SSG Digital Attendance record.")
    lines.append("- Complete the assessment and be assessed as 'Competent'")
    lines.append("")

    # ── TRAQOM SURVEY ──
    lines.append("## TRAQOM Survey")
    lines.append("- Access the survey here (link will be provided)")
    lines.append("- Key in your Last four NRIC/FIN characters and the Six-digit course run ID to complete the survey")
    lines.append("")

    # ── SECTION 6: DETAILED TOPIC CONTENT (MAIN BODY) ──
    # Each LU becomes a major section with deep, expanded topic slides
    assessment_details = context.get('Assessment_Methods_Details', [])
    for lu_idx, lu in enumerate(context.get('Learning_Units', []), 1):
        lu_title = lu.get('LU_Title', 'Learning Unit')
        lu_num = lu.get('LU_Number', f'LU{lu_idx}')
        lo_num = lu.get('LO_Number', f'LO{lu_idx}')

        lines.append(f"## {lu_num}: {lu_title}")
        lines.append("")

        if lu.get('LO'):
            lines.append(f"**{lo_num} - Learning Outcome:** {lu['LO']}")
            lines.append("")

        for t_idx, topic in enumerate(lu.get('Topics', []), 1):
            topic_title = topic.get('Topic_Title', 'Topic')
            lines.append(f"### T{t_idx}: {topic_title}")
            lines.append("")

            # Include all bullet points as detailed content
            bullet_points = topic.get('Bullet_Points', [])
            for bp in bullet_points:
                lines.append(f"- {bp}")

            lines.append("")

            # Add Activity/Lab slide after each topic
            lines.extend(_build_activity_text(topic_title, bullet_points, assessment_details))

        # K&A statements for this LU
        k_statements = lu.get('K_numbering_description', [])
        if k_statements:
            lines.append(f"### {lu_num} - Knowledge Statements")
            for k in k_statements:
                lines.append(f"- **{k.get('K_number', '')}:** {k.get('Description', '')}")
            lines.append("")

        a_statements = lu.get('A_numbering_description', [])
        if a_statements:
            lines.append(f"### {lu_num} - Ability Statements")
            for a in a_statements:
                lines.append(f"- **{a.get('A_number', '')}:** {a.get('Description', '')}")
            lines.append("")

        # Instructional methods
        methods = lu.get('Instructional_Methods', [])
        if methods:
            lines.append(f"**Instructional Methods:** {', '.join(methods)}")
            lines.append("")

    # ── CLOSING SLIDES — EXACT supervisor order ──
    # 1. Summary & Q&A
    lines.append("## Summary & Q&A")
    lines.append("Recap of all learning outcomes and key concepts covered.")
    for lu_idx, lu in enumerate(context.get('Learning_Units', []), 1):
        lo_num = lu.get('LO_Number', f'LO{lu_idx}')
        lu_title = lu.get('LU_Title', '')
        if lu.get('LO'):
            lines.append(f"- **{lo_num} ({lu_title}):** {lu['LO']}")
    lines.append("")

    # 2. TRAQOM Survey
    lines.append("## TRAQOM Survey")
    lines.append("- Access the survey here (link will be provided)")
    lines.append("- Key in your Last four NRIC/FIN characters and the Six-digit course run ID to complete the survey")
    lines.append("")

    # 3. Certificate of Accomplishment
    lines.append("## Certificate of Accomplishment")
    lines.append("Please provide the following details to facilitate the issuance of your certificate after the class and to ensure the accuracy of the information printed on it.")
    lines.append("")
    lines.append("### Certificate Template:")
    lines.append("CERTIFICATE OF ACCOMPLISHMENT")
    lines.append("This Certificate is proudly presented to")
    lines.append("[Student Name]")
    lines.append("For completing the course")
    lines.append("[Course Name]")
    lines.append("Held on [Course Dates]")
    lines.append("")
    lines.append("Dr. Alfred Ang")
    lines.append("Managing Director")
    lines.append("Tertiary Infotech Academy")
    lines.append("")

    # 4. Digital Attendance
    lines.append("## Digital Attendance")
    lines.append("- It is mandatory for you to take both AM, PM and Assessment digital attendance for WSQ funded courses.")
    lines.append("- Please scan the QR code to submit your attendance.")
    lines.append("")

    # 5. Final Assessment (section header)
    lines.append("## Final Assessment")
    lines.append("(Final Assessment section)")
    lines.append("")

    # 6. Recommended Courses
    lines.append("## Recommended Courses")
    lines.append("- Explore other WSQ courses to continue your professional development.")
    lines.append("")

    # 7. Support
    lines.append("## Support")
    lines.append("- If you have any enquiries during and after the class, you can contact us:")
    lines.append("- Email: enquiry@tertiaryinfotech.com")
    lines.append("- Tel: 61000613")
    lines.append("- Website: www.tertiaryinfotech.com")
    lines.append("")

    # 8. Thank You
    lines.append("## Thank You")
    lines.append("Thank you for attending this course!")
    lines.append("We wish you all the best in applying what you have learned.")
    lines.append("")

    return "\n".join(lines)


def _format_lu_source_text(context: dict, lu_index: int, num_lus: int) -> str:
    """
    Format source text for a single Learning Unit's slide deck.

    Args:
        context: Full course context dict.
        lu_index: 0-based index of the LU to format.
        num_lus: Total number of LUs.

    Returns:
        Formatted text for one LU's NotebookLM source.
    """
    lines = []
    lus = context.get('Learning_Units', [])
    lu = lus[lu_index]

    course_title = context.get('Course_Title', 'Course')
    tgs_ref = context.get('TGS_Ref_No', '')
    tsc_title = context.get('TSC_Title', '')
    tsc_code = context.get('TSC_Code', '')
    lu_title = lu.get('LU_Title', 'Learning Unit')
    lu_num = lu.get('LU_Number', f'LU{lu_index + 1}')
    lo_num = lu.get('LO_Number', f'LO{lu_index + 1}')

    is_first = lu_index == 0
    is_last = lu_index == num_lus - 1

    # ── INTRO STANDARD PAGES (first LU only) ──
    if is_first:
        lines.append(f"# {course_title}")
        lines.append("")
        if tgs_ref:
            lines.append(f"Course Code: {tgs_ref}")
        if tsc_code:
            lines.append(f"TSC Code: {tsc_code}")
        lines.append("")

        lines.append("## Digital Attendance")
        lines.append("- It is mandatory for you to take both AM, PM and Assessment digital attendance for WSQ funded courses.")
        lines.append("- The trainer or administrator will show you the digital attendance QR code generated from SSG portal.")
        lines.append("- Please scan the QR code from your mobile phone camera and submit your attendance.")
        lines.append("")

        lines.append("## About the Trainer")
        lines.append("- Just a professional trainer illustration image only")
        lines.append("- Do NOT write any text — no names, no credentials, no experience, no specialization")
        lines.append("- Keep this slide completely plain with only the title and an illustration")
        lines.append("")

        lines.append("## Let's Know Each Other")
        lines.append("- Name")
        lines.append("- Current role / company")
        lines.append("- Experience in the field")
        lines.append("- What do you hope to learn from this course?")
        lines.append("")

        lines.append("## Ground Rules")
        lines.append("- Set your mobile phone to silent mode")
        lines.append("- Participate actively in the class. No question is stupid.")
        lines.append("- Mutual respect. Agree to disagree.")
        lines.append("- One conversation at one time.")
        lines.append("- Be punctual. Back from breaks on time.")
        lines.append("- Exit the class silently if you need to step out for phone call, toilet break etc.")
        lines.append("- 75% attendance is required for WSQ funding eligibility.")
        lines.append("")

        # Lesson Plan (before Skills Framework — matching supervisor order)
        lines.append("## Lesson Plan")
        lines.extend(_build_lesson_plan_text(context))

        # Skills Framework
        lines.append("## Skills Framework")
        if tsc_title:
            lines.append(f"**TSC Title:** {tsc_title}")
        if tsc_code:
            lines.append(f"**TSC Code:** {tsc_code}")
        lines.append("")

        description = context.get('TSC_Description') or context.get('Proficiency_Description')
        if description:
            lines.append("### TSC Description")
            lines.append(description)
            lines.append("")

        # All K&A statements
        all_k = []
        all_a = []
        for u in lus:
            for k in u.get('K_numbering_description', []):
                entry = f"**{k.get('K_number', '')}:** {k.get('Description', '')}"
                if entry not in all_k:
                    all_k.append(entry)
            for a in u.get('A_numbering_description', []):
                entry = f"**{a.get('A_number', '')}:** {a.get('Description', '')}"
                if entry not in all_a:
                    all_a.append(entry)
        if all_k:
            lines.append("### TSC Knowledge Statements")
            for k in all_k:
                lines.append(f"- {k}")
            lines.append("")
        if all_a:
            lines.append("### TSC Ability Statements")
            for a in all_a:
                lines.append(f"- {a}")
            lines.append("")

        # Learning Outcomes
        lines.append("## Learning Outcomes")
        for i, u in enumerate(lus, 1):
            lo = u.get('LO_Number', f'LO{i}')
            if u.get('LO'):
                lines.append(f"- **{lo}:** {u['LO']}")
        lines.append("")

        # Course Outline
        lines.append("## Course Outline")
        for i, u in enumerate(lus, 1):
            un = u.get('LU_Number', f'LU{i}')
            ut = u.get('LU_Title', 'Learning Unit')
            k_nums = [k.get('K_number', '') for k in u.get('K_numbering_description', [])]
            a_nums = [a.get('A_number', '') for a in u.get('A_numbering_description', [])]
            ka_refs = ', '.join(k_nums + a_nums)
            ka_str = f" ({ka_refs})" if ka_refs else ""
            lines.append(f"### {un}: {ut}{ka_str}")
            for j, t in enumerate(u.get('Topics', []), 1):
                lines.append(f"- T{j}: {t.get('Topic_Title', 'Topic')}")
            lines.append("")

    # ── LU CONTENT (always included) ──
    assessment_details = context.get('Assessment_Methods_Details', [])
    lines.append(f"## {lu_num}: {lu_title}")
    lines.append("")
    if lu.get('LO'):
        lines.append(f"**{lo_num} - Learning Outcome:** {lu['LO']}")
        lines.append("")

    lines.append("**IMPORTANT INSTRUCTION FOR SLIDE GENERATION:** For each topic below, "
                 "create 15-20 detailed content slides. Each bullet point represents a concept "
                 "that needs MULTIPLE slides with full explanations, definitions, examples, "
                 "step-by-step processes, and practical applications. Do NOT just list the bullet "
                 "points on one slide — EXPAND each into detailed teaching content that a learner "
                 "can study from. Think of each bullet as a section heading that needs 2-4 slides of content.")
    lines.append("")

    for t_idx, topic in enumerate(lu.get('Topics', []), 1):
        topic_title = topic.get('Topic_Title', 'Topic')
        lines.append(f"### T{t_idx}: {topic_title}")
        lines.append("")
        lines.append(f"Create detailed slides covering the following concepts for '{topic_title}':")
        lines.append("")
        bullet_points = topic.get('Bullet_Points', [])
        for bp_idx, bp in enumerate(bullet_points, 1):
            lines.append(f"- **Concept {bp_idx}: {bp}**")
            lines.append(f"  - Define what '{bp}' means in the context of {topic_title}")
            lines.append(f"  - Explain the process, methodology, or framework involved")
            lines.append(f"  - Provide a real-world example or case study")
            lines.append(f"  - List the key steps or components")
        lines.append("")

        # Add Activity/Lab after each topic
        lines.extend(_build_activity_text(topic_title, bullet_points, assessment_details))

    # K&A for this LU
    k_statements = lu.get('K_numbering_description', [])
    if k_statements:
        lines.append(f"### {lu_num} - Knowledge Statements")
        for k in k_statements:
            lines.append(f"- **{k.get('K_number', '')}:** {k.get('Description', '')}")
        lines.append("")

    a_statements = lu.get('A_numbering_description', [])
    if a_statements:
        lines.append(f"### {lu_num} - Ability Statements")
        for a in a_statements:
            lines.append(f"- **{a.get('A_number', '')}:** {a.get('Description', '')}")
        lines.append("")

    methods = lu.get('Instructional_Methods', [])
    if methods:
        lines.append(f"**Instructional Methods:** {', '.join(methods)}")
        lines.append("")

    # ── CLOSING PAGES (last LU only) — EXACT supervisor order ──
    if is_last:
        # 1. Summary & Q&A
        lines.append("## Summary & Q&A")
        lines.append("Recap of all learning outcomes and key concepts covered.")
        for i, u in enumerate(lus, 1):
            lo = u.get('LO_Number', f'LO{i}')
            ut = u.get('LU_Title', '')
            if u.get('LO'):
                lines.append(f"- **{lo} ({ut}):** {u['LO']}")
        lines.append("")

        # 2. TRAQOM Survey
        lines.append("## TRAQOM Survey")
        lines.append("- Access the survey here (link will be provided)")
        lines.append("- Key in your Last four NRIC/FIN characters and the Six-digit course run ID to complete the survey")
        lines.append("")

        # 3. Certificate of Accomplishment
        lines.append("## Certificate of Accomplishment")
        lines.append("(Company certificate template slide)")
        lines.append("")

        # 4. Digital Attendance
        lines.append("## Digital Attendance")
        lines.append("- It is mandatory for you to take both AM, PM and Assessment digital attendance for WSQ funded courses.")
        lines.append("- Please scan the QR code to submit your attendance.")
        lines.append("")

        # 5. Final Assessment
        assessment_details = context.get('Assessment_Methods_Details', [])
        if assessment_details:
            lines.append("## Final Assessment")
            for am in assessment_details:
                method = am.get('Assessment_Method', '')
                abbr = am.get('Method_Abbreviation', '')
                dur = am.get('Total_Delivery_Hours', '')
                label = f"{method} ({abbr})" if abbr else method
                lines.append(f"- {label} - {dur}")
            lines.append("- Assessment format: Open Book")
            lines.append("- Open book assessment ONLY includes Slides, Learner Guide or any approved materials.")
            lines.append("")

        # 6. Recommended Courses
        lines.append("## Recommended Courses")
        lines.append("- Explore other WSQ courses to continue your professional development.")
        lines.append("")

        # 7. Support
        lines.append("## Support")
        lines.append("- If you have any enquiries during and after the class, you can contact us:")
        lines.append("- Email: enquiry@tertiaryinfotech.com")
        lines.append("- Tel: 61000613")
        lines.append("- Website: www.tertiaryinfotech.com")
        lines.append("")

        # 8. Thank You
        lines.append("## Thank You")
        lines.append("Thank you for attending this course!")
        lines.append("We wish you all the best in applying what you have learned.")
        lines.append("")

    return "\n".join(lines)


def _format_chunk_source_text(context: dict, lu_index: int, num_lus: int,
                               topic_indices: list, chunk_label: str,
                               is_first_chunk_of_course: bool,
                               is_last_chunk_of_course: bool) -> str:
    """
    Format source text for a CHUNK of topics within a Learning Unit.

    This is used when an LU has many topics and we split them across
    multiple NotebookLM decks (~4-5 topics per deck for ~20 slides each).

    Args:
        context: Full course context dict.
        lu_index: 0-based index of the LU.
        num_lus: Total number of LUs.
        topic_indices: List of 0-based topic indices within this LU to include.
        chunk_label: Label like "Deck 1A", "Deck 1B" etc.
        is_first_chunk_of_course: Whether this is the very first chunk (include intro pages).
        is_last_chunk_of_course: Whether this is the very last chunk (include closing pages).

    Returns:
        Formatted text for this chunk's NotebookLM source.
    """
    lines = []
    lus = context.get('Learning_Units', [])
    lu = lus[lu_index]

    course_title = context.get('Course_Title', 'Course')
    tgs_ref = context.get('TGS_Ref_No', '')
    tsc_title = context.get('TSC_Title', '')
    tsc_code = context.get('TSC_Code', '')
    lu_title = lu.get('LU_Title', 'Learning Unit')
    lu_num = lu.get('LU_Number', f'LU{lu_index + 1}')
    lo_num = lu.get('LO_Number', f'LO{lu_index + 1}')

    # ── INTRO STANDARD PAGES (first chunk of course only) ──
    if is_first_chunk_of_course:
        lines.append(f"# {course_title}")
        lines.append("")
        if tgs_ref:
            lines.append(f"Course Code: {tgs_ref}")
        if tsc_code:
            lines.append(f"TSC Code: {tsc_code}")
        lines.append("")

        lines.append("## Digital Attendance")
        lines.append("- It is mandatory for you to take both AM, PM and Assessment digital attendance for WSQ funded courses.")
        lines.append("- The trainer or administrator will show you the digital attendance QR code generated from SSG portal.")
        lines.append("- Please scan the QR code from your mobile phone camera and submit your attendance.")
        lines.append("")

        lines.append("## About the Trainer")
        lines.append("- Just a professional trainer illustration image only")
        lines.append("- Do NOT write any text — no names, no credentials, no experience, no specialization")
        lines.append("- Keep this slide completely plain with only the title and an illustration")
        lines.append("")

        lines.append("## Let's Know Each Other")
        lines.append("- Name")
        lines.append("- Current role / company")
        lines.append("- Experience in the field")
        lines.append("- What do you hope to learn from this course?")
        lines.append("")

        lines.append("## Ground Rules")
        lines.append("- Set your mobile phone to silent mode")
        lines.append("- Participate actively in the class. No question is stupid.")
        lines.append("- Mutual respect. Agree to disagree.")
        lines.append("- One conversation at one time.")
        lines.append("- Be punctual. Back from breaks on time.")
        lines.append("- Exit the class silently if you need to step out for phone call, toilet break etc.")
        lines.append("- 75% attendance is required for WSQ funding eligibility.")
        lines.append("")

        # Lesson Plan (before Skills Framework — matching supervisor order)
        lines.append("## Lesson Plan")
        lines.extend(_build_lesson_plan_text(context))

        lines.append("## Skills Framework")
        if tsc_title:
            lines.append(f"**TSC Title:** {tsc_title}")
        if tsc_code:
            lines.append(f"**TSC Code:** {tsc_code}")
        lines.append("")

        description = context.get('TSC_Description') or context.get('Proficiency_Description')
        if description:
            lines.append("### TSC Description")
            lines.append(description)
            lines.append("")

        all_k = []
        all_a = []
        for u in lus:
            for k in u.get('K_numbering_description', []):
                entry = f"**{k.get('K_number', '')}:** {k.get('Description', '')}"
                if entry not in all_k:
                    all_k.append(entry)
            for a in u.get('A_numbering_description', []):
                entry = f"**{a.get('A_number', '')}:** {a.get('Description', '')}"
                if entry not in all_a:
                    all_a.append(entry)
        if all_k:
            lines.append("### TSC Knowledge Statements")
            for k in all_k:
                lines.append(f"- {k}")
            lines.append("")
        if all_a:
            lines.append("### TSC Ability Statements")
            for a in all_a:
                lines.append(f"- {a}")
            lines.append("")

        lines.append("## Learning Outcomes")
        for i, u in enumerate(lus, 1):
            lo = u.get('LO_Number', f'LO{i}')
            if u.get('LO'):
                lines.append(f"- **{lo}:** {u['LO']}")
        lines.append("")

        lines.append("## Course Outline")
        for i, u in enumerate(lus, 1):
            un = u.get('LU_Number', f'LU{i}')
            ut = u.get('LU_Title', 'Learning Unit')
            k_nums = [k.get('K_number', '') for k in u.get('K_numbering_description', [])]
            a_nums = [a.get('A_number', '') for a in u.get('A_numbering_description', [])]
            ka_refs = ', '.join(k_nums + a_nums)
            ka_str = f" ({ka_refs})" if ka_refs else ""
            lines.append(f"### {un}: {ut}{ka_str}")
            for j, t in enumerate(u.get('Topics', []), 1):
                lines.append(f"- T{j}: {t.get('Topic_Title', 'Topic')}")
            lines.append("")

    # ── LU + TOPIC CHUNK CONTENT ──
    assessment_details = context.get('Assessment_Methods_Details', [])
    lines.append(f"## {lu_num}: {lu_title}")
    lines.append("")
    if lu.get('LO'):
        lines.append(f"**{lo_num} - Learning Outcome:** {lu['LO']}")
        lines.append("")

    lines.append("**IMPORTANT INSTRUCTION FOR SLIDE GENERATION:** For each topic below, "
                 "create 15-20 detailed content slides. Each bullet point represents a concept "
                 "that needs MULTIPLE slides with full explanations, definitions, examples, "
                 "step-by-step processes, and practical applications. Do NOT just list the bullet "
                 "points on one slide — EXPAND each into detailed teaching content that a learner "
                 "can study from. Think of each bullet as a section heading that needs 2-4 slides of content.")
    lines.append("")

    topics = lu.get('Topics', [])
    for t_idx in topic_indices:
        if t_idx >= len(topics):
            continue
        topic = topics[t_idx]
        topic_title = topic.get('Topic_Title', 'Topic')

        # K&A refs for this topic
        k_refs = [k.get('K_number', '') for k in lu.get('K_numbering_description', [])]
        a_refs = [a.get('A_number', '') for a in lu.get('A_numbering_description', [])]
        ka_str = f" ({', '.join(k_refs + a_refs)})" if (k_refs or a_refs) else ""

        lines.append(f"### T{t_idx + 1}: {topic_title}{ka_str}")
        lines.append("")
        lines.append(f"Create detailed slides covering the following concepts for '{topic_title}':")
        lines.append("")
        bullet_points = topic.get('Bullet_Points', [])
        for bp_idx, bp in enumerate(bullet_points, 1):
            lines.append(f"- **Concept {bp_idx}: {bp}**")
            lines.append(f"  - Define what '{bp}' means in the context of {topic_title}")
            lines.append(f"  - Explain the process, methodology, or framework involved")
            lines.append(f"  - Provide a real-world example or case study")
            lines.append(f"  - List the key steps or components")
        lines.append("")

        # Activity/Lab after each topic (proper structured format)
        lines.extend(_build_activity_text(topic_title, bullet_points, assessment_details))
        lines.append("")

    # K&A for this LU
    k_statements = lu.get('K_numbering_description', [])
    if k_statements:
        lines.append(f"### {lu_num} - Knowledge Statements")
        for k in k_statements:
            lines.append(f"- **{k.get('K_number', '')}:** {k.get('Description', '')}")
        lines.append("")

    a_statements = lu.get('A_numbering_description', [])
    if a_statements:
        lines.append(f"### {lu_num} - Ability Statements")
        for a in a_statements:
            lines.append(f"- **{a.get('A_number', '')}:** {a.get('Description', '')}")
        lines.append("")

    # ── CLOSING PAGES (last chunk of course only) ──
    if is_last_chunk_of_course:
        # Closing slides in EXACT supervisor order
        # 1. Summary & Q&A
        lines.append("## Summary & Q&A")
        lines.append("Recap of all learning outcomes and key concepts covered.")
        for i, u in enumerate(lus, 1):
            lo = u.get('LO_Number', f'LO{i}')
            ut = u.get('LU_Title', '')
            if u.get('LO'):
                lines.append(f"- **{lo} ({ut}):** {u['LO']}")
        lines.append("")

        # 2. TRAQOM Survey
        lines.append("## TRAQOM Survey")
        lines.append("- Access the survey here (link will be provided)")
        lines.append("- Key in your Last four NRIC/FIN characters and the Six-digit course run ID to complete the survey")
        lines.append("")

        # 3. Certificate of Accomplishment
        lines.append("## Certificate of Accomplishment")
        lines.append("(Company certificate template slide)")
        lines.append("")

        # 4. Digital Attendance
        lines.append("## Digital Attendance")
        lines.append("- It is mandatory for you to take both AM, PM and Assessment digital attendance for WSQ funded courses.")
        lines.append("- Please scan the QR code to submit your attendance.")
        lines.append("")

        # 5. Final Assessment
        assessment_details = context.get('Assessment_Methods_Details', [])
        if assessment_details:
            lines.append("## Final Assessment")
            for am in assessment_details:
                method = am.get('Assessment_Method', '')
                abbr = am.get('Method_Abbreviation', '')
                dur = am.get('Total_Delivery_Hours', '')
                label = f"{method} ({abbr})" if abbr else method
                lines.append(f"- {label} - {dur}")
            lines.append("- Assessment format: Open Book")
            lines.append("- Open book assessment ONLY includes Slides, Learner Guide or any approved materials.")
            lines.append("")

        # 6. Recommended Courses
        lines.append("## Recommended Courses")
        lines.append("- Explore other WSQ courses to continue your professional development.")
        lines.append("")

        # 7. Support
        lines.append("## Support")
        lines.append("- If you have any enquiries during and after the class, you can contact us:")
        lines.append("- Email: enquiry@tertiaryinfotech.com")
        lines.append("- Tel: 61000613")
        lines.append("- Website: www.tertiaryinfotech.com")
        lines.append("")

        # 8. Thank You
        lines.append("## Thank You")
        lines.append("Thank you for attending this course!")
        lines.append("We wish you all the best in applying what you have learned.")
        lines.append("")

    return "\n".join(lines)


# =============================================================================
# Multi-source helpers — build 3-5 reliable sources per notebook
# =============================================================================

def _build_course_context_source(context: dict, cm: dict) -> str:
    """Build a concise course context source with framework, outcomes, and structure.
    This gives NotebookLM a high-level overview separate from the detailed topic content."""
    lines = []
    course_title = context.get('Course_Title', 'Course')
    tsc_title = context.get('TSC_Title', '')
    tsc_code = context.get('TSC_Code', '')
    description = context.get('TSC_Description') or context.get('Proficiency_Description') or ''
    lus = context.get('Learning_Units', [])

    lines.append(f"# Course Framework: {course_title}")
    lines.append("")
    if tsc_title:
        lines.append(f"**TSC Title:** {tsc_title}")
    if tsc_code:
        lines.append(f"**TSC Code:** {tsc_code}")
    if description:
        lines.append(f"\n**Description:** {description}")
    lines.append("")

    # All K&A statements across the course
    lines.append("## Knowledge & Ability Statements")
    for lu in lus:
        lu_num = lu.get('LU_Number', '')
        for k in lu.get('K_numbering_description', []):
            lines.append(f"- **{k.get('K_number', '')}:** {k.get('Description', '')}")
        for a in lu.get('A_numbering_description', []):
            lines.append(f"- **{a.get('A_number', '')}:** {a.get('Description', '')}")
    lines.append("")

    # All learning outcomes
    lines.append("## Learning Outcomes")
    for i, lu in enumerate(lus, 1):
        lo = lu.get('LO_Number', f'LO{i}')
        if lu.get('LO'):
            lines.append(f"- **{lo}:** {lu['LO']}")
    lines.append("")

    # Assessment info
    assessment_details = context.get('Assessment_Methods_Details', [])
    if assessment_details:
        lines.append("## Assessment Methods")
        for am in assessment_details:
            method = am.get('Assessment_Method', '')
            abbr = am.get('Method_Abbreviation', '')
            label = f"{method} ({abbr})" if abbr else method
            lines.append(f"- {label}")
        lines.append("")

    # Full course outline
    lines.append("## Complete Course Outline")
    for i, lu in enumerate(lus, 1):
        un = lu.get('LU_Number', f'LU{i}')
        ut = lu.get('LU_Title', '')
        lines.append(f"### {un}: {ut}")
        for j, t in enumerate(lu.get('Topics', []), 1):
            lines.append(f"  - T{j}: {t.get('Topic_Title', '')}")
        lines.append("")

    return "\n".join(lines)[:15000]  # Keep concise


def _extract_topic_names_from_cm(cm: dict) -> List[str]:
    """Extract clean topic names from chunk metadata for Wikipedia lookups."""
    content = cm.get('content', '')
    topics = []
    for line in content.split('\n'):
        # Match topic headings like "### T1: Topic Title (K1, A1)"
        m = re.match(r'^###\s+T\d+:\s*(.+?)(?:\s*\(.*\))?\s*$', line)
        if m:
            topic = m.group(1).strip()
            if topic and len(topic) > 3:
                topics.append(topic)
    return topics[:5]  # Max 5 topics


def _build_topic_search_url(topic_names: list, course_title: str) -> List[str]:
    """Build direct Wikipedia article URLs from topic names for web source imports.

    Uses direct article URLs (not search pages) — search pages redirect
    and often fail as NotebookLM sources.
    """
    urls = []
    for topic in topic_names[:2]:  # Max 2 URL sources
        # Clean topic name for Wikipedia article URL
        clean = re.sub(r'\([^)]*\)', '', topic).strip()
        clean = re.sub(r'^(Topic\s*\d+[:\-\s]*)', '', clean, flags=re.IGNORECASE).strip()
        if len(clean) > 5:
            wiki_title = clean.replace(' ', '_')
            urls.append(f"https://en.wikipedia.org/wiki/{wiki_title}")
    return urls


async def _add_multi_sources(client, nb_id: str, cm: dict, course_title: str,
                              context: dict, logger_obj=None) -> List[str]:
    """Add sources to a notebook:
    1. Main course material (text source - instant)
    2. Course framework overview (text source - instant)
    3-4. Web URL sources from Wikipedia for topic content (internet sources)

    Returns list of all source IDs added successfully.
    """
    source_ids = []
    log = logger_obj or logger

    # Source 1: Main course material (the detailed content)
    try:
        src1 = await client.sources.add_text(
            nb_id,
            f"{course_title} - {cm['label']} (Course Material)",
            cm['content'][:50000]
        )
        source_ids.append(src1.id)
    except Exception as e:
        log.warning(f"Failed to add main source: {e}")
        return source_ids  # Can't proceed without main source

    # Source 2: Course framework overview (K&A statements, LOs, outline)
    try:
        ctx_text = _build_course_context_source(context, cm)
        if ctx_text and len(ctx_text) > 100:
            src2 = await client.sources.add_text(
                nb_id,
                f"{course_title} - Course Framework & K&A Statements",
                ctx_text
            )
            source_ids.append(src2.id)
    except Exception as e:
        log.info(f"Skipped framework source: {e}")

    # Sources 3-4: Web URL sources for internet content quality
    # Uses direct Wikipedia article URLs (not search pages) for reliability
    try:
        topic_names = cm.get('topic_names', [])
        # Build reliable direct URLs (Wikipedia articles, not search pages)
        direct_urls = []
        for topic in topic_names[:2]:
            clean = re.sub(r'\([^)]*\)', '', topic).strip()
            clean = re.sub(r'^(T\d+[:\-\s]*)', '', clean, flags=re.IGNORECASE).strip()
            # Use direct Wikipedia article URLs for common BCM concepts
            if len(clean) > 5:
                wiki_title = clean.replace(' ', '_')
                direct_urls.append(f"https://en.wikipedia.org/wiki/{wiki_title}")

        for url in direct_urls:
            try:
                src_url = await client.sources.add_url(nb_id, url)
                source_ids.append(src_url.id)
                log.info(f"Added URL source: {url}")
            except Exception as e:
                log.info(f"Skipped URL source {url}: {e}")
    except Exception as e:
        log.info(f"Skipped URL sources: {e}")

    return source_ids


# =============================================================================
# NotebookLM helpers
# =============================================================================

def _check_notebooklm_available() -> bool:
    """Check if notebooklm-py library is installed."""
    try:
        from notebooklm import NotebookLMClient  # noqa: F401
        return True
    except ImportError:
        return False


def _extract_research_queries(content: str, course_title: str,
                               num_queries: int = 2) -> List[str]:
    """
    Extract research queries from document content by finding key topics.
    Generates diverse queries for comprehensive, high-quality source material.
    Pure text parsing — no LLM needed.
    """
    queries = []
    lines = content.split('\n')

    # Extract topic headings
    topic_candidates = []
    for line in lines:
        stripped = line.strip()
        if not stripped or len(stripped) < 5:
            continue
        if re.match(r'^#{2,3}\s+\w', stripped) and len(stripped) < 120:
            topic_candidates.append(re.sub(r'^#{2,3}\s+', '', stripped))
        elif re.match(r'^\d+\.?\d*\s+\w', stripped) and len(stripped) < 120:
            topic_candidates.append(re.sub(r'^\d+\.?\d*\s+', '', stripped))

    # Extract bullet points as sub-topics
    bullet_candidates = []
    for line in lines:
        stripped = line.strip()
        if stripped.startswith('- ') and len(stripped) > 10 and len(stripped) < 120:
            bullet_text = stripped[2:].strip()
            # Skip generic bullets
            if not any(skip in bullet_text.lower() for skip in
                       ('trainer', 'attendance', 'lunch', 'break', 'phone', 'assessment',
                        'submit', 'note:', 'create multiple', 'each major')):
                bullet_candidates.append(bullet_text)

    seen = set()
    unique_topics = []
    for t in topic_candidates:
        normalized = t.lower().strip()
        if normalized in seen or normalized in ('introduction', 'conclusion', 'summary',
                                                  'references', 'appendix', 'course description',
                                                  'assessment methods', 'digital attendance',
                                                  'ground rules', 'lesson plan', 'course outline'):
            continue
        seen.add(normalized)
        unique_topics.append(t)

    unique_bullets = []
    for b in bullet_candidates:
        normalized = b.lower().strip()
        if normalized not in seen and len(normalized) > 15:
            seen.add(normalized)
            unique_bullets.append(b)

    # Generate diverse query types for best content quality
    query_templates = [
        "{topic} comprehensive guide with examples and best practices",
        "{topic} real world applications case studies and industry examples",
        "{topic} key concepts definitions explained with diagrams",
        "{topic} latest developments trends and standards 2024 2025",
        "{topic} practical exercises activities for training workshop",
        "{topic} comparison framework methodology approaches",
    ]

    template_idx = 0
    # Add queries for main topics
    for topic in unique_topics[:num_queries]:
        template = query_templates[template_idx % len(query_templates)]
        queries.append(template.format(topic=topic))
        template_idx += 1

    # Add queries for sub-topics if we need more
    for bullet in unique_bullets:
        if len(queries) >= num_queries:
            break
        template = query_templates[template_idx % len(query_templates)]
        queries.append(template.format(topic=bullet))
        template_idx += 1

    # Fill remaining with course-level queries
    if len(queries) < num_queries:
        queries.append(f"{course_title} professional training content best practices examples")
    if len(queries) < num_queries:
        queries.append(f"{course_title} engaging activities exercises for classroom learners")

    return queries[:num_queries]


def _build_platform_urls(topics: List[str]) -> List[Dict[str, str]]:
    """Build direct Wikipedia article URLs for the given topics.

    Uses direct article URLs (not search pages) for reliability — search pages
    often fail as NotebookLM sources because they redirect.
    """
    urls = []
    for topic in topics[:5]:
        clean_topic = re.sub(r'[^\w\s\-]', '', topic).strip()
        if not clean_topic or len(clean_topic) < 4:
            continue
        # Direct article URL — replace spaces with underscores for Wikipedia
        wiki_title = clean_topic.replace(' ', '_')
        urls.append({
            "url": f"https://en.wikipedia.org/wiki/{wiki_title}",
            "title": f"Wikipedia: {clean_topic[:60]}"
        })
    return urls


async def _add_platform_sources(client, notebook_id: str, topics: List[str],
                                 progress_callback=None) -> List[str]:
    """Add sources from Wikipedia to the notebook."""
    platform_urls = _build_platform_urls(topics)
    added_source_ids = []

    if not platform_urls:
        return added_source_ids

    if progress_callback:
        progress_callback(
            f"Adding {len(platform_urls)} Wikipedia sources...",
            30
        )

    for url_info in platform_urls:
        try:
            source = await client.sources.add_url(
                notebook_id, url_info["url"], wait=False
            )
            if source and hasattr(source, 'id'):
                added_source_ids.append(source.id)
                logger.info(f"Added platform source: {url_info['title']} -> {source.id}")
        except Exception as e:
            logger.info(f"Skipped platform source {url_info['url']}: {e}")

    return added_source_ids


async def _do_internet_research(client, notebook_id: str, queries: List[str],
                                 progress_callback=None) -> List[str]:
    """Perform web research using NotebookLM's Research API and import sources.

    Runs all queries concurrently for speed. Each query can yield up to 5 sources.
    """
    all_imported_source_ids = []
    total_queries = len(queries)

    if progress_callback:
        progress_callback(
            f"Researching {total_queries} topics concurrently...",
            35
        )

    async def _run_single_query(query: str, idx: int) -> List[str]:
        """Run a single research query and return imported source IDs."""
        imported_ids = []
        try:
            logger.info(f"Research query {idx+1}/{total_queries}: {query}")
            task = await client.research.start(notebook_id, query, source="web", mode="fast")
            task_id = task.get("task_id") or task.get("report_id", "")

            poll_timeout = 15
            elapsed = 0
            poll_interval = 2
            research_result = None

            while elapsed < poll_timeout:
                await asyncio.sleep(poll_interval)
                elapsed += poll_interval
                result = await client.research.poll(notebook_id)
                status = result.get("status", "")
                if status == "completed":
                    research_result = result
                    break
                elif status == "no_research":
                    break

            if not research_result:
                return imported_ids

            found_sources = research_result.get("sources", [])
            sources_to_import = [s for s in found_sources if s.get("url")][:4]

            if sources_to_import:
                try:
                    imported = await client.research.import_sources(
                        notebook_id, task_id, sources_to_import
                    )
                    for src in imported:
                        src_id = src.get("id") or src.get("source_id", "")
                        if src_id:
                            imported_ids.append(src_id)
                except Exception as e:
                    logger.warning(f"Import failed for query '{query}': {e}")

        except Exception as e:
            logger.warning(f"Research failed for query '{query}': {e}")
        return imported_ids

    # Run all queries concurrently
    query_tasks = [_run_single_query(q, i) for i, q in enumerate(queries)]
    results = await asyncio.gather(*query_tasks, return_exceptions=True)

    for res in results:
        if isinstance(res, list):
            all_imported_source_ids.extend(res)
        elif isinstance(res, Exception):
            logger.warning(f"Research query exception: {res}")

    if progress_callback:
        progress_callback(f"Added {len(all_imported_source_ids)} research sources", 50)

    if all_imported_source_ids:
        try:
            await client.sources.wait_for_sources(
                notebook_id, all_imported_source_ids, timeout=5.0
            )
        except Exception as e:
            logger.warning(f"Some research sources may not be ready: {e}")

    return all_imported_source_ids


async def _generate_slides_direct(content: str, course_title: str, config: Dict[str, Any],
                                   progress_callback=None) -> Dict[str, Any]:
    """
    Generate slides by calling NotebookLM directly.

    Args:
        content: Formatted course text for NotebookLM source
        course_title: Course title for the notebook name
        config: Slide configuration options
        progress_callback: Optional callback to update progress
    """
    try:
        from notebooklm import NotebookLMClient
    except ImportError:
        return {
            "success": False,
            "message": ("**notebooklm-py is not installed.**\n\n"
                        "Run: `pip install notebooklm-py[browser]`\n\n"
                        "Then authenticate: `python -m notebooklm login`")
        }

    enable_research = config.get('enable_research', False)
    num_queries = config.get('num_queries', 0)
    total_steps = 8 if enable_research else 5

    try:
        if progress_callback:
            progress_callback(f"Step 1/{total_steps}: Connecting to NotebookLM...", 5)

        client = await NotebookLMClient.from_storage()

        async with client:
            if progress_callback:
                progress_callback(f"Step 2/{total_steps}: Creating notebook...", 10)

            notebook_title = f"{course_title} - Slides"
            notebook = await client.notebooks.create(notebook_title)
            notebook_id = notebook.id

            if progress_callback:
                progress_callback(f"Step 3/{total_steps}: Uploading course content...", 15)

            source_title = f"{course_title} (Course Material)"
            source_text = content[:50000]
            source = await client.sources.add_text(notebook_id, source_title, source_text)
            source_id = source.id

            if progress_callback:
                progress_callback(f"Step 3/{total_steps}: Adding sources (course + framework + Wikipedia)...", 20)

            # Add course framework overview as 2nd source
            ctx = config.get('_context', {})
            all_source_ids = [source_id]
            research_sources_count = 0
            platform_sources_count = 0

            try:
                ctx_text = _build_course_context_source(ctx, {'content': content})
                if ctx_text and len(ctx_text) > 100:
                    src2 = await client.sources.add_text(
                        notebook_id,
                        f"{course_title} - Course Framework & K&A Statements",
                        ctx_text
                    )
                    all_source_ids.append(src2.id)
            except Exception:
                pass

            # Wait for sources to be ready (URL sources may take longer)
            try:
                await client.sources.wait_for_sources(notebook_id, all_source_ids, timeout=30.0)
            except Exception:
                pass

            if enable_research:
                if progress_callback:
                    progress_callback(f"Step 4/{total_steps}: Extracting research topics...", 25)

                queries = _extract_research_queries(content, course_title, num_queries)

                if queries:
                    if progress_callback:
                        progress_callback(f"Step 4/{total_steps}: Adding Wikipedia sources...", 28)

                    topic_names = []
                    for q in queries:
                        core = re.sub(
                            r'\s+(training|best practices|latest|industry|standards|developments|course content).*$',
                            '', q, flags=re.IGNORECASE
                        ).strip()
                        if core:
                            topic_names.append(core)

                    platform_source_ids = await _add_platform_sources(
                        client, notebook_id, topic_names, progress_callback
                    )
                    all_source_ids.extend(platform_source_ids)
                    platform_sources_count = len(platform_source_ids)

                    research_source_ids = await _do_internet_research(
                        client, notebook_id, queries, progress_callback
                    )
                    all_source_ids.extend(research_source_ids)
                    research_sources_count = len(research_source_ids)

                    if platform_source_ids:
                        try:
                            await client.sources.wait_for_sources(
                                notebook_id, platform_source_ids, timeout=5.0
                            )
                        except Exception:
                            pass

            slide_step = 7 if enable_research else 4
            if progress_callback:
                progress_callback(f"Step {slide_step}/{total_steps}: Generating slide deck...", 65)

            include_notes = config.get('include_notes', False)
            include_summaries = config.get('include_summaries', False)
            slide_style = config.get('slide_style', 'Professional')

            # Determine course duration for slide count target
            ctx = config.get('_context', {})
            raw_hours = (
                ctx.get('Total_Course_Duration_Hours', '')
                or ctx.get('Total_Training_Hours', '')
                or ''
            )
            _hrs_str = str(raw_hours).lower().replace('hours', '').replace('hrs', '').replace('hr', '').replace('h', '').strip()
            import re as _re_d
            _hrs_m = _re_d.search(r'[\d.]+', _hrs_str)
            try:
                course_hours = float(_hrs_m.group()) if _hrs_m else 8
            except (ValueError, TypeError):
                course_hours = 8

            num_topics = sum(
                len(lu.get('Topics', []))
                for lu in ctx.get('Learning_Units', [])
            ) or 10

            # Target: 1 day = 60-100, 2 days = 120-160, 3 days = 180-220
            course_days = max(1, round(course_hours / 8))
            min_slides = course_days * 60
            max_slides = course_days * 100
            course_type = f"{course_hours}-hour"

            instructions = (
                f"Create a comprehensive {slide_style.lower()} slide deck for this WSQ training course. "
                f"This is for PROFESSIONAL CLIENT DELIVERY at Tertiary Infotech — "
                f"quality, depth, and completeness are CRITICAL. "
                f"The client expects slides matching the standard of Dr. Alfred Ang's training materials.\n\n"
                f"TITLE: Use exactly \"{course_title}\" on the cover slide. "
                "Do NOT rephrase, shorten, or change the course title.\n\n"

                "CONTENT RULES:\n"
                "- STRICTLY follow the Course Proposal (CP) source material provided.\n"
                "- All slide content must come from the source.\n"
                "- Do NOT fabricate or invent information not in the source.\n"
                "- EXPAND each topic with detailed explanations, examples, "
                "diagrams, comparisons, and practical applications.\n\n"

                f"SLIDE COUNT: This is a {course_type} course ({course_days} day{'s' if course_days > 1 else ''}). "
                f"Generate between {min_slides} to {max_slides} slides total. "
                f"There are {num_topics} topics across {course_days} day{'s' if course_days > 1 else ''}.\n\n"

                "MANDATORY INTRO SLIDES (in this exact order):\n"
                "1. COVER SLIDE — Course title large and centered, version number, "
                "trainer name, course code. At bottom: ONE single line "
                "'Website: www.tertiaryinfotech.com | Email: enquiry@tertiaryinfotech.com | Tel: +65 6123 4567' "
                "(do NOT repeat this line)\n"
                "2. DIGITAL ATTENDANCE — Title 'Digital Attendance' centered. Body: "
                "'It is mandatory for you to take both AM, PM and Assessment digital attendance for WSQ funded courses. "
                "The trainer or administrator will show you the digital attendance QR code generated from SSG portal. "
                "Please scan the QR code from your mobile phone camera and submit your attendance.'\n"
                "3. ABOUT THE TRAINER — Title centered, space for trainer credentials\n"
                "4. LET'S KNOW EACH OTHER — Title centered, name/role/experience/expectations\n"
                "5. GROUND RULES — Title centered. Bullets: Silent phone, active participation, "
                "respect views, exit silently for breaks, 75% attendance\n"
                "6. LESSON PLAN — ONE slide PER DAY. Title 'Lesson Plan' centered. Format per day:\n"
                "   Day N (bold)\n"
                "   • Digital Attendance (AM) — red color, bold\n"
                "   • Trainer and Learners Introduction (Day 1 only)\n"
                "   • Learning Outcomes (Day 1 only)\n"
                "   • Course Outline (Day 1 only)\n"
                "   • Topic N: [Topic Title] (morning topics)\n"
                "   • Lunch Break\n"
                "   • Digital Attendance (PM) — red color, bold\n"
                "   • Topic N: [Topic Title] (afternoon topics)\n"
                "   • End of Day N\n"
                "   Last day adds before End: Course Feedback and TRAQOM Survey (red, bold), "
                "Digital Attendance (Assessment) (red, bold), Final Assessment & End of Class\n"
                "7. SKILLS FRAMEWORK — Title 'Skills Framework' centered bold 36pt. "
                "Body: TSC Title and TSC Code\n"
                "8. SKILLS FRAMEWORK TSC — Title 'Skills Framework TSC' centered bold 36pt. "
                "Body: All TSC Abilities (A1-AN) and Knowledge (K1-KN) statements listed\n"
                "9. LEARNING OUTCOMES — Title centered. Body: all LOs listed as bullet points\n"
                "10. COURSE OUTLINE — Title 'Course Outline' centered 36pt. Body:\n"
                "    - Each topic in bold 24pt with K&A refs: 'Topic 1: [Title] (K2, K4, A1, A2)'\n"
                "    - Sub-topics indented in bold 18pt below each topic\n"
                "    - May span 2 slides if many topics\n"
                "11. FINAL ASSESSMENT — Title centered. Body: assessment methods bold 24pt, "
                "'Assessment format: Open Book', approved materials note, appeal process\n"
                "12. BRIEFING FOR ASSESSMENT — Title centered. Body: rules as bullet points "
                "(phones away, no photos, no discussion, use black/blue pen, no liquid paper, scripts collected)\n"
                "13. CRITERIA FOR FUNDING — Title centered. Body: "
                "'Minimum attendance rate of 75%', 'Complete assessment and be assessed as Competent'\n\n"

                "TOPIC CONTENT STRUCTURE (MANDATORY for each topic):\n"
                "a) SECTION HEADER SLIDE — Large centered text: 'Topic N' on first line, "
                "topic title on second line. Clean divider layout, no body text.\n"
                "b) CONCEPT SLIDES (10-15 per topic) — MUST use text boxes:\n"
                "   - Title text box at top (centered, 36pt bold)\n"
                "   - Body text box with bullet points (left-aligned, 24pt)\n"
                "   - Each concept gets its own slide with clear title\n"
                "   - Include diagrams, tables, comparison layouts where relevant\n"
                "   - Use two-column layouts for comparisons\n"
                "   - Progressive disclosure from simple to complex\n"
                "c) LAB/ACTIVITY SLIDE — MANDATORY after each topic. Title format: "
                "'Lab - [Topic Name]' or 'Activity: [Topic Name]'. Body includes:\n"
                "   - Clear hands-on activity instructions\n"
                "   - URLs/tools to use if applicable\n"
                "   - 'Submit your result as part of practical test' or "
                "'Submit your analysis as part of theory test'\n"
                "   - This is NOT optional — every topic MUST end with a lab/activity slide\n\n"

                "MANDATORY CLOSING SLIDES (at end, in this exact order):\n"
                "1. SUMMARY & Q&A — Section header: 'Summary & Q&A' centered\n"
                "2. TRAQOM SURVEY — Title 'TRAQOM Survey'. Body: 'Access the survey here' (teal link), "
                "'Key in your Last four NRIC/FIN characters and the Six-digit course run ID to complete the survey'\n"
                "3. CERTIFICATE OF ACCOMPLISHMENT — Two-column layout. Left: 'Please provide the following details to facilitate the issuance of your certificate after the class and to ensure the accuracy of the information printed on it.' Right: Certificate template: 'CERTIFICATE OF ACCOMPLISHMENT', 'This Certificate is proudly presented to [Student Name]', 'For completing the course [Course Name]', 'Held on [Course Dates]', 'Dr. Alfred Ang, Managing Director, Tertiary Infotech Academy'.\n"
                "4. DIGITAL ATTENDANCE — Repeat of attendance instructions\n"
                "5. FINAL ASSESSMENT — Section header: 'Final Assessment' centered\n"
                "6. RECOMMENDED COURSES — Title centered, list related WSQ courses\n"
                "7. SUPPORT — Title 'Support', contact: enquiry@tertiaryinfotech.com, Tel: 61000613\n"
                "8. THANK YOU — Section header: 'Thank You!' centered\n\n"

                "SLIDE DESIGN (follow this exactly):\n"
                "- Every slide MUST use text boxes — title text box + body text box\n"
                "- Title: centered, bold, 36pt (457200 EMU)\n"
                "- Body text: left-aligned, 24pt (304800 EMU), bullet points\n"
                "- Small text where needed: 18pt (228600 EMU)\n"
                "- BACKGROUND: EVERY slide MUST have the SAME solid white background. No exceptions.\n"
                "- HEADINGS: Dark navy blue. ACCENTS: Teal.\n"
                "- BODY TEXT: Dark grey, clean sans-serif font.\n"
                "- SECTION HEADERS: Large centered text on clean background, no body text.\n"
                "- Digital Attendance (AM), Digital Attendance (PM), Digital Attendance (Assessment), "
                "Course Feedback and TRAQOM Survey: make these red color and bold text.\n"
                "- K&A references shown in bold after topic titles: '(K2, K4, A1, A2)'\n"
                "- Slide numbers on every slide (bottom-right corner).\n"
                "- Use icons, tables, flowcharts, two-column layouts for visual variety.\n"
                "- Do NOT use large decorative background images.\n"
            )
            if include_notes:
                instructions += "\nInclude detailed speaker/facilitator notes for every slide."
            if include_summaries:
                instructions += "\nAdd a summary slide at the end of each topic section."
            if enable_research and research_sources_count > 0:
                instructions += (
                    "\nUse research sources to supplement CP content "
                    "with latest industry practices and real-world examples. "
                    "CP content takes priority — research adds depth."
                )
            instructions += (
                f"\n\nREMEMBER: {course_type} course ({course_days} day{'s' if course_days > 1 else ''}). "
                f"Generate {min_slides}-{max_slides} slides. "
                "Quality AND quantity required for professional client delivery."
            )

            # SINGLE attempt — no retry to prevent duplicate slide decks
            gen_result = await client.artifacts.generate_slide_deck(
                notebook_id,
                source_ids=all_source_ids,
                instructions=instructions,
            )
            task_id = gen_result.task_id if hasattr(gen_result, 'task_id') else str(gen_result)

            wait_step = 8 if enable_research else 5
            if progress_callback:
                progress_callback(
                    f"Step {wait_step}/{total_steps}: Waiting for slides (1-3 min)...",
                    80
                )

            generation_status = "triggered"
            if task_id:
                try:
                    await client.artifacts.wait_for_completion(
                        notebook_id, task_id, timeout=300.0
                    )
                    generation_status = "completed"
                except TimeoutError:
                    generation_status = "timeout"
                except Exception as e:
                    generation_status = f"wait_error: {e}"

            if progress_callback:
                progress_callback("Slides generated successfully!", 95)

            return {
                "success": True,
                "message": "Slide deck generated successfully!",
                "notebook_id": notebook_id,
                "notebook_title": notebook_title,
                "task_id": task_id,
                "generation_status": generation_status,
                "research_enabled": enable_research,
                "research_sources_count": research_sources_count,
                "platform_sources_count": platform_sources_count,
                "total_sources": len(all_source_ids),
            }

    except FileNotFoundError:
        return {
            "success": False,
            "message": ("**NotebookLM authentication not found.**\n\n"
                        "Run in your terminal:\n"
                        "```\nuv run notebooklm login\n```")
        }
    except Exception as e:
        error_msg = str(e)
        if "auth" in error_msg.lower() or "login" in error_msg.lower():
            return {
                "success": False,
                "message": (f"**NotebookLM authentication error:** {error_msg}\n\n"
                            "Please re-authenticate:\n"
                            "```\nuv run notebooklm login\n```")
            }
        return {
            "success": False,
            "message": f"**Error generating slides:** {error_msg}"
        }


# =============================================================================
# Browser-based generation fallback (bypasses API rate limits)
# =============================================================================

async def _generate_via_browser(notebook_id: str, label: str, progress_callback=None) -> str:
    """Trigger slide generation via browser UI when API is rate limited.
    Opens notebook in visible browser, clicks 'Slide Deck', waits for completion.
    Returns generation_status string.
    """
    from pathlib import Path as _Path
    _profile_dir = _Path.home() / ".notebooklm" / "browser_profile"
    _profile_dir.mkdir(parents=True, exist_ok=True)

    try:
        from playwright.async_api import async_playwright
        async with async_playwright() as pw:
            ctx = await pw.chromium.launch_persistent_context(
                user_data_dir=str(_profile_dir),
                headless=False,
                channel="chrome",
                accept_downloads=True,
                args=[
                    "--disable-blink-features=AutomationControlled",
                    "--no-first-run",
                    "--no-default-browser-check",
                ],
                ignore_default_args=["--enable-automation"],
            )
            pg = ctx.pages[0] if ctx.pages else await ctx.new_page()
            nb_url = f"https://notebooklm.google.com/notebook/{notebook_id}"

            if progress_callback:
                progress_callback(f"[{label}] Opening notebook in browser...", None)
            await pg.goto(nb_url, timeout=60000, wait_until="domcontentloaded")
            await asyncio.sleep(5)

            # Check if redirected to login
            if "accounts.google.com" in pg.url:
                logger.warning(f"[{label}] Browser redirected to login")
                await ctx.close()
                return "failed: browser_auth_expired"

            # Check if we're on the notebook page (not redirected to homepage)
            if notebook_id not in pg.url:
                logger.warning(f"[{label}] Redirected to {pg.url} instead of notebook")
                # Try navigating directly again
                await pg.goto(nb_url, timeout=60000, wait_until="domcontentloaded")
                await asyncio.sleep(5)

            # First delete any failed slide deck
            if progress_callback:
                progress_callback(f"[{label}] Checking for failed slide decks...", None)
            try:
                delete_btns = await pg.query_selector_all("text=Delete")
                for btn in delete_btns:
                    try:
                        parent = await btn.evaluate("el => el.closest('[class]')?.textContent || ''")
                        if "slide" in parent.lower() or "deck" in parent.lower():
                            await btn.click()
                            await asyncio.sleep(2)
                            logger.info(f"[{label}] Deleted failed slide deck")
                            break
                    except Exception:
                        pass
            except Exception:
                pass

            # Click 'Slide deck' button to trigger generation
            if progress_callback:
                progress_callback(f"[{label}] Clicking 'Slide deck' to generate...", None)
            try:
                # Try multiple selectors for the Slide Deck button
                slide_btn = None
                for selector in ["text=Slide deck", "text=Slide Deck", "text='Slide deck'"]:
                    slide_btn = await pg.query_selector(selector)
                    if slide_btn:
                        break
                if not slide_btn:
                    # Try finding by aria label or button with slide text
                    all_btns = await pg.query_selector_all("button, [role='button'], [class*='card']")
                    for btn in all_btns:
                        txt = (await btn.text_content() or "").lower()
                        if "slide" in txt and "deck" in txt:
                            slide_btn = btn
                            break

                if slide_btn:
                    await slide_btn.click()
                    await asyncio.sleep(3)
                    if progress_callback:
                        progress_callback(f"[{label}] Generation triggered!", None)
                else:
                    logger.warning(f"[{label}] Could not find 'Slide deck' button")
                    await ctx.close()
                    return "failed: no_slide_deck_button"
            except Exception as e:
                logger.warning(f"[{label}] Click failed: {e}")
                await ctx.close()
                return f"failed: click_error: {e}"

            # Wait for generation to complete (poll every 5s, up to 10 min)
            if progress_callback:
                progress_callback(f"[{label}] Waiting for slide generation (up to 10 min)...", None)
            for poll_i in range(120):
                try:
                    page_text = await pg.text_content("body") or ""
                    if "Generating" in page_text and ("Slide" in page_text or "slide" in page_text):
                        pass  # Still generating
                    elif poll_i > 3:  # Give it a few polls before checking completion
                        # Check for completed or failed indicators
                        if "generation" not in page_text.lower() or "failed" in page_text.lower():
                            # Check if there's a completed slide deck
                            error_els = await pg.query_selector_all("text=Slide deck generation")
                            generating_els = await pg.query_selector_all("text=Generating Slide")
                            if not generating_els and not error_els:
                                if progress_callback:
                                    progress_callback(f"[{label}] Slide deck ready! ({poll_i * 5}s)", None)
                                await ctx.close()
                                return "completed"
                except Exception:
                    pass
                if poll_i > 0 and poll_i % 12 == 0 and progress_callback:
                    progress_callback(f"[{label}] Still generating... ({poll_i * 5}s)", None)
                await asyncio.sleep(5)

            await ctx.close()
            return "incomplete: browser_timeout"
    except Exception as e:
        logger.warning(f"[{label}] Browser generation failed: {e}")
        return f"failed: browser_error: {e}"


# =============================================================================
# Module-level chunk deck generator (used by both single & multi-account modes)
# =============================================================================

async def _generate_chunk_deck_impl(client, cm: dict, notebook_id: str,
                                     nb_title: str, source_id,
                                     course_title: str, config: dict,
                                     progress_callback=None) -> dict:
    """Generate slides for a single chunk using the given NotebookLM client.

    source_id: either a single source ID string or a list of source IDs.
    """
    label = cm['label']
    enable_research = config.get('enable_research', False)
    num_queries = config.get('num_queries', 0)
    slide_style = config.get('slide_style', 'Professional')
    include_notes = config.get('include_notes', False)

    # Accept single ID or list of IDs
    if isinstance(source_id, list):
        all_source_ids = list(source_id)
    else:
        all_source_ids = [source_id]
    research_count = 0

    if enable_research:
        if progress_callback:
            progress_callback(f"[{label}] Researching {num_queries} topics...", None)
        try:
            queries = _extract_research_queries(cm['content'], course_title, num_queries)
            if queries:
                research_ids = await _do_internet_research(
                    client, notebook_id, queries
                )
                all_source_ids.extend(research_ids)
                research_count = len(research_ids)
                if progress_callback:
                    progress_callback(f"[{label}] Added {research_count} research sources", None)
        except Exception as e:
            logger.warning(f"{label} research failed: {e}")

    # Build instructions — matching supervisor's reference PPTX format
    # Target: 1 day = 60-100, 2 days = 120-160, 3 days = 180-220
    ctx = config.get('_context', {})
    raw_hrs = (
        ctx.get('Total_Course_Duration_Hours', '')
        or ctx.get('Total_Training_Hours', '')
        or ''
    )
    _hrs_str2 = str(raw_hrs).lower().replace('hours', '').replace('hrs', '').replace('hr', '').replace('h', '').strip()
    import re as _re_c
    _hrs_m2 = _re_c.search(r'[\d.]+', _hrs_str2)
    try:
        course_hours = float(_hrs_m2.group()) if _hrs_m2 else 8
    except (ValueError, TypeError):
        course_hours = 8
    course_days = max(1, round(course_hours / 8))
    # Target: 1 day = 60-100 pages, 2 day = 120-160 pages
    course_min_slides = course_days * 60
    course_max_slides = course_days * 100 if course_days == 1 else course_days * 80

    # Calculate this chunk's share: proportional to its topic count
    total_topics_in_course = sum(
        len(lu.get('Topics', [])) for lu in ctx.get('Learning_Units', [])
    ) or 1
    my_topic_count = cm.get('num_topics', len(cm['topic_names']))
    # Use midpoint-to-upper range for chunk allocation
    course_target = (course_min_slides + course_max_slides) // 2
    content_slides = max(15 * my_topic_count, round(
        (course_target - 13 - 8) * my_topic_count / total_topics_in_course
    ))
    # Add intro/closing overhead
    slides_target = content_slides
    if cm['is_first']:
        slides_target += 13  # intro pages
    if cm['is_last']:
        slides_target += 8   # closing pages
    topic_list = ", ".join(cm['topic_names'])

    # ── Build instructions for NotebookLM slide generation ──
    # NOTE: NotebookLM has a character limit on instructions (~3000 chars).
    # Keep instructions concise but comprehensive.
    instructions = (
        f"Create a professional slide deck for WSQ training: \"{course_title}\". "
        f"Target ~{slides_target} slides. Topics: {topic_list}.\n\n"

        "DESIGN: White background, navy blue headings (Arial bold centered), dark grey body text. "
        "Two-column layout: detailed bullet points LEFT, colorful realistic image RIGHT. "
        "NO monochrome blue icons — use vibrant photos, diagrams, infographics. "
        "Footer on every slide: 'This material belongs to Tertiary Infotech Academy Pte Ltd (UEN: 20120096W). All Rights Reserved'. "
        "Slide numbers bottom-right. K&A references bold after topic titles. "
        "Digital Attendance, Course Feedback, TRAQOM Survey in RED BOLD text. No logos on any slide.\n\n"

        "CONTENT: Focus on VISUAL slides — this deck provides the images and diagrams. "
        "For EVERY concept, create slides with:\n"
        "- DIAGRAMS: flowcharts, process flows, architecture diagrams\n"
        "- INFOGRAPHICS: comparison charts, data visualizations, statistics\n"
        "- TABLES: feature comparisons, framework summaries, pros/cons\n"
        "- IMAGES: realistic photos, icons, illustrations related to each concept\n"
        "Every slide MUST have a visual element (image, diagram, chart, or table). "
        "Use vibrant colors, professional graphics, real photographs. "
        "Include brief text with the visuals but prioritize the VISUAL elements.\n\n"
    )

    if cm['is_first']:
        instructions += (
            "INTRO SLIDES (this exact order): "
            f"1. COVER: '{course_title}' centered, Version 1.0, "
            f"Course Code: {ctx.get('TGS_Ref_No', 'TGS-XXXXXXXX')}, 'Trainer:' only (no name/brackets). "
            "No logos. Bottom: www.tertiaryinfotech.com | enquiry@tertiaryinfotech.com | +65 6123 4567. "
            "2. DIGITAL ATTENDANCE: mandatory AM/PM/Assessment attendance. "
            "3. ABOUT TRAINER: title + professional illustration only, NO text/names/credentials. "
            "4. LET'S KNOW EACH OTHER: name/role/experience/expectations. "
            "5. GROUND RULES: silent phone, participate, respect, punctual, 75% attendance. "
            "6. LESSON PLAN: one slide per day. "
            "7-8. SKILLS FRAMEWORK: TSC Title/Code, K&A statements. "
            "9. LEARNING OUTCOMES. 10-11. COURSE OUTLINE with K&A refs. "
            "12. FINAL ASSESSMENT: methods, Open Book. "
            "13. BRIEFING FOR ASSESSMENT. 14. CRITERIA FOR FUNDING.\n\n"
        )

    if cm['is_last']:
        instructions += (
            "CLOSING SLIDES: "
            "1. SUMMARY & Q&A. 2. TRAQOM SURVEY. 3. CERTIFICATE OF ACCOMPLISHMENT. "
            "4. DIGITAL ATTENDANCE. 5. FINAL ASSESSMENT. 6. RECOMMENDED COURSES. "
            "7. SUPPORT: enquiry@tertiaryinfotech.com, Tel: 61000613. 8. THANK YOU.\n\n"
        )

    instructions += (
        "PER TOPIC: a) Section header with topic title. "
        "b) 5-8 VISUAL slides per topic — EVERY slide must have a diagram, flowchart, "
        "image, table, chart, or infographic. Create: "
        "1) Overview diagram — visual map of the concept, "
        "2) Process flowchart — how it works step by step, "
        "3) Comparison table — key types or approaches, "
        "4) Framework diagram — model or architecture visual, "
        "5) Real-world image — photo or illustration of application. "
        "c) Activity slide after each topic.\n"
    )

    if include_notes:
        instructions += "\nInclude detailed speaker/facilitator notes for every slide."

    # Clean up any failed/error sources before generating slides
    # Failed URL sources can cause generation to fail
    try:
        sources_list = await client.sources.list(notebook_id)
        for src in sources_list:
            src_id = getattr(src, 'id', None) or (src.get('id') if isinstance(src, dict) else None)
            status = getattr(src, 'status', None) or (src.get('status') if isinstance(src, dict) else None)
            status_str = str(status).lower() if status else ''
            if status_str in ('error', 'failed', 'invalid', 'unprocessable'):
                try:
                    await client.sources.delete(notebook_id, src_id)
                    if src_id in all_source_ids:
                        all_source_ids.remove(src_id)
                    logger.info(f"[{label}] Removed failed source {src_id} (status: {status_str})")
                except Exception:
                    pass
    except Exception as e:
        logger.warning(f"[{label}] Source cleanup check failed: {e}")

    if progress_callback:
        progress_callback(f"[{label}] Sending slide generation request...", None)

    # Try API generation first; if rate limited, flag for browser-based generation
    # Truncate instructions to stay within NotebookLM's limit (~3000 chars)
    if len(instructions) > 3000:
        instructions = instructions[:2950] + "\n"
        logger.info(f"[{label}] Instructions truncated to {len(instructions)} chars")

    # Use DETAILED_DECK format + DEFAULT length for visual slides
    # NotebookLM provides images/diagrams/flowcharts; Claude AI provides the deep content
    # DEFAULT length is faster than LONG while still producing good visual slides
    try:
        from notebooklm import SlideDeckFormat, SlideDeckLength
        _slide_format = SlideDeckFormat.DETAILED_DECK
        _slide_length = SlideDeckLength.DEFAULT
    except (ImportError, AttributeError):
        _slide_format = None
        class _DefaultLength:
            value = 1
        _slide_length = _DefaultLength()

    rate_limited = False
    generation_status = "pending"
    try:
        gen_kwargs = dict(
            notebook_id=notebook_id,
            source_ids=all_source_ids,
            instructions=instructions,
        )
        if _slide_format is not None:
            gen_kwargs['slide_format'] = _slide_format
        gen_kwargs['slide_length'] = _slide_length
        gen_result = await client.artifacts.generate_slide_deck(**gen_kwargs)
        status_val = getattr(gen_result, 'status', '') or ''
        logger.info(f"[{label}] API generation: status={status_val!r}")
        # Fix: Check rate limit with OR — either flag can indicate rate limiting
        if getattr(gen_result, 'is_rate_limited', False) or (
            getattr(gen_result, 'is_failed', False) and 'rate' in str(getattr(gen_result, 'error', '')).lower()
        ):
            rate_limited = True
        elif status_val in ('completed', 'in_progress', 'pending'):
            generation_status = "api_triggered"
    except Exception as e:
        error_str = str(e).lower()
        if 'rate' in error_str or 'limit' in error_str or 'quota' in error_str or '429' in error_str:
            rate_limited = True
        else:
            # Retry once with minimal instructions on any other failure
            logger.warning(f"[{label}] First attempt failed: {e}. Retrying with shorter instructions...")
            if progress_callback:
                progress_callback(f"[{label}] Retrying generation...", None)
            try:
                short_instructions = (
                    f"Create a professional training slide deck for \"{course_title}\". "
                    f"Target ~{slides_target} slides. White background, navy headings, two-column layout. "
                    "Detailed content with definitions, examples, processes. Activity after each topic. "
                    "Footer: 'This material belongs to Tertiary Infotech Academy Pte Ltd (UEN: 20120096W). All Rights Reserved'"
                )
                gen_kwargs2 = dict(
                    notebook_id=notebook_id,
                    source_ids=all_source_ids,
                    instructions=short_instructions,
                )
                if _slide_format is not None:
                    gen_kwargs2['slide_format'] = _slide_format
                gen_kwargs2['slide_length'] = _slide_length
                gen_result = await client.artifacts.generate_slide_deck(**gen_kwargs2)
                status_val = getattr(gen_result, 'status', '') or ''
                logger.info(f"[{label}] Retry generation: status={status_val!r}")
                if status_val in ('completed', 'in_progress', 'pending'):
                    generation_status = "api_triggered"
            except Exception as e2:
                logger.error(f"[{label}] Retry also failed: {e2}")
                raise

    # If API worked, poll for completion (max 4 minutes per deck)
    if generation_status == "api_triggered":
        if progress_callback:
            progress_callback(f"[{label}] API generation triggered — waiting for completion...", None)
        for poll_i in range(120):
            await asyncio.sleep(2)
            try:
                decks = await client.artifacts.list_slide_decks(notebook_id)
                if decks:
                    for deck in decks:
                        # Status can be int (3=COMPLETED, 4=FAILED) or string
                        raw_status = getattr(deck, 'status', '')
                        ds = str(raw_status or '').lower()
                        is_failed = getattr(deck, 'is_failed', False)
                        is_completed = getattr(deck, 'is_completed', False)
                        is_generating = getattr(deck, 'is_generating', None)
                        # Check for actual failure first
                        if is_failed or raw_status == 4 or ds in ('failed', 'error', '4'):
                            logger.warning(f"[{label}] Deck generation failed: status={raw_status}")
                            generation_status = f"failed: {ds}"
                            break
                        # Check for completion — handle int status 3, string "completed", or bool flags
                        if (raw_status == 3 or ds == 'completed' or ds == '3'
                                or is_completed
                                or (is_generating is not None and not is_generating and not is_failed)):
                            generation_status = "completed"
                            break
                    if generation_status == "completed" or generation_status.startswith("failed"):
                        break
            except Exception:
                pass
            if poll_i > 0 and poll_i % 10 == 0:
                print(f"[SLIDES] [{label}] Polling... ({poll_i * 2}s) status={generation_status}")
                if progress_callback:
                    progress_callback(f"[{label}] Still generating... ({poll_i * 2}s)", None)

    # If API rate limited, use BROWSER to trigger generation (bypasses API limits)
    if rate_limited:
        if progress_callback:
            progress_callback(f"[{label}] API rate limited — using browser to generate...", None)
        generation_status = await _generate_via_browser(notebook_id, label, progress_callback)

    # Download as PPTX — try direct PPTX first, fall back to PDF→PPTX
    pptx_path = None
    original_pdf_path = None  # Keep PDF for editable conversion
    print(f"[SLIDES] [{label}] Download phase: generation_status={generation_status}")
    if generation_status == "completed":
        try:
            import tempfile
            pptx_path = tempfile.mktemp(suffix=f"_{label}.pptx")

            # Step 1: Try direct PPTX download from NotebookLM
            direct_ok = await _download_pptx_direct(client, notebook_id, pptx_path)
            print(f"[SLIDES] [{label}] Direct PPTX download: {direct_ok}")

            if direct_ok:
                if progress_callback:
                    progress_callback(f"[{label}] Direct PPTX downloaded!", None)
            else:
                # Step 2: Fallback — download PDF, convert to PPTX
                if progress_callback:
                    progress_callback(f"[{label}] Downloading PDF, converting to PPTX...", None)
                pdf_path = tempfile.mktemp(suffix=f"_{label}.pdf")
                await client.artifacts.download_slide_deck(notebook_id, pdf_path)
                _pdf_size = os.path.getsize(pdf_path) // 1024 if os.path.exists(pdf_path) else 0
                print(f"[SLIDES] [{label}] PDF downloaded: {_pdf_size}KB at {pdf_path}")
                _pdf_to_pptx(pdf_path, pptx_path)
                # Keep PDF for editable conversion (Claude Agent SDK extracts text from images)
                original_pdf_path = pdf_path

            # Stamp logos and replace certificate on PPTX
            if pptx_path and os.path.exists(pptx_path):
                _stamp_logos_on_pptx(pptx_path)
                _replace_certificate_in_pptx(pptx_path)
                if progress_callback:
                    progress_callback(f"[{label}] PPTX ready with logos!", None)

        except Exception as e:
            print(f"[SLIDES] [{label}] DOWNLOAD FAILED: {type(e).__name__}: {e}")
            if progress_callback:
                progress_callback(f"[{label}] Download failed: {e}", None)
            pptx_path = None
    else:
        print(f"[SLIDES] [{label}] SKIPPING DOWNLOAD: status={generation_status} (not completed)")

    # Count slides in the downloaded PPTX
    _slide_count = 0
    if pptx_path and os.path.exists(pptx_path):
        try:
            from pptx import Presentation as _CountPrs
            _cprs = _CountPrs(pptx_path)
            _slide_count = len(_cprs.slides)
        except Exception:
            pass
    _pptx_size = os.path.getsize(pptx_path) // 1024 if pptx_path and os.path.exists(pptx_path) else 0

    print(f"[SLIDES] [{label}] DONE: status={generation_status}, slides={_slide_count}, size={_pptx_size}KB, path={pptx_path}")
    if progress_callback:
        progress_callback(f"[{label}] Done ({generation_status}, {_slide_count} slides)", None)

    return {
        "lu_num": cm['label'],
        "lu_num_base": cm['lu_num'],  # Base LU number (e.g., "LU1") for aggregation
        "lu_title": f"{cm['lu_title']} ({cm['topic_range']})",
        "notebook_id": notebook_id,
        "notebook_title": nb_title,
        "task_id": "",
        "generation_status": generation_status,
        "research_sources_count": research_count,
        "total_sources": len(all_source_ids),
        "chunk_idx": cm['chunk_idx'],
        "original_pdf_path": original_pdf_path,  # Original NotebookLM PDF (for editable conversion)
        "pdf_path": pptx_path,  # Legacy: points to PPTX
        "pptx_path": pptx_path,
        "slide_count": _slide_count,
    }


# =============================================================================
# Multi-account batch runner
# =============================================================================

async def _run_account_batch(account, chunk_meta_list: list,
                              course_title: str, config: dict,
                              progress_callback=None, context: dict = None) -> list:
    """Run a batch of deck generations on a single NotebookLM account.

    Creates one client per account, generates up to N decks concurrently.
    Returns list of result dicts.
    """
    from notebooklm import NotebookLMClient

    results = []
    try:
        client = await NotebookLMClient.from_storage(
            path=str(account.storage_state_path)
        )
        acct_name = account.email.split("@")[0]
        async with client:
            # Phase 1: Create notebooks
            if progress_callback:
                progress_callback(f"[{acct_name}] Creating {len(chunk_meta_list)} notebooks...", None)

            async def _create_nb(cm):
                nb_title = f"{course_title} - {cm['label']}: {cm['lu_title']} ({cm['topic_range']})"
                notebook = await client.notebooks.create(nb_title)
                return (cm, notebook.id, nb_title)

            nb_tasks = [_create_nb(cm) for cm in chunk_meta_list]
            nb_results = await asyncio.gather(*nb_tasks, return_exceptions=True)

            active = []
            for res in nb_results:
                if isinstance(res, Exception):
                    logger.warning(f"[{acct_name}] Notebook creation failed: {res}")
                    account.decks_failed += 1
                    continue
                active.append(res)

            if not active:
                account.error = "Failed to create any notebooks"
                return results

            # Phase 2: Add multiple sources per notebook (course material + framework + Wikipedia)
            ctx = context or config.get('_context', {})
            if progress_callback:
                progress_callback(f"[{acct_name}] Adding sources to {len(active)} notebooks (3-5 per deck)...", None)

            async def _add_multi(cm, nb_id):
                return await _add_multi_sources(client, nb_id, cm, course_title, ctx)

            src_tasks = [_add_multi(cm, nb_id) for (cm, nb_id, _) in active]
            src_results = await asyncio.gather(*src_tasks, return_exceptions=True)

            # Phase 2b: Wait for ALL sources to be ready (prevents generation failures)
            # URL sources take longer to process than text sources — use 60s timeout
            all_wait_tasks = []
            for i, (cm, nb_id, nb_title) in enumerate(active):
                src_ids = src_results[i] if not isinstance(src_results[i], Exception) else []
                if src_ids:
                    all_wait_tasks.append(client.sources.wait_for_sources(nb_id, src_ids, timeout=15.0))
            if all_wait_tasks:
                await asyncio.gather(*all_wait_tasks, return_exceptions=True)

            # Phase 3: Generate slide decks IN PARALLEL
            if progress_callback:
                progress_callback(f"[{acct_name}] Generating {len(active)} slide decks in parallel...", None)

            # Build valid (cm, nb_id, nb_title, src_ids) tuples
            gen_inputs = []
            for i, (cm, nb_id, nb_title) in enumerate(active):
                src_ids = src_results[i] if not isinstance(src_results[i], Exception) else []
                if not src_ids:
                    logger.warning(f"[{acct_name}] Skipping {cm['label']} — no sources added")
                    account.decks_failed += 1
                    continue
                gen_inputs.append((cm, nb_id, nb_title, src_ids))

            # Generate decks SEQUENTIALLY with brief delays
            gen_results = []
            for gi, (cm, nb_id, nb_title, src_ids) in enumerate(gen_inputs):
                if gi > 0:
                    # Brief pause between generation requests
                    if progress_callback:
                        progress_callback(f"[{acct_name}] Starting next deck...", None)
                    await asyncio.sleep(3)
                try:
                    res = await _generate_chunk_deck_impl(
                        client, cm, nb_id, nb_title, src_ids, course_title, config,
                        progress_callback=progress_callback
                    )
                    gen_results.append(res)
                except Exception as e:
                    gen_results.append(e)

            hit_rate_limit = False
            error_details = []
            for i, gen_res in enumerate(gen_results):
                cm = gen_inputs[i][0]
                if isinstance(gen_res, Exception):
                    error_str = str(gen_res).lower()
                    if 'rate' in error_str or 'limit' in error_str or 'quota' in error_str or '429' in error_str:
                        logger.warning(f"[{acct_name}] Rate limit on {cm['label']}: {gen_res}")
                        account.decks_failed += 1
                        hit_rate_limit = True
                        error_details.append(f"{cm['label']}: rate_limited")
                    else:
                        logger.warning(f"[{acct_name}] Generation failed for {cm['label']}: {gen_res}")
                        account.decks_failed += 1
                        error_details.append(f"{cm['label']}: {type(gen_res).__name__}: {gen_res}")
                    continue

                status = gen_res.get("generation_status", "")
                if status == "rate_limited" or "rate" in status.lower():
                    logger.warning(f"[{acct_name}] Rate limited on {cm['label']}")
                    account.decks_failed += 1
                    hit_rate_limit = True
                    error_details.append(f"{cm['label']}: rate_limited")
                    if progress_callback:
                        progress_callback(f"[{acct_name}] {cm['label']} rate limited", None)
                elif status.startswith("failed") or status.startswith("incomplete"):
                    logger.warning(f"[{acct_name}] Failed: {cm['label']}: {status}")
                    account.decks_failed += 1
                    error_details.append(f"{cm['label']}: {status}")
                else:
                    results.append(gen_res)
                    account.decks_completed += 1
                    if progress_callback:
                        has_pptx = "PPTX ready" if gen_res.get("pptx_path") else "no PPTX"
                        progress_callback(f"[{acct_name}] {gen_res['lu_num']} completed ({has_pptx})", None)

            if error_details:
                account.error = "; ".join(error_details)
            elif hit_rate_limit:
                account.error = "rate_limited"

    except FileNotFoundError as e:
        account.error = f"Auth file not found: {account.storage_state_path}"
        logger.error(f"[{account.email}] Storage state file missing: {e}")
        if progress_callback:
            progress_callback(f"[{account.email.split('@')[0]}] Auth expired or missing — run 'uv run notebooklm login'", None)
    except Exception as e:
        account.error = f"{type(e).__name__}: {e}"
        logger.error(f"[{account.email}] Account batch failed: {e}")
        if progress_callback:
            progress_callback(f"[{account.email.split('@')[0]}] Error: {type(e).__name__}: {e}", None)

    return results


# =============================================================================
# Multi-account orchestrator
# =============================================================================

async def _generate_slides_multi_account(context: dict, course_title: str,
                                          config: Dict[str, Any],
                                          progress_callback=None,
                                          skip_lu_indices: set = None) -> Dict[str, Any]:
    """Generate slides across multiple NotebookLM accounts.

    Distributes chunks across accounts (max N per account) and runs
    all accounts concurrently. Falls back to single-account if no
    multi-account config is available.
    """
    try:
        from notebooklm import NotebookLMClient
    except ImportError:
        return {
            "success": False,
            "message": ("**notebooklm-py is not installed.**\n\n"
                        "Run: `pip install notebooklm-py[browser]`")
        }

    from generate_slides.account_pool import AccountPool, AccountInfo

    lus = context.get('Learning_Units', [])
    num_lus = len(lus)
    if num_lus == 0:
        return {"success": False, "message": "No Learning Units found in course info."}

    # ── Calculate target decks (same logic as _generate_slides_per_lu) ──
    total_hours = (
        context.get('Total_Course_Duration_Hours', '')
        or context.get('Total_Training_Hours', '')
        or ''
    )
    _hrs_raw = str(total_hours).lower().replace('hours', '').replace('hrs', '').replace('hr', '').replace('h', '').strip()
    _hrs_match = re.search(r'[\d.]+', _hrs_raw)
    try:
        hours = float(_hrs_match.group()) if _hrs_match else 8
    except (ValueError, TypeError):
        hours = 8
    days = max(1, round(hours / 8))
    total_topics = sum(len(lu.get('Topics', [])) for lu in lus)

    # ── 1 deck per TOPIC — each topic gets its own NotebookLM notebook ──
    # This generates ~20-30 NotebookLM slides per topic.
    # For 10 topics: 200-300 slides total → plenty for 100+ slide target.
    lu_indices_with_topics = [i for i in range(num_lus) if len(lus[i].get('Topics', [])) > 0]

    all_chunks = []
    chunk_idx = 0
    for lu_idx in lu_indices_with_topics:
        lu = lus[lu_idx]
        lu_num = lu.get('LU_Number', f'LU{lu_idx + 1}')
        lu_title = lu.get('LU_Title', 'Learning Unit')
        topics = lu.get('Topics', [])
        for t_idx, topic in enumerate(topics):
            t_title = topic.get('Topic_Title', f'Topic {t_idx + 1}')
            label = f"{lu_num}_T{t_idx + 1}"
            all_chunks.append((chunk_idx, lu_idx, lu_num, lu_title, [t_idx], label))
            chunk_idx += 1

    total_chunks = len(all_chunks)
    if total_chunks == 0:
        return {"success": False, "message": "No topics found in any Learning Unit."}

    if skip_lu_indices is None:
        skip_lu_indices = set()

    # ── Build chunk metadata ──
    # is_first/is_last based on FULL course order, NOT just selected decks
    # This ensures intro slides only appear on the actual first deck
    # and closing slides only on the actual last deck of the course
    first_chunk_idx = all_chunks[0][0] if all_chunks else -1
    last_chunk_idx = all_chunks[-1][0] if all_chunks else -1

    chunk_meta = []
    for chunk in all_chunks:
        c_idx, lu_idx, lu_num, lu_title, t_indices, label = chunk
        if c_idx in skip_lu_indices:
            continue

        is_first = (c_idx == first_chunk_idx)
        is_last = (c_idx == last_chunk_idx)

        content = _format_chunk_source_text(
            context, lu_idx, num_lus, t_indices, label,
            is_first_chunk_of_course=is_first,
            is_last_chunk_of_course=is_last
        )

        topics = lus[lu_idx].get('Topics', [])
        topic_names = [topics[ti].get('Topic_Title', f'T{ti+1}') for ti in t_indices if ti < len(topics)]
        topic_range = f"T{t_indices[0]+1}-T{t_indices[-1]+1}" if len(t_indices) > 1 else f"T{t_indices[0]+1}"

        chunk_meta.append({
            'chunk_idx': c_idx,
            'lu_idx': lu_idx,
            'lu_num': lu_num,
            'lu_title': lu_title,
            'label': label,
            'topic_indices': t_indices,
            'topic_range': topic_range,
            'topic_names': topic_names,
            'num_topics': len(t_indices),
            'content': content,
            'is_first': is_first,
            'is_last': is_last,
        })

    if not chunk_meta:
        return {"success": True, "message": "All decks already completed.", "lu_results": [], "num_lus": num_lus}

    # ── Multi-account distribution ──
    pool = AccountPool()
    authenticated = pool.get_authenticated()

    # Build a list of all available accounts (pool + default)
    import pathlib
    all_accounts = list(authenticated)  # copy pool accounts

    # Always add the default account if it exists
    default_storage = pathlib.Path.home() / ".notebooklm" / "storage_state.json"
    if default_storage.exists():
        _default_email = "default-account"
        try:
            import json
            _ss = json.loads(default_storage.read_text())
            for origin in _ss.get("origins", []):
                for item in origin.get("localStorage", []):
                    val = item.get("value", "")
                    if "@gmail.com" in val or "@google.com" in val:
                        import re as _re_email
                        _m = _re_email.search(r'[\w.+-]+@[\w.-]+', val)
                        if _m:
                            _default_email = _m.group()
                            break
        except Exception:
            pass

        # Check if default account is already in the pool (by path)
        default_already_in_pool = any(
            str(a.storage_state_path) == str(default_storage) for a in all_accounts
        )
        if not default_already_in_pool:
            default_account = AccountInfo(
                email=_default_email,
                password="",
                account_key="default",
                storage_dir=default_storage.parent,
                storage_state_path=default_storage,
                is_authenticated=True,
            )
            # Insert default account at the beginning (primary account)
            all_accounts.insert(0, default_account)

    if not all_accounts:
        return {
            "success": False,
            "message": ("**NotebookLM authentication not found.**\n\n"
                        "Run: `uv run notebooklm login`")
        }

    num_accounts = len(all_accounts)
    if progress_callback:
        acct_names = [a.email.split('@')[0] for a in all_accounts]
        progress_callback(
            f"{len(chunk_meta)} decks to generate, {num_accounts} account(s) available: {', '.join(acct_names)}", 5
        )

    # ── Automatic account switching: run decks sequentially, switch on rate limit ──
    try:
        remaining_chunks = list(chunk_meta)
        lu_results = []
        account_idx = 0

        while remaining_chunks and account_idx < len(all_accounts):
            account = all_accounts[account_idx]
            acct_name = account.email.split('@')[0]
            batch_chunks = list(remaining_chunks)  # Try all remaining on this account
            account.decks_assigned = len(batch_chunks)

            if progress_callback:
                progress_callback(
                    f"[{acct_name}] Starting {len(batch_chunks)} deck(s) on account {account_idx+1}/{num_accounts}...", None
                )

            batch_results = await _run_account_batch(
                account, batch_chunks, course_title, config, progress_callback,
                context=config.get('_context', {})
            )

            if isinstance(batch_results, list) and len(batch_results) > 0:
                # Separate successful and rate-limited results
                successful = []
                rate_limited_chunks = []

                for res in batch_results:
                    status = res.get("generation_status", "")
                    if status == "rate_limited" or "rate" in status.lower():
                        # Find the chunk_meta for this result to retry
                        for cm in batch_chunks:
                            if cm['chunk_idx'] == res.get('chunk_idx'):
                                rate_limited_chunks.append(cm)
                                break
                        if progress_callback:
                            progress_callback(
                                f"[{acct_name}] {res['lu_num']} rate limited — will retry on next account", None
                            )
                    else:
                        successful.append(res)

                lu_results.extend(successful)

                # Also find chunks that weren't even attempted (not in results)
                attempted_indices = {res.get('chunk_idx') for res in batch_results}
                not_attempted = [cm for cm in batch_chunks if cm['chunk_idx'] not in attempted_indices]
                rate_limited_chunks.extend(not_attempted)

                if rate_limited_chunks:
                    remaining_chunks = rate_limited_chunks
                    account_idx += 1
                    if account_idx < len(all_accounts):
                        next_acct = all_accounts[account_idx].email.split('@')[0]
                        if progress_callback:
                            progress_callback(
                                f"[{acct_name}] Rate limited! Auto-switching to {next_acct} for {len(remaining_chunks)} remaining deck(s)...", None
                            )
                        await asyncio.sleep(1)  # Brief pause before switching
                    else:
                        if progress_callback:
                            progress_callback(
                                f"[{acct_name}] Rate limited but no more accounts available. {len(remaining_chunks)} deck(s) not generated.", None
                            )
                        break
                else:
                    remaining_chunks = []  # All done!
            else:
                # Batch returned empty list or non-list — account failed completely
                error_msg = account.error or "Unknown error"
                if progress_callback:
                    progress_callback(
                        f"[{acct_name}] Account failed: {error_msg}. Trying next account...", None
                    )
                account_idx += 1
                if progress_callback:
                    progress_callback(f"[{acct_name}] Batch failed, trying next account...", None)

        # Sort by chunk_idx to maintain order
        lu_results.sort(key=lambda r: r.get("chunk_idx", 0))

        if progress_callback:
            progress_callback("Finalizing PPTX files...", 95)

        # Logo stamping and certificate replacement already done in _generate_chunk_deck_impl
        # Just collect the PPTX paths
        pptx_paths = [r["pptx_path"] for r in lu_results if r.get("pptx_path")]

        # Build account status from all accounts that were used
        account_status = {
            "total": len(all_accounts),
            "authenticated": len(all_accounts),
            "unauthenticated": 0,
            "max_decks_per_account": len(chunk_meta),
            "accounts": [{
                "email": acct.email,
                "authenticated": True,
                "decks_assigned": acct.decks_assigned,
                "decks_completed": acct.decks_completed,
                "decks_failed": acct.decks_failed,
                "error": acct.error,
            } for acct in all_accounts if acct.decks_assigned > 0],
        }

        generated_count = len(lu_results)
        skipped_count = len(skip_lu_indices)
        downloaded_count = len(pptx_paths)
        msg = f"Generated {generated_count}/{len(chunk_meta)} decks across {num_accounts} account(s)."
        if downloaded_count:
            msg += f" Downloaded {downloaded_count} PPTX file(s) with logos."
        elif generated_count > 0:
            msg += " PPTX auto-download pending — check progress messages."
        if skipped_count:
            msg += f" ({skipped_count} decks not selected.)"

        # If 0 generated, surface account errors so user knows WHY
        if generated_count == 0:
            acct_errors = [
                f"{a.email.split('@')[0]}: {a.error}"
                for a in all_accounts if a.error
            ]
            if acct_errors:
                msg += "\n\n**Account errors:**\n" + "\n".join(acct_errors)
            else:
                msg += "\n\nAll accounts returned empty results — check NotebookLM authentication."

        return {
            "success": generated_count > 0,
            "message": msg,
            "lu_results": lu_results,
            "num_lus": num_lus,
            "total_chunks": total_chunks,
            "is_resume": bool(skip_lu_indices),
            "account_status": account_status,
            "pptx_paths": pptx_paths,
        }

    except FileNotFoundError:
        return {
            "success": False,
            "message": ("**NotebookLM authentication not found.**\n\n"
                        "Run: `python -m generate_slides.authenticate_accounts`")
        }
    except Exception as e:
        error_msg = str(e)
        return {
            "success": False,
            "message": f"**Error generating slides:** {error_msg}",
        }


# =============================================================================
# Single-account slide generation (original, kept as fallback)
# =============================================================================

async def _generate_slides_per_lu(context: dict, course_title: str,
                                   config: Dict[str, Any],
                                   progress_callback=None,
                                   skip_lu_indices: set = None) -> Dict[str, Any]:
    """
    Generate slides by splitting LU topics into chunks of ~4-5 topics each.

    Each chunk gets its own NotebookLM deck (~20 slides). Topics are split
    across multiple decks so every topic gets detailed coverage without
    repetition. E.g. LU1 with 13 topics → Deck 1A (T1-T5), 1B (T6-T10), 1C (T11-T13).

    Args:
        context: Full course context dict.
        course_title: Course title string.
        config: Slide generation config.
        progress_callback: Optional callback for UI progress updates.
        skip_lu_indices: Optional set of LU indices (0-based) to skip.
                         Used when resuming after a partial failure.

    Returns:
        Dict with success status and list of per-LU results.
    """
    try:
        from notebooklm import NotebookLMClient
    except ImportError:
        return {
            "success": False,
            "message": ("**notebooklm-py is not installed.**\n\n"
                        "Run: `pip install notebooklm-py[browser]`\n\n"
                        "Then authenticate: `python -m notebooklm login`")
        }

    lus = context.get('Learning_Units', [])
    num_lus = len(lus)
    if num_lus == 0:
        return {"success": False, "message": "No Learning Units found in course info."}

    enable_research = config.get('enable_research', False)
    num_queries = config.get('num_queries', 0)
    slide_style = config.get('slide_style', 'Professional')
    include_notes = config.get('include_notes', False)

    # ── Calculate target decks based on course duration ──
    # 1-day (~8hrs): ~80 slides = 4 decks × 20
    # 2-day (~16hrs): ~140 slides = 7 decks × 20
    # 3-day (~24hrs): ~200 slides = 10 decks × 20
    total_hours = (
        context.get('Total_Course_Duration_Hours', '')
        or context.get('Total_Training_Hours', '')
        or ''
    )
    _hrs_raw = str(total_hours).lower().replace('hours', '').replace('hrs', '').replace('hr', '').replace('h', '').strip()
    _hrs_match = re.search(r'[\d.]+', _hrs_raw)
    try:
        hours = float(_hrs_match.group()) if _hrs_match else 8
    except (ValueError, TypeError):
        hours = 8
    days = max(1, round(hours / 8))
    total_topics = sum(len(lu.get('Topics', [])) for lu in lus)
    if days >= 3:
        target_decks = 11  # ~220 slides
    elif days >= 2:
        target_decks = 7
    else:
        target_decks = 4
    # Cap target to total topics (can't have more decks than topics)
    target_decks = min(target_decks, total_topics)

    # ── Distribute decks proportionally across LUs ──
    topic_counts = [len(lu.get('Topics', [])) for lu in lus if len(lu.get('Topics', [])) > 0]
    lu_indices_with_topics = [i for i in range(num_lus) if len(lus[i].get('Topics', [])) > 0]

    # Proportional allocation: each LU gets decks proportional to its topic count
    raw_alloc = [(tc / total_topics) * target_decks for tc in topic_counts]
    chunks_per_lu = [max(1, int(r)) for r in raw_alloc]
    # Distribute remaining decks to LUs with largest fractional parts
    remaining = target_decks - sum(chunks_per_lu)
    if remaining > 0:
        fracs = [(raw_alloc[i] - chunks_per_lu[i], i) for i in range(len(topic_counts))]
        fracs.sort(reverse=True)
        for _, idx in fracs[:remaining]:
            chunks_per_lu[idx] += 1
    # Ensure no LU has more chunks than topics
    for i, tc in enumerate(topic_counts):
        chunks_per_lu[i] = min(chunks_per_lu[i], tc)

    all_chunks = []
    chunk_idx = 0
    for alloc_idx, lu_idx in enumerate(lu_indices_with_topics):
        lu = lus[lu_idx]
        lu_num = lu.get('LU_Number', f'LU{lu_idx + 1}')
        lu_title = lu.get('LU_Title', 'Learning Unit')
        topics = lu.get('Topics', [])
        num_topics = len(topics)
        num_chunks = chunks_per_lu[alloc_idx]

        # Distribute topics evenly across this LU's chunks
        base = num_topics // num_chunks
        extra = num_topics % num_chunks
        topic_chunks = []
        start = 0
        for ci in range(num_chunks):
            size = base + (1 if ci < extra else 0)
            topic_chunks.append(list(range(start, start + size)))
            start += size

        # Label chunks: if only 1 chunk, just "LU1"; if multiple, "LU1-A", "LU1-B", etc.
        for ci, t_indices in enumerate(topic_chunks):
            if len(topic_chunks) == 1:
                label = lu_num
            else:
                letter = chr(65 + ci)  # A, B, C, ...
                label = f"{lu_num}-{letter}"
            all_chunks.append((chunk_idx, lu_idx, lu_num, lu_title, t_indices, label))
            chunk_idx += 1

    total_chunks = len(all_chunks)
    if total_chunks == 0:
        return {"success": False, "message": "No topics found in any Learning Unit."}

    lu_results = []

    try:
        if progress_callback:
            progress_callback("Connecting to NotebookLM...", 2)

        client = await NotebookLMClient.from_storage()

        if skip_lu_indices is None:
            skip_lu_indices = set()

        async with client:
            # ── Phase 1: Build chunk metadata (skip already completed) ──
            chunk_meta = []
            # is_first/is_last based on FULL course order, not just selected decks
            first_chunk_idx = all_chunks[0][0] if all_chunks else -1
            last_chunk_idx = all_chunks[-1][0] if all_chunks else -1
            for chunk in all_chunks:
                c_idx, lu_idx, lu_num, lu_title, t_indices, label = chunk
                if c_idx in skip_lu_indices:
                    if progress_callback:
                        progress_callback(f"Skipping {label} (already completed)...", 5)
                    continue

                is_first = (c_idx == first_chunk_idx)
                is_last = (c_idx == last_chunk_idx)

                content = _format_chunk_source_text(
                    context, lu_idx, num_lus, t_indices, label,
                    is_first_chunk_of_course=is_first,
                    is_last_chunk_of_course=is_last
                )

                # Topic names for this chunk
                topics = lus[lu_idx].get('Topics', [])
                topic_names = [topics[ti].get('Topic_Title', f'T{ti+1}') for ti in t_indices if ti < len(topics)]
                topic_range = f"T{t_indices[0]+1}-T{t_indices[-1]+1}" if len(t_indices) > 1 else f"T{t_indices[0]+1}"

                chunk_meta.append({
                    'chunk_idx': c_idx,
                    'lu_idx': lu_idx,
                    'lu_num': lu_num,
                    'lu_title': lu_title,
                    'label': label,
                    'topic_indices': t_indices,
                    'topic_range': topic_range,
                    'topic_names': topic_names,
                    'num_topics': len(t_indices),
                    'content': content,
                    'is_first': is_first,
                    'is_last': is_last,
                })

            if not chunk_meta:
                return {"success": True, "message": "All decks already completed.", "lu_results": [], "num_lus": num_lus}

            if progress_callback:
                progress_callback(f"Creating {len(chunk_meta)} notebooks ({total_chunks} total decks)...", 5)

            # ── Create all notebooks concurrently ──
            async def _create_notebook(cm):
                nb_title = f"{course_title} - {cm['label']}: {cm['lu_title']} ({cm['topic_range']})"
                notebook = await client.notebooks.create(nb_title)
                return (cm, notebook.id, nb_title)

            nb_tasks = [_create_notebook(cm) for cm in chunk_meta]
            nb_results = await asyncio.gather(*nb_tasks, return_exceptions=True)

            active_chunks = []
            for res in nb_results:
                if isinstance(res, Exception):
                    logger.warning(f"Notebook creation failed: {res}")
                    continue
                active_chunks.append(res)

            if not active_chunks:
                return {"success": False, "message": "Failed to create any notebooks. Check NotebookLM quota."}

            if progress_callback:
                progress_callback(f"Adding sources to {len(active_chunks)} notebooks (3-5 per deck)...", 15)

            # ── Add multiple sources per notebook ──
            ctx = config.get('_context', {})
            async def _add_sources(cm, notebook_id):
                return await _add_multi_sources(client, notebook_id, cm, course_title, ctx)

            src_tasks = [_add_sources(cm, nb_id) for (cm, nb_id, nb_title) in active_chunks]
            src_results = await asyncio.gather(*src_tasks, return_exceptions=True)

            # Wait for all sources to be ready before generating (prevents failures)
            # URL sources take longer to process — use 60s timeout
            wait_tasks = []
            for i, (cm, nb_id, nb_title) in enumerate(active_chunks):
                src_ids = src_results[i] if not isinstance(src_results[i], Exception) else []
                if src_ids:
                    wait_tasks.append(client.sources.wait_for_sources(nb_id, src_ids, timeout=30.0))
            if wait_tasks:
                await asyncio.gather(*wait_tasks, return_exceptions=True)

            if progress_callback:
                progress_callback(f"Generating {len(active_chunks)} slide decks...", 25)

            # ── Phase 2: Generate slide decks ──
            # Uses the shared _generate_chunk_deck_impl for consistent instructions
            gen_tasks = []
            for i, (cm, nb_id, nb_title) in enumerate(active_chunks):
                src_ids = src_results[i] if not isinstance(src_results[i], Exception) else []
                if not src_ids:
                    logger.warning(f"Skipping {cm['label']} — no sources added")
                    continue
                gen_tasks.append(_generate_chunk_deck_impl(
                    client, cm, nb_id, nb_title, src_ids, course_title, config,
                    progress_callback=progress_callback
                ))

            gen_results = await asyncio.gather(*gen_tasks, return_exceptions=True)

            for i, res in enumerate(gen_results):
                if isinstance(res, Exception):
                    logger.warning(f"Chunk generation failed: {res}")
                    continue
                lu_results.append(res)
                if progress_callback:
                    pct = int(((i + 1) / len(gen_results)) * 70) + 25
                    progress_callback(
                        f"{res['lu_num']} slides generated ({res['generation_status']})",
                        pct
                    )

        if progress_callback:
            progress_callback("All slide decks generated!", 95)

        generated_count = len(lu_results)
        skipped_count = len(skip_lu_indices)
        msg = f"Generated {generated_count} slide deck(s) across {num_lus} LUs."
        if skipped_count:
            msg += f" ({skipped_count} decks not selected.)"

        return {
            "success": True,
            "message": msg,
            "lu_results": lu_results,
            "num_lus": num_lus,
            "total_chunks": total_chunks,
            "is_resume": bool(skip_lu_indices),
        }

    except FileNotFoundError:
        return {
            "success": False,
            "message": ("**NotebookLM authentication not found.**\n\n"
                        "Run in your terminal:\n```\nuv run notebooklm login\n```")
        }
    except Exception as e:
        error_msg = str(e)
        completed = [r["lu_num"] for r in lu_results]
        if "auth" in error_msg.lower() or "login" in error_msg.lower():
            return {
                "success": False,
                "message": (f"**NotebookLM authentication error:** {error_msg}\n\n"
                            "Please re-authenticate:\n```\nuv run notebooklm login\n```"),
                "lu_results": lu_results,
            }
        return {
            "success": False,
            "message": (f"**Error generating slides:** {error_msg}\n\n"
                        f"Completed decks: {', '.join(completed) if completed else 'None'}"),
            "lu_results": lu_results,
        }


# =============================================================================
# NotebookLM logo removal from PPTX
# =============================================================================

def _remove_notebooklm_logo(pptx_bytes: bytes, progress_container=None) -> bytes:
    """
    Remove NotebookLM branding/logo from a PPTX file permanently.

    Simple, reliable approach:
    1. ZIP-level: Edit ALL slide images — paint over bottom-right logo area
    2. XML-level: Remove any shape elements containing NotebookLM text

    Args:
        pptx_bytes: Raw bytes of the uploaded PPTX file.
        progress_container: Optional Streamlit container for debug output.

    Returns:
        Cleaned PPTX file as bytes with NotebookLM branding removed.
    """
    import io
    import zipfile
    from PIL import Image
    from lxml import etree

    removed_count = 0
    details = []

    def _log(msg):
        """Log to both logger and optional Streamlit container."""
        logger.info(msg)
        if progress_container:
            try:
                progress_container.caption(msg)
            except Exception:
                pass

    _log("Starting NotebookLM logo removal...")

    # ═══════════════════════════════════════════════════════════════════
    # SINGLE PASS: Process the entire PPTX as a ZIP archive
    # Edit images (paint over logo) + clean XML (remove text shapes)
    # ═══════════════════════════════════════════════════════════════════
    input_zip = zipfile.ZipFile(io.BytesIO(pptx_bytes), 'r')
    output_buf = io.BytesIO()
    output_zip = zipfile.ZipFile(output_buf, 'w', zipfile.ZIP_DEFLATED)

    images_found = 0
    images_processed = 0
    images_failed = 0

    for item in input_zip.infolist():
        data = input_zip.read(item.filename)
        fname_lower = item.filename.lower()

        # ── Detect images ──
        is_media = 'ppt/media/' in fname_lower
        is_image_ext = fname_lower.endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif', '.webp'))
        is_image_magic = (
            len(data) > 4 and (
                data[:4] == b'\x89PNG' or
                data[:2] == b'\xff\xd8' or
                data[:6] in (b'GIF87a', b'GIF89a') or
                data[:2] == b'BM'
            )
        )

        if is_media or is_image_ext or is_image_magic:
            try:
                img = Image.open(io.BytesIO(data))
                w, h = img.size
                images_found += 1

                # Process slide-sized images (skip tiny icons/thumbnails)
                if w >= 400 and h >= 200:
                    new_data = _erase_logo_region(img)
                    if new_data:
                        # Verify the output is valid
                        verify = Image.open(io.BytesIO(new_data))
                        if verify.size == (w, h):
                            data = new_data
                            images_processed += 1
                            removed_count += 1
                        else:
                            _log(f"  WARNING: {item.filename} — output size mismatch {verify.size} != {(w, h)}")
                            images_failed += 1
                    else:
                        _log(f"  WARNING: {item.filename} — _erase_logo_region returned None")
                        images_failed += 1
            except Exception as e:
                details.append(f"Cannot open {item.filename}: {type(e).__name__}")
                _log(f"  Cannot open {item.filename}: {e}")

        # ── Clean XML: remove shapes with NotebookLM text ──
        elif fname_lower.endswith(('.xml', '.rels')):
            try:
                text_content = data.decode('utf-8', errors='replace')
                if 'notebooklm' in text_content.lower() or 'notebook lm' in text_content.lower():
                    _log(f"  Found NotebookLM text in {item.filename}")

                    # Try XML parsing first
                    try:
                        root = etree.fromstring(data)
                        modified = False

                        for text_elem in root.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}t'):
                            if text_elem.text and 'notebooklm' in text_elem.text.lower():
                                parent = text_elem
                                sp_elem = None
                                for _ in range(20):
                                    parent = parent.getparent()
                                    if parent is None:
                                        break
                                    tag = parent.tag.split('}')[-1] if '}' in parent.tag else parent.tag
                                    if tag in ('sp', 'grpSp', 'pic', 'cxnSp'):
                                        sp_elem = parent
                                        break
                                if sp_elem is not None and sp_elem.getparent() is not None:
                                    sp_elem.getparent().remove(sp_elem)
                                    modified = True
                                    removed_count += 1
                                    _log(f"  Removed XML shape from {item.filename}")
                                    break

                        if modified:
                            data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)
                    except etree.XMLSyntaxError:
                        pass

                    # String replacement fallback
                    text_content = data.decode('utf-8', errors='replace')
                    if 'notebooklm' in text_content.lower():
                        text_content = text_content.replace('NotebookLM', '').replace('notebooklm', '')
                        text_content = text_content.replace('Notebook LM', '').replace('notebook lm', '')
                        data = text_content.encode('utf-8')
                        removed_count += 1
                        _log(f"  String-replaced in {item.filename}")
            except Exception:
                pass

        output_zip.writestr(item, data)

    input_zip.close()
    output_zip.close()
    output_buf.seek(0)

    summary = (
        f"Logo removal complete: {images_found} images found, "
        f"{images_processed} processed, {images_failed} failed"
    )
    _log(summary)
    logger.info(f"Details: {'; '.join(details[:20])}")

    return output_buf.read(), removed_count


def _diagnose_pptx(pptx_bytes: bytes) -> str:
    """Inspect a PPTX file and return diagnostic info about its contents."""
    import io
    import zipfile
    from PIL import Image

    lines = []
    try:
        z = zipfile.ZipFile(io.BytesIO(pptx_bytes), 'r')
        lines.append(f"Total files in PPTX: {len(z.namelist())}")

        images = []
        xmls_with_nlm = []

        for name in z.namelist():
            lower = name.lower()
            # Check images
            if any(lower.endswith(ext) for ext in ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.emf', '.wmf', '.svg')):
                data = z.read(name)
                try:
                    img = Image.open(io.BytesIO(data))
                    images.append(f"  {name}: {img.size[0]}x{img.size[1]} {img.mode} ({img.format})")
                except Exception:
                    images.append(f"  {name}: {len(data)} bytes (cannot open as image)")

            # Check XML for NotebookLM references
            elif lower.endswith('.xml'):
                try:
                    text = z.read(name).decode('utf-8', errors='replace')
                    if 'notebooklm' in text.lower():
                        # Find context around the match
                        idx = text.lower().find('notebooklm')
                        snippet = text[max(0, idx-100):idx+100]
                        xmls_with_nlm.append(f"  {name}: ...{snippet}...")
                except Exception:
                    pass

        lines.append(f"\nImages found: {len(images)}")
        lines.extend(images[:30])

        lines.append(f"\nXML files with 'NotebookLM': {len(xmls_with_nlm)}")
        lines.extend(xmls_with_nlm[:20])

        if not xmls_with_nlm and not images:
            lines.append("\nWARNING: No images and no NotebookLM text found!")

        z.close()
    except Exception as e:
        lines.append(f"Error: {e}")

    return "\n".join(lines)


def _erase_logo_region(img):
    """
    Erase the NotebookLM logo from a slide image.

    Simple left-clone approach with a VERY TIGHT region:
    - Region: 8% width × 3% height (just the logo, nothing else)
    - For each row, clone the pixel from just left of the region
    - No above-reference (avoids text streak artifacts)

    Returns modified image as bytes, or None on failure.
    """
    import io

    try:
        w, h = img.size
        original_format = img.format or 'PNG'

        if img.mode not in ('RGB', 'RGBA'):
            img = img.convert('RGBA')

        # Logo region: 8% width × 3% height — covers NotebookLM branding only
        logo_region_w = max(int(w * 8 / 100), 70)
        logo_region_h = max(int(h * 0.03), 23)
        logo_left = w - logo_region_w
        logo_top = h - logo_region_h

        pixels = img.load()

        # Get the background color at the right edge (above the logo)
        # This tells us what the slide background is at the far right
        right_bg = pixels[min(w - 1, w - 3), max(0, logo_top - 3)]

        for y in range(logo_top, h):
            left_color = pixels[max(0, logo_left - 1), y]

            # Check if left color matches the right-edge background
            dist = (abs(left_color[0] - right_bg[0]) +
                    abs(left_color[1] - right_bg[1]) +
                    abs(left_color[2] - right_bg[2]))

            if dist < 50:
                # Uniform background — simple left-clone
                for x in range(logo_left, w):
                    pixels[x, y] = left_color
            else:
                # Bar on the left doesn't extend to right edge
                # Use left-clone but check each X: if the pixel ABOVE
                # matches the bar → continue bar, otherwise → use right_bg
                for x in range(logo_left, w):
                    above_c = pixels[x, max(0, logo_top - 3)]
                    dist_to_bar = (abs(above_c[0] - left_color[0]) +
                                   abs(above_c[1] - left_color[1]) +
                                   abs(above_c[2] - left_color[2]))
                    dist_to_bg = (abs(above_c[0] - right_bg[0]) +
                                  abs(above_c[1] - right_bg[1]) +
                                  abs(above_c[2] - right_bg[2]))

                    if dist_to_bar <= dist_to_bg:
                        pixels[x, y] = left_color  # bar continues here
                    else:
                        pixels[x, y] = right_bg  # bar ended, use background

        # ── Save ──
        out = io.BytesIO()
        if original_format.upper() == 'JPEG':
            save_img = img.convert('RGB') if img.mode == 'RGBA' else img
            save_img.save(out, format='JPEG', quality=95)
        else:
            img.save(out, format='PNG')
        out.seek(0)
        return out.read()

    except Exception as e:
        logger.error(f"_erase_logo_region failed: {type(e).__name__}: {e}")
        import traceback
        traceback.print_exc()
        return None


def _combine_pptx_files(pptx_files: List[bytes], remove_logo: bool = True) -> tuple:
    """
    Combine multiple PPTX files into one and optionally remove NotebookLM logos.

    Uses the first file as the base presentation, then appends slides from
    subsequent files by copying slide XML and relationships.

    Args:
        pptx_files: List of PPTX file bytes to combine.
        remove_logo: Whether to remove NotebookLM branding.

    Returns:
        Tuple of (combined PPTX bytes, total slide count, logos removed count).
    """
    import io
    from pptx import Presentation
    from copy import deepcopy
    from lxml import etree

    if not pptx_files:
        raise ValueError("No PPTX files provided.")

    if len(pptx_files) == 1:
        if remove_logo:
            clean, count = _remove_notebooklm_logo(pptx_files[0])
            prs = Presentation(io.BytesIO(clean))
            return clean, len(prs.slides), count
        else:
            prs = Presentation(io.BytesIO(pptx_files[0]))
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            return output.read(), len(prs.slides), 0

    # Load base presentation (first file)
    base_prs = Presentation(io.BytesIO(pptx_files[0]))

    # Append slides from remaining files
    for file_bytes in pptx_files[1:]:
        src_prs = Presentation(io.BytesIO(file_bytes))
        for slide in src_prs.slides:
            # Copy slide layout from source — use blank or first available in base
            slide_layout = base_prs.slide_layouts[6] if len(base_prs.slide_layouts) > 6 else base_prs.slide_layouts[0]
            new_slide = base_prs.slides.add_slide(slide_layout)

            # Clear default placeholders from new slide
            for ph in list(new_slide.placeholders):
                sp = ph._element
                sp.getparent().remove(sp)

            # Copy all shapes from source slide
            for shape in slide.shapes:
                el = deepcopy(shape._element)
                new_slide.shapes._spTree.append(el)

            # Copy slide background if present
            src_bg = slide._element.find(
                '{http://schemas.openxmlformats.org/presentationml/2006/main}bg'
            )
            if src_bg is not None:
                new_bg = deepcopy(src_bg)
                existing_bg = new_slide._element.find(
                    '{http://schemas.openxmlformats.org/presentationml/2006/main}bg'
                )
                if existing_bg is not None:
                    new_slide._element.replace(existing_bg, new_bg)
                else:
                    new_slide._element.insert(0, new_bg)

    total_slides = len(base_prs.slides)

    # Save combined, then remove logos
    output = io.BytesIO()
    base_prs.save(output)
    output.seek(0)
    combined_bytes = output.read()

    logos_removed = 0
    if remove_logo:
        combined_bytes, logos_removed = _remove_notebooklm_logo(combined_bytes)

    logger.info(f"Combined {len(pptx_files)} PPTX files: {total_slides} slides, {logos_removed} logos removed")
    return combined_bytes, total_slides, logos_removed


def _add_branded_cover_slide(prs, course_title: str):
    """Add a branded cover slide as the FIRST slide of the presentation.

    Matches the supervisor's reference format:
    - WSQ logo (top-left)
    - Tertiary Infotech Academy logo (top-right)
    - Course title (centered, bold, dark navy)
    - 'Trainer Slides' subtitle
    - Copyright footer

    Args:
        prs: python-pptx Presentation object.
        course_title: Course title to display.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pathlib import Path

    logo_dir = Path(__file__).resolve().parent.parent / "assets" / "slide_logos"
    wsq_logo = logo_dir / "wsq_logo.png"
    tia_logo = logo_dir / "tertiary_infotech_logo.png"

    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    cover = prs.slides.add_slide(blank_layout)

    # Clear default placeholders
    for ph in list(cover.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # WSQ logo — top left
    if wsq_logo.exists():
        cover.shapes.add_picture(str(wsq_logo), Inches(0.5), Inches(0.3), Inches(2.5), None)

    # Tertiary Infotech Academy logo — top right
    if tia_logo.exists():
        cover.shapes.add_picture(str(tia_logo), Inches(9.5), Inches(0.2), Inches(3.3), None)

    # Course title — centered
    title_box = cover.shapes.add_textbox(
        Inches(1.5), Inches(2.8), Inches(10.333), Inches(2.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = course_title
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x26, 0x5A)  # Dark navy
    run.font.name = 'Calibri'

    # Subtitle
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)
    run2 = p2.add_run()
    run2.text = 'Trainer Slides'
    run2.font.size = Pt(24)
    run2.font.color.rgb = RGBColor(0x44, 0x72, 0xC4)  # Steel blue
    run2.font.name = 'Calibri'

    # Copyright footer
    footer_box = cover.shapes.add_textbox(
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.5)
    )
    ftf = footer_box.text_frame
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    frun = fp.add_run()
    frun.text = '\u00a9 2025 Tertiary Infotech Pte Ltd. All Rights Reserved.'
    frun.font.size = Pt(10)
    frun.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    frun.font.name = 'Calibri'

    # Move cover slide to position 0 (first slide)
    # python-pptx adds slides at the end; we need to move it to the front
    from lxml import etree
    sldIdLst = prs.slides._sldIdLst
    sldId_elements = list(sldIdLst)
    if len(sldId_elements) > 1:
        # The cover slide is the last element; move it to first position
        cover_sldId = sldId_elements[-1]
        sldIdLst.remove(cover_sldId)
        sldIdLst.insert(0, cover_sldId)

    return cover


def _merge_pptx_to_single(pptx_paths: list, output_path: str) -> tuple:
    """Merge multiple PPTX files into ONE single PPTX, preserving images and text.

    Properly transfers image blobs (not just XML references) so image backgrounds
    and editable text boxes are both preserved in the combined file.

    Args:
        pptx_paths: List of PPTX file paths to merge (in order).
        output_path: Where to save the combined PPTX.

    Returns:
        (output_path, total_slide_count) tuple.
    """
    import shutil
    import tempfile
    from copy import deepcopy
    from pptx import Presentation

    if not pptx_paths:
        return (None, 0)

    # Filter to existing files only
    valid_paths = [p for p in pptx_paths if p and os.path.exists(p)]
    if not valid_paths:
        return (None, 0)

    if len(valid_paths) == 1:
        shutil.copy2(valid_paths[0], output_path)
        prs = Presentation(output_path)
        return (output_path, len(prs.slides))

    # Use the first file as the base presentation
    base_prs = Presentation(valid_paths[0])

    for src_path in valid_paths[1:]:
        src_prs = Presentation(src_path)

        for slide in src_prs.slides:
            # Add a blank slide
            blank_layout = base_prs.slide_layouts[6] if len(base_prs.slide_layouts) > 6 else base_prs.slide_layouts[0]
            new_slide = base_prs.slides.add_slide(blank_layout)

            # Clear default placeholders
            for ph in list(new_slide.placeholders):
                sp = ph._element
                sp.getparent().remove(sp)

            # Copy shapes — handle images specially to preserve blobs
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture shape
                    try:
                        img_blob = shape.image.blob
                        img_ct = shape.image.content_type or 'image/png'
                        ext = img_ct.split('/')[-1].split('+')[0]
                        if ext not in ('png', 'jpeg', 'jpg', 'gif', 'bmp', 'tiff'):
                            ext = 'png'

                        tmp_img = tempfile.mktemp(suffix=f'.{ext}')
                        with open(tmp_img, 'wb') as f:
                            f.write(img_blob)

                        new_slide.shapes.add_picture(
                            tmp_img,
                            shape.left, shape.top,
                            shape.width, shape.height,
                        )

                        try:
                            os.unlink(tmp_img)
                        except OSError:
                            pass
                    except Exception as e:
                        logger.warning(f"Failed to copy image shape: {e}")
                        # Fallback: deepcopy XML (may break image)
                        new_slide.shapes._spTree.append(deepcopy(shape._element))
                else:
                    # Text boxes, groups, etc. — deepcopy works fine
                    new_slide.shapes._spTree.append(deepcopy(shape._element))

    total_slides = len(base_prs.slides)
    base_prs.save(output_path)
    logger.info(f"Merged {len(valid_paths)} PPTX files into single deck: {total_slides} slides")
    return (output_path, total_slides)


# =============================================================================
# Hybrid Pipeline: NotebookLM (images) + Claude Vision + Editable PPTX
# =============================================================================

async def _extract_slide_text_vision(image_path: str) -> dict:
    """Extract text and layout from a slide image using Claude Vision API.

    Args:
        image_path: Path to the slide image (PNG).

    Returns:
        Dict with 'title', 'bullets', 'layout', 'is_section_header'.
    """
    import anthropic
    import base64

    try:
        with open(image_path, 'rb') as f:
            img_data = base64.standard_b64encode(f.read()).decode('utf-8')

        client = anthropic.Anthropic()
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=1024,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": "image/png",
                            "data": img_data,
                        },
                    },
                    {
                        "type": "text",
                        "text": (
                            "Extract all text from this presentation slide. Return ONLY valid JSON:\n"
                            "{\n"
                            '  "title": "The slide title text",\n'
                            '  "bullets": ["bullet point 1", "bullet point 2", ...],\n'
                            '  "layout": "text-left" | "text-right" | "text-full" | "title-only" | "image-full",\n'
                            '  "is_section_header": true/false,\n'
                            '  "has_diagram": true/false\n'
                            "}\n\n"
                            "Layout guide:\n"
                            '- "text-left": Text on left side, image/diagram on right\n'
                            '- "text-right": Text on right side, image/diagram on left\n'
                            '- "text-full": Text spans full width (no significant images)\n'
                            '- "title-only": Just a title or section header with minimal text\n'
                            '- "image-full": Mostly image/diagram with minimal text\n\n'
                            "If there are no bullet points, set bullets to [].\n"
                            "Output ONLY the JSON, nothing else."
                        ),
                    },
                ],
            }],
        )

        text = response.content[0].text.strip()
        # Parse JSON from response
        import json
        # Try direct parse
        try:
            return json.loads(text)
        except json.JSONDecodeError:
            # Try extracting from code block
            json_match = re.search(r'```(?:json)?\s*\n(.*?)\n```', text, re.DOTALL)
            if json_match:
                return json.loads(json_match.group(1).strip())
            # Try finding JSON object
            brace_start = text.find('{')
            if brace_start != -1:
                depth = 0
                for i in range(brace_start, len(text)):
                    if text[i] == '{':
                        depth += 1
                    elif text[i] == '}':
                        depth -= 1
                        if depth == 0:
                            return json.loads(text[brace_start:i + 1])

        return {"title": "", "bullets": [], "layout": "text-full", "is_section_header": False, "has_diagram": False}

    except Exception as e:
        logger.warning(f"Vision extraction failed for {image_path}: {e}")
        return {"title": "", "bullets": [], "layout": "text-full", "is_section_header": False, "has_diagram": False}


async def _extract_slides_text_batch(image_paths: list) -> list:
    """Extract text from MULTIPLE slides in a single Claude Vision API call.

    Sends up to 10 slide images in one request, returns a list of
    extracted data dicts. This is 10x faster than calling per-slide.

    Args:
        image_paths: List of paths to slide images (PNG).

    Returns:
        List of dicts with 'title', 'bullets', 'layout', etc.
    """
    import anthropic
    import base64
    import json

    empty = {"title": "", "bullets": [], "layout": "text-full", "is_section_header": False, "has_diagram": False}

    if not image_paths:
        return []

    try:
        content_blocks = []
        for idx, img_path in enumerate(image_paths):
            with open(img_path, 'rb') as f:
                img_data = base64.standard_b64encode(f.read()).decode('utf-8')
            content_blocks.append({
                "type": "text",
                "text": f"--- SLIDE {idx + 1} ---",
            })
            content_blocks.append({
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": "image/png",
                    "data": img_data,
                },
            })

        content_blocks.append({
            "type": "text",
            "text": (
                f"\nExtract text from ALL {len(image_paths)} slides above. "
                "Return ONLY a JSON array with one object per slide:\n"
                "[\n"
                '  {"title": "...", "bullets": ["..."], "layout": "text-full|text-left|text-right|title-only|image-full", '
                '"is_section_header": false, "has_diagram": false},\n'
                "  ...\n"
                "]\n\n"
                f"Array MUST have exactly {len(image_paths)} objects, one per slide in order.\n"
                "Output ONLY the JSON array, nothing else."
            ),
        })

        client = anthropic.Anthropic()
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=4096,
            messages=[{"role": "user", "content": content_blocks}],
        )

        text = response.content[0].text.strip()

        # Parse JSON array
        try:
            results = json.loads(text)
            if isinstance(results, list):
                # Pad or trim to match input count
                while len(results) < len(image_paths):
                    results.append(dict(empty))
                return results[:len(image_paths)]
        except json.JSONDecodeError:
            pass

        # Try extracting from code block
        json_match = re.search(r'```(?:json)?\s*\n(.*?)\n```', text, re.DOTALL)
        if json_match:
            try:
                results = json.loads(json_match.group(1).strip())
                if isinstance(results, list):
                    while len(results) < len(image_paths):
                        results.append(dict(empty))
                    return results[:len(image_paths)]
            except json.JSONDecodeError:
                pass

        # Try finding JSON array
        bracket_start = text.find('[')
        if bracket_start != -1:
            depth = 0
            for i in range(bracket_start, len(text)):
                if text[i] == '[':
                    depth += 1
                elif text[i] == ']':
                    depth -= 1
                    if depth == 0:
                        try:
                            results = json.loads(text[bracket_start:i + 1])
                            if isinstance(results, list):
                                while len(results) < len(image_paths):
                                    results.append(dict(empty))
                                return results[:len(image_paths)]
                        except json.JSONDecodeError:
                            break

        logger.warning(f"Batch vision: could not parse response for {len(image_paths)} slides")
        return [dict(empty) for _ in image_paths]

    except Exception as e:
        logger.warning(f"Batch vision extraction failed: {e}")
        return [dict(empty) for _ in image_paths]


def _build_editable_slide_from_image(prs, slide_image_path: str, slide_data: dict):
    """Build a single editable PPTX slide with NotebookLM image background + text boxes.

    Places the original NotebookLM slide image as background, then adds
    perfectly aligned editable text boxes on top. All text stays within
    safe margins (0.6" from edges) for a clean, professional look.

    Slide dimensions: 13.333" x 7.5" (widescreen 16:9)
    Safe area: 0.6" margin all sides, 0.5" footer reserve

    Args:
        prs: python-pptx Presentation object.
        slide_image_path: Path to the slide image (PNG).
        slide_data: Dict from _extract_slide_text_vision with title, bullets, layout.
    """
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    SLIDE_W = prs.slide_width   # 13.333"
    SLIDE_H = prs.slide_height  # 7.5"

    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # 1. Full-page background image — sent to back layer
    pic = slide.shapes.add_picture(
        slide_image_path, Emu(0), Emu(0), SLIDE_W, SLIDE_H
    )
    sp_tree = slide.shapes._spTree
    sp_tree.remove(pic._element)
    sp_tree.insert(2, pic._element)

    title = slide_data.get('title', '')
    bullets = slide_data.get('bullets', [])
    layout = slide_data.get('layout', 'text-full')
    is_section = slide_data.get('is_section_header', False)
    has_diagram = slide_data.get('has_diagram', False)

    # Image-only slides: no text overlay needed
    if layout == 'image-full' and not title:
        return
    # Section headers: centered title only
    if is_section and not bullets:
        if title:
            _add_transparent_title(slide, title, SLIDE_W, SLIDE_H)
        return

    # 2. Safe area layout — all text guaranteed within slide boundaries
    #    Widescreen: 13.333" x 7.5"
    #    Safe margins: 0.6" left/right, 0.5" top, 0.5" footer reserve
    MARGIN_LR = Inches(0.6)     # Left/right margin
    MARGIN_TOP = Inches(0.5)    # Top margin
    FOOTER_H = Inches(0.5)     # Footer reserve at bottom
    SAFE_BOTTOM = SLIDE_H - FOOTER_H  # 7.0" from top
    TITLE_H = Inches(0.85)     # Title box height
    GAP = Inches(0.1)          # Gap between title and body

    # Calculate text area based on layout
    if layout == 'text-left' or (has_diagram and layout != 'text-right'):
        # Left half: 0.6" to 6.4" (5.8" wide)
        text_left = MARGIN_LR
        text_width = Inches(5.8)
    elif layout == 'text-right':
        # Right half: 7.0" to 12.733" (5.733" wide)
        text_left = Inches(7.0)
        text_width = SLIDE_W - text_left - MARGIN_LR
    else:
        # Full width: 0.6" to 12.733" (12.133" wide)
        text_left = MARGIN_LR
        text_width = SLIDE_W - MARGIN_LR * 2

    # --- Title box ---
    if title:
        title_box = slide.shapes.add_textbox(
            text_left, MARGIN_TOP, text_width, TITLE_H
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.auto_size = None  # Fixed size, no auto-grow
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.04)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.name = 'Calibri'
        run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)  # Dark navy

        # Clean white fill — 93% opaque for readability
        _set_shape_fill(title_box, RGBColor(255, 255, 255), alpha=93000)

    # --- Body box (bullets) ---
    if bullets:
        body_top = MARGIN_TOP + TITLE_H + GAP if title else MARGIN_TOP
        body_height = SAFE_BOTTOM - body_top

        # Clamp: ensure body doesn't go below safe area
        if body_height < Inches(0.5):
            body_height = Inches(0.5)

        body_box = slide.shapes.add_textbox(
            text_left, body_top, text_width, body_height
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        tf.auto_size = None  # Fixed size
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.08)

        # Limit bullets to prevent overflow (max ~12 per slide)
        display_bullets = bullets[:12]
        if len(bullets) > 12:
            display_bullets.append(f"... and {len(bullets) - 12} more")

        # Dynamic font size: fewer bullets = larger text, more = smaller
        if len(display_bullets) <= 5:
            bullet_size = Pt(16)
            spacing_after = Pt(8)
        elif len(display_bullets) <= 8:
            bullet_size = Pt(14)
            spacing_after = Pt(6)
        else:
            bullet_size = Pt(12)
            spacing_after = Pt(4)

        for i, bullet in enumerate(display_bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = spacing_after
            p.space_before = Pt(1)
            run = p.add_run()
            run.text = f"\u2022  {bullet}"
            run.font.size = bullet_size
            run.font.name = 'Calibri'
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

        # Clean white fill — 93% opaque
        _set_shape_fill(body_box, RGBColor(255, 255, 255), alpha=93000)


def _set_shape_fill(shape, color: 'RGBColor', alpha: int = 100000):
    """Set shape fill to a solid color with optional transparency.

    Args:
        shape: python-pptx shape object.
        color: RGBColor for fill.
        alpha: Transparency (0 = fully transparent, 100000 = fully opaque).
              90000 = 90% opaque (slightly transparent).
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    spPr = shape._element.spPr
    # Remove any existing fill
    for child in list(spPr):
        if child.tag.endswith('}solidFill') or child.tag.endswith('}noFill'):
            spPr.remove(child)

    solidFill = etree.SubElement(spPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))

    if alpha < 100000:
        alphaElem = etree.SubElement(srgbClr, qn('a:alpha'))
        alphaElem.set('val', str(alpha))


def _add_transparent_title(slide, title: str, slide_w, slide_h):
    """Add a clean centered title overlay for section header slides.

    Centered both horizontally and vertically within the safe area.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    # Centered box: 1.5" margin each side, vertically centered
    box_width = Inches(10.333)
    box_height = Inches(1.8)
    box_left = Inches(1.5)
    box_top = Inches(2.85)  # Visually centered in 7.5" slide

    title_box = slide.shapes.add_textbox(box_left, box_top, box_width, box_height)
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.name = 'Calibri'
    run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    _set_shape_fill(title_box, RGBColor(255, 255, 255), alpha=90000)


async def _convert_notebooklm_to_editable(pdf_path: str, output_path: str,
                                            context: dict = None,
                                            progress_callback=None) -> tuple:
    """Convert a NotebookLM PDF to editable PPTX.

    Each PDF page becomes a full-slide image in the PPTX. This preserves
    all NotebookLM visual elements (images, diagrams, flowcharts, formatting)
    while producing a PPTX that can be opened and edited in PowerPoint.

    Args:
        pdf_path: Path to the NotebookLM PDF.
        output_path: Where to save the editable PPTX.
        context: Optional course context (unused, kept for API compat).
        progress_callback: Optional callback(msg, pct).

    Returns:
        (output_path, slide_count) tuple.
    """
    import tempfile

    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
        import fitz  # pymupdf
    except ImportError as e:
        logger.error(f"Missing dependency for editable conversion: {e}")
        return (pdf_path, 0)

    if progress_callback:
        progress_callback("Converting NotebookLM PDF to editable PPTX...", None)

    # Step 1: Convert PDF pages to high-resolution images
    doc = fitz.open(pdf_path)
    num_pages = len(doc)
    slide_images = []
    temp_dir = tempfile.mkdtemp(prefix="slide_images_")

    for page_num in range(num_pages):
        page = doc[page_num]
        mat = fitz.Matrix(2.0, 2.0)  # 2x resolution for crisp slides
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img_path = os.path.join(temp_dir, f"slide_{page_num:03d}.png")
        pix.save(img_path)
        slide_images.append(img_path)

    doc.close()

    if progress_callback:
        progress_callback(f"Converted {num_pages} pages. Building PPTX...", 50)

    # Step 2: Build PPTX — each page becomes a full-slide image
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for img_path in slide_images:
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        # Remove any default placeholders
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)
        # Add image spanning full slide
        slide.shapes.add_picture(
            img_path, Emu(0), Emu(0), prs.slide_width, prs.slide_height
        )

    prs.save(output_path)
    slide_count = len(prs.slides)

    # Cleanup temp images
    for img_path in slide_images:
        try:
            os.unlink(img_path)
        except OSError:
            pass
    try:
        os.rmdir(temp_dir)
    except OSError:
        pass

    if progress_callback:
        progress_callback(f"Editable PPTX ready: {slide_count} slides with images!", 90)

    logger.info(f"Converted NotebookLM PDF to PPTX: {output_path} ({slide_count} slides)")
    return (output_path, slide_count)


def _extract_cached_nblm_images(progress_callback=None) -> list:
    """Find cached NotebookLM PPTX files and extract images from them.

    Looks in the temp directory for notebooklm_*.pptx files from previous runs.
    Extracts slide images (>50KB) and returns a list of image file paths.
    """
    import glob
    import tempfile

    tmp = tempfile.gettempdir()
    nblm_files = sorted(
        glob.glob(os.path.join(tmp, "notebooklm_*.pptx")),
        key=os.path.getmtime, reverse=True
    )

    if not nblm_files:
        print("[SLIDES] No cached NotebookLM PPTX files found")
        return []

    all_images = []
    _tmp_dir = tempfile.mkdtemp(prefix="nblm_cache_")

    for nblm_path in nblm_files:
        try:
            from pptx import Presentation as _CachePrs
            prs = _CachePrs(nblm_path)
            fname = os.path.basename(nblm_path).replace('.pptx', '')
            extracted = 0
            for si, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        blob = shape.image.blob
                        if len(blob) > 50000:  # >50KB = real slide image
                            ext = shape.image.content_type.split('/')[-1]
                            if ext == 'jpeg':
                                ext = 'jpg'
                            img_path = os.path.join(
                                _tmp_dir, f"{fname}_s{si:03d}.{ext}"
                            )
                            with open(img_path, 'wb') as f:
                                f.write(blob)
                            all_images.append(img_path)
                            extracted += 1
                        break  # one image per slide
            print(f"[SLIDES] Cached {fname}: {extracted} images extracted")
        except Exception as e:
            print(f"[SLIDES] Failed to read cached {nblm_path}: {e}")

    if progress_callback and all_images:
        progress_callback(f"Found {len(all_images)} cached NotebookLM images from previous runs", None)
    print(f"[SLIDES] Total cached images: {len(all_images)}")
    return all_images


async def _generate_hybrid_slides(context: dict, course_title: str,
                                    config: Dict[str, Any],
                                    progress_callback=None,
                                    skip_lu_indices: set = None) -> Dict[str, Any]:
    """Generate slides using NotebookLM + Claude AI → combined editable PPTX.

    Pipeline (uses Claude Code subscription + NotebookLM):
    1. IN PARALLEL:
       a) NotebookLM generates visual slides (images, diagrams, flowcharts)
       b) Claude AI generates knowledge/concept content slides (editable text)
    2. Convert NotebookLM PDF → PPTX image slides (preserves all visuals)
    3. Merge: NotebookLM visual slides + Claude AI knowledge slides → one PPTX
    4. Add branded cover slide with logos

    Result: 100+ slides with BOTH visual elements AND deep knowledge content.
    """
    import tempfile

    if progress_callback:
        progress_callback("Step 1: Generating slides (NotebookLM visuals + Claude AI content) in parallel...", 5)

    # Step 1: Run BOTH in parallel for speed
    # NotebookLM has an 8-minute timeout — if it hangs, we use Claude AI results
    NBLM_TIMEOUT = 480  # 8 minutes max for NotebookLM

    nblm_raw_task = _generate_slides_multi_account(
        context, course_title, config,
        progress_callback=progress_callback,
        skip_lu_indices=skip_lu_indices
    )
    # Wrap NotebookLM with timeout so it can't block the pipeline forever
    async def _nblm_with_timeout():
        try:
            return await asyncio.wait_for(nblm_raw_task, timeout=NBLM_TIMEOUT)
        except asyncio.TimeoutError:
            print(f"[SLIDES] NotebookLM TIMED OUT after {NBLM_TIMEOUT}s")
            if progress_callback:
                progress_callback(f"NotebookLM timed out after {NBLM_TIMEOUT}s — using Claude AI slides", None)
            return {"success": False, "message": f"NotebookLM timed out after {NBLM_TIMEOUT//60} minutes", "lu_results": []}

    claude_task = _generate_editable_pptx(
        context, course_title, config,
        progress_callback=progress_callback,
        skip_lu_indices=skip_lu_indices
    )

    nblm_result, claude_result = await asyncio.gather(
        _nblm_with_timeout(), claude_task, return_exceptions=True
    )

    # Handle exceptions
    if isinstance(nblm_result, Exception):
        logger.warning(f"NotebookLM generation failed: {nblm_result}")
        if progress_callback:
            progress_callback(f"NotebookLM FAILED: {nblm_result}", None)
        print(f"[SLIDES] NotebookLM FAILED: {nblm_result}")
        nblm_result = {"success": False, "lu_results": []}
    if isinstance(claude_result, Exception):
        logger.error(f"Claude AI generation failed: {claude_result}")
        if progress_callback:
            progress_callback(f"Claude AI FAILED: {claude_result}", None)
        print(f"[SLIDES] Claude AI FAILED: {claude_result}")
        claude_result = {"success": False, "lu_results": []}

    nblm_ok = isinstance(nblm_result, dict) and nblm_result.get("success")
    claude_ok = isinstance(claude_result, dict) and claude_result.get("success")

    # Log status for debugging
    nblm_msg = nblm_result.get("message", "") if isinstance(nblm_result, dict) else str(nblm_result)
    claude_msg = claude_result.get("message", "") if isinstance(claude_result, dict) else str(claude_result)
    print(f"[SLIDES] NotebookLM: {'OK' if nblm_ok else 'FAILED'} — {nblm_msg[:200]}")
    print(f"[SLIDES] Claude AI: {'OK' if claude_ok else 'FAILED'} — {claude_msg[:200]}")
    if progress_callback:
        progress_callback(
            f"NotebookLM: {'OK' if nblm_ok else 'FAILED'} | Claude AI: {'OK' if claude_ok else 'FAILED'}", None
        )

    # --- CASE 1: Both succeeded — Rebuild PPTX with NotebookLM images in right column ---
    if claude_ok and nblm_ok:
        if progress_callback:
            progress_callback("Both generators done! Rebuilding slides with NotebookLM images...", 65)

        claude_lu_results = claude_result.get("lu_results", [])
        nblm_lu_results = nblm_result.get("lu_results", [])

        try:
            from generate_slides.build_pptx import build_lu_deck

            # --- STEP 1: Extract ALL NotebookLM images (from PDF or PPTX) ---
            # With per-topic chunks, multiple results share the same base LU.
            # Aggregate images and temp dirs by base LU for the rebuild step.
            all_nblm_images = []
            _tmp_dirs = []
            nblm_img_dirs_per_lu = {}   # base lu_num -> [temp dirs with images]
            nblm_images_per_lu = {}     # base lu_num -> [image_paths]
            nblm_images_per_topic = {}  # topic label e.g. "LU1_T1" -> [image_paths]
            for lr in nblm_lu_results:
                pdf_path = lr.get("original_pdf_path")
                pptx_path = lr.get("pptx_path")
                lu_label = lr.get('lu_num', 'LU')  # e.g., "LU1_T1"
                lu_base = lr.get('lu_num_base', lu_label.split('_T')[0])  # e.g., "LU1"
                _tmp_dir = tempfile.mkdtemp(prefix="nblm_img_")
                _tmp_dirs.append(_tmp_dir)
                extracted = 0

                # Method 1: Extract from PDF (highest quality via fitz)
                if pdf_path and os.path.exists(pdf_path):
                    try:
                        import fitz
                        if progress_callback:
                            progress_callback(f"Extracting {lu_label} illustrations from PDF...", None)
                        doc = fitz.open(pdf_path)
                        for page_num in range(len(doc)):
                            page = doc[page_num]
                            mat = fitz.Matrix(2.0, 2.0)
                            pix = page.get_pixmap(matrix=mat, alpha=False)
                            img_path = os.path.join(_tmp_dir, f"{lu_label}_p{page_num:03d}.png")
                            pix.save(img_path)
                            all_nblm_images.append(img_path)
                            extracted += 1
                        doc.close()
                        print(f"[SLIDES] {lu_label}: Extracted {extracted} images from PDF")
                    except Exception as e:
                        print(f"[SLIDES] {lu_label}: PDF extraction failed: {e}")

                # Method 2: Extract from PPTX (when PDF not available — direct download)
                if extracted == 0 and pptx_path and os.path.exists(pptx_path):
                    try:
                        from pptx import Presentation as _ExPrs
                        if progress_callback:
                            progress_callback(f"Extracting {lu_label} illustrations from PPTX...", None)
                        _eprs = _ExPrs(pptx_path)
                        for si, slide in enumerate(_eprs.slides):
                            for shape in slide.shapes:
                                if shape.shape_type == 13:  # Picture
                                    img_blob = shape.image.blob
                                    img_ext = shape.image.content_type.split('/')[-1]
                                    if img_ext == 'jpeg':
                                        img_ext = 'jpg'
                                    img_path = os.path.join(
                                        _tmp_dir, f"{lu_label}_s{si:03d}.{img_ext}"
                                    )
                                    with open(img_path, 'wb') as f:
                                        f.write(img_blob)
                                    # Only use images that are large enough (skip tiny logos)
                                    if len(img_blob) > 50000:  # >50KB = real slide image
                                        all_nblm_images.append(img_path)
                                        extracted += 1
                                    break  # one image per slide
                        print(f"[SLIDES] {lu_label}: Extracted {extracted} images from PPTX")
                    except Exception as e:
                        print(f"[SLIDES] {lu_label}: PPTX extraction failed: {e}")

                if extracted == 0:
                    print(f"[SLIDES] {lu_label}: WARNING - No images extracted")

                # Track per-topic and per-LU images
                lu_imgs = [p for p in all_nblm_images if _tmp_dir in p.replace("\\", "/") or os.path.dirname(p) == _tmp_dir]
                if not lu_imgs:
                    lu_imgs = all_nblm_images[-extracted:] if extracted > 0 else []
                # Per-topic tracking (e.g. "LU1_T1" -> images from that topic only)
                nblm_images_per_topic[lu_label] = lu_imgs
                # Aggregate: multiple topic chunks feed into the same base LU
                nblm_images_per_lu.setdefault(lu_base, []).extend(lu_imgs)
                nblm_img_dirs_per_lu.setdefault(lu_base, []).append(_tmp_dir)

            if all_nblm_images:
                n_images = len(all_nblm_images)
                if progress_callback:
                    progress_callback(
                        f"Extracted {n_images} illustrations. Now extracting NotebookLM text content...", 68
                    )

                # --- STEP 1.5: Extract TEXT from NotebookLM images (OCR via Claude) ---
                # nblm_img_dirs_per_lu is {base_lu: [dir1, dir2, ...]} (multiple topic dirs per LU)
                from courseware_agents.slides.slides_agent import extract_slides_text
                nblm_text_per_lu = {}
                text_tasks = []
                text_task_labels = []
                for lu_base, img_dirs in nblm_img_dirs_per_lu.items():
                    for img_dir in img_dirs:
                        if os.path.exists(img_dir):
                            png_files = [f for f in os.listdir(img_dir) if f.endswith('.png')]
                            if png_files:
                                text_tasks.append(extract_slides_text(img_dir))
                                text_task_labels.append(lu_base)

                if text_tasks:
                    if progress_callback:
                        progress_callback(
                            f"Reading text from {len(text_tasks)} NotebookLM topic decks...", 70
                        )
                    # Run all text extractions in parallel for speed
                    results = await asyncio.gather(
                        *text_tasks, return_exceptions=True
                    )
                    for label, result in zip(text_task_labels, results):
                        if isinstance(result, Exception):
                            logger.warning(f"NotebookLM text extraction for {label} failed: {result}")
                        else:
                            extracted_text = result or []
                            nblm_text_per_lu.setdefault(label, []).extend(extracted_text)
                            print(f"[SLIDES] {label}: Extracted text from {len(extracted_text)} NotebookLM slides")

                total_nblm_text = sum(len(v) for v in nblm_text_per_lu.values())
                if progress_callback:
                    progress_callback(
                        f"Extracted text from {total_nblm_text} NotebookLM slides + {n_images} images. Building editable PPTX...", 75
                    )

                # --- STEP 2: Rebuild each LU deck using ONLY NotebookLM content ---
                # Supervisor requirement: every slide must be NotebookLM-generated.
                # NotebookLM text (extracted via OCR) + NotebookLM images.
                # Claude AI slides are only used as fallback if NotebookLM text
                # extraction fails for a particular LU.
                rebuilt_paths = []
                total_slides = 0
                for lr in claude_lu_results:
                    claude_slides_data = lr.get("slides_data")
                    lu_idx = lr.get("lu_idx")
                    is_first = lr.get("is_first", False)
                    is_last = lr.get("is_last", False)
                    lu_label = lr.get("lu_num", "LU")

                    if lu_idx is None:
                        if lr.get("pptx_path"):
                            rebuilt_paths.append(lr["pptx_path"])
                            total_slides += lr.get("slide_count", 0)
                        continue

                    # --- Build slides_data from NotebookLM text ONLY ---
                    nblm_text = nblm_text_per_lu.get(lu_label, [])
                    # Use ONLY images from this LU's topics — never mix with other courses
                    lu_images = nblm_images_per_lu.get(lu_label, [])

                    if nblm_text:
                        # Group NotebookLM slides into topics by section headers
                        nblm_topics = []
                        current_topic_title = f"{lu_label} Content"
                        current_slides = []
                        for s in nblm_text:
                            if s.get("is_section_header", False):
                                # Save previous topic if it has slides
                                if current_slides:
                                    nblm_topics.append({
                                        "title": current_topic_title,
                                        "ref": "",
                                        "slides": current_slides,
                                        "activity": [],
                                    })
                                current_topic_title = s.get("title", "Section")
                                current_slides = []
                            elif s.get("bullets"):
                                current_slides.append({
                                    "title": s.get("title", ""),
                                    "bullets": s.get("bullets", []),
                                })
                        # Don't forget last group
                        if current_slides:
                            nblm_topics.append({
                                "title": current_topic_title,
                                "ref": "",
                                "slides": current_slides,
                                "activity": [],
                            })

                        slides_data = {"topics": nblm_topics}
                        print(
                            f"[SLIDES] {lu_label}: Using NotebookLM content ONLY — "
                            f"{len(nblm_topics)} topics, "
                            f"{sum(len(t['slides']) for t in nblm_topics)} slides"
                        )
                    elif claude_slides_data:
                        # Fallback: NotebookLM text extraction failed, use Claude AI
                        slides_data = claude_slides_data
                        print(f"[SLIDES] {lu_label}: NotebookLM text unavailable, falling back to Claude AI")
                    else:
                        # Skip — never use original NotebookLM PPTX (may have wrong course images)
                        print(f"[SLIDES] {lu_label}: No usable content — skipping")
                        continue

                    # Build per-topic image mapping for this LU
                    # nblm_images_per_topic has keys like "LU1_T1", "LU1_T2" etc.
                    _lu_topic_images = {}
                    topics_list = slides_data.get("topics", [])
                    for ti in range(len(topics_list)):
                        topic_label = f"{lu_label}_T{ti + 1}"
                        topic_imgs = nblm_images_per_topic.get(topic_label, [])
                        if topic_imgs:
                            _lu_topic_images[ti] = topic_imgs

                    if progress_callback:
                        progress_callback(f"Building {lu_label} editable PPTX from NotebookLM content...", None)

                    try:
                        pptx_path, slide_count = build_lu_deck(
                            context, lu_idx, slides_data,
                            is_first=is_first, is_last=is_last,
                            image_paths=lu_images if lu_images else None,
                            images_per_topic=_lu_topic_images if _lu_topic_images else None,
                        )
                        rebuilt_paths.append(pptx_path)
                        total_slides += slide_count
                        print(f"[SLIDES] {lu_label}: NotebookLM editable PPTX -> {slide_count} slides")
                    except Exception as e:
                        logger.warning(f"Rebuild {lu_label} failed: {e}")
                        # Never fall back to original NotebookLM PPTX (wrong course images)

                # --- STEP 3: Merge rebuilt decks into one PPTX ---
                if rebuilt_paths:
                    if progress_callback:
                        progress_callback(f"Merging {len(rebuilt_paths)} decks with images...", 85)

                    if len(rebuilt_paths) == 1:
                        merged_path = rebuilt_paths[0]
                    else:
                        merged_path = tempfile.mktemp(suffix="_ALL_LUs_with_images.pptx")
                        merged_path, total_slides = _merge_pptx_to_single(rebuilt_paths, merged_path)

                    claude_result["merged_pptx_path"] = merged_path
                    claude_result["pptx_paths"] = [merged_path] + rebuilt_paths
                    claude_result["message"] = (
                        f"{total_slides} editable slides — ALL content from NotebookLM "
                        f"with {n_images} images!"
                    )
                    if progress_callback:
                        progress_callback(
                            f"Done! {total_slides} NotebookLM slides (editable + images)!", 90
                        )
                    print(f"[SLIDES] CASE 1 complete: {total_slides} NotebookLM-only slides with {n_images} images")
            else:
                # No NotebookLM images available — build text-only slides
                # NEVER use cached images from other courses (causes mismatch)
                print("[SLIDES] No NotebookLM images — building text-only editable slides")
                rebuilt_paths = []
                total_slides = 0
                for lr in claude_lu_results:
                    slides_data = lr.get("slides_data")
                    lu_idx = lr.get("lu_idx")
                    is_first = lr.get("is_first", False)
                    is_last = lr.get("is_last", False)
                    lu_label = lr.get("lu_num", "LU")
                    if slides_data is None or lu_idx is None:
                        if lr.get("pptx_path"):
                            rebuilt_paths.append(lr["pptx_path"])
                            total_slides += lr.get("slide_count", 0)
                        continue
                    try:
                        pptx_path, slide_count = build_lu_deck(
                            context, lu_idx, slides_data,
                            is_first=is_first, is_last=is_last,
                        )
                        rebuilt_paths.append(pptx_path)
                        total_slides += slide_count
                    except Exception as e:
                        if lr.get("pptx_path"):
                            rebuilt_paths.append(lr["pptx_path"])
                            total_slides += lr.get("slide_count", 0)
                if rebuilt_paths:
                    if len(rebuilt_paths) == 1:
                        merged_path = rebuilt_paths[0]
                    else:
                        merged_path = tempfile.mktemp(suffix="_ALL_LUs.pptx")
                        merged_path, total_slides = _merge_pptx_to_single(rebuilt_paths, merged_path)
                    claude_result["merged_pptx_path"] = merged_path
                    claude_result["pptx_paths"] = [merged_path] + rebuilt_paths
                    claude_result["message"] = f"{total_slides} editable text-only slides!"
                    print(f"[SLIDES] Text-only: {total_slides} slides")

            # Note: do NOT cleanup _tmp_dirs yet — images are referenced during merge
            # They will be cleaned up by OS temp cleanup
        except Exception as e:
            logger.warning(f"CASE 1 rebuild failed (non-fatal): {e}")
            import traceback
            traceback.print_exc()

        return claude_result

    # --- CASE 2: Only Claude AI succeeded — use editable slides as-is ---
    # NEVER use cached NotebookLM images from previous runs (causes cross-course contamination)
    if claude_ok:
        nblm_error = nblm_result.get("message", "Unknown error") if isinstance(nblm_result, dict) else str(nblm_result)
        claude_result["nblm_failed"] = True
        claude_result["nblm_error"] = nblm_error
        original_msg = claude_result.get("message", "")
        claude_result["message"] = (
            f"{original_msg}\n\n"
            f"**Note:** NotebookLM images unavailable ({nblm_error}). "
            f"Slides contain editable text + shape diagrams only."
        )
        print(f"[SLIDES] CASE 2: NotebookLM FAILED ({nblm_error}), returning Claude-only editable slides")
        if progress_callback:
            progress_callback(
                f"NotebookLM failed — returning editable text + diagrams (no realistic images)", 95
            )
        return claude_result

    # --- CASE 3: Only NotebookLM succeeded — convert to editable PPTX ---
    if nblm_ok:
        lu_results = nblm_result.get("lu_results", [])
        if not lu_results:
            return nblm_result

        if progress_callback:
            progress_callback(
                f"Step 2: Converting {len(lu_results)} NotebookLM deck(s) to editable PPTX...", 60
            )

        updated_results = []
        for lr_idx, lr in enumerate(lu_results):
            pptx_path = lr.get("pptx_path")
            pdf_path = lr.get("original_pdf_path")
            lu_num = lr.get("lu_num", "LU")

            if not pptx_path or not os.path.exists(pptx_path):
                updated_results.append(lr)
                continue

            source_for_conversion = pdf_path if (pdf_path and os.path.exists(pdf_path)) else None

            if source_for_conversion:
                if progress_callback:
                    progress_callback(f"[{lu_num}] Converting to editable PPTX...", None)

                editable_path = tempfile.mktemp(suffix=f"_{lu_num}_editable.pptx")
                try:
                    editable_path, slide_count = await _convert_notebooklm_to_editable(
                        source_for_conversion, editable_path,
                        context=context,
                        progress_callback=progress_callback
                    )
                    lr["pptx_path"] = editable_path
                    lr["slide_count"] = slide_count
                    lr["generation_status"] = "completed"
                    print(f"[SLIDES] [{lu_num}] Editable PPTX ready: {slide_count} slides")
                    if progress_callback:
                        progress_callback(f"[{lu_num}] {slide_count} editable slides ready!", None)
                except Exception as e:
                    logger.warning(f"[{lu_num}] Editable conversion failed: {e}. Using original.")
                    if not lr.get("slide_count"):
                        try:
                            from pptx import Presentation as _SPrs
                            lr["slide_count"] = len(_SPrs(pptx_path).slides)
                        except Exception:
                            pass
            else:
                if not lr.get("slide_count"):
                    try:
                        from pptx import Presentation as _SPrs
                        lr["slide_count"] = len(_SPrs(pptx_path).slides)
                        lr["generation_status"] = "completed"
                    except Exception:
                        pass

            updated_results.append(lr)

        # Collect individual PPTX paths
        individual_paths = [r["pptx_path"] for r in updated_results if r.get("pptx_path")]
        total_slides = sum(r.get("slide_count", 0) for r in updated_results)

        # Step 3: Merge ALL LU decks into ONE single PPTX
        merged_path = None
        if len(individual_paths) > 1:
            if progress_callback:
                progress_callback(f"Step 3: Merging {len(individual_paths)} decks into one PPTX...", 90)

            merged_path = tempfile.mktemp(suffix="_ALL_LUs.pptx")
            try:
                merged_path, total_slides = _merge_pptx_to_single(individual_paths, merged_path)
                if progress_callback:
                    progress_callback(f"Merged into single PPTX: {total_slides} slides!", 95)
            except Exception as e:
                logger.error(f"Merge failed: {e}. Individual files still available.")
                merged_path = None
        elif len(individual_paths) == 1:
            merged_path = individual_paths[0]

        # Step 4: Add branded cover slide
        if merged_path and os.path.exists(merged_path):
            try:
                if progress_callback:
                    progress_callback("Adding branded cover slide...", 97)
                from pptx import Presentation as _Prs
                _mprs = _Prs(merged_path)
                _add_branded_cover_slide(_mprs, course_title)
                _mprs.save(merged_path)
                total_slides = len(_mprs.slides)
                if progress_callback:
                    progress_callback(f"Cover slide added! Total: {total_slides} slides.", 98)
            except Exception as e:
                logger.warning(f"Cover slide failed (non-fatal): {e}")

        pptx_paths = []
        if merged_path and os.path.exists(merged_path):
            pptx_paths.append(merged_path)
        pptx_paths.extend(individual_paths)

        msg = f"Generated {len(updated_results)} LU deck(s) → {total_slides} total slides!"
        if progress_callback:
            progress_callback(msg, 100)

        return {
            "success": len(individual_paths) > 0,
            "message": msg,
            "lu_results": updated_results,
            "num_lus": nblm_result.get("num_lus", 0),
            "total_chunks": nblm_result.get("total_chunks", 0),
            "is_resume": nblm_result.get("is_resume", False),
            "pptx_paths": pptx_paths,
            "merged_pptx_path": merged_path,
        }

    # --- CASE 4: Both failed ---
    return {
        "success": False,
        "message": "Both NotebookLM and Claude AI generation failed.",
        "lu_results": [],
        "num_lus": 0,
        "total_chunks": 0,
        "is_resume": False,
        "pptx_paths": [],
        "merged_pptx_path": None,
    }


async def _convert_pptx_to_editable(pptx_path: str, output_path: str,
                                      context: dict = None,
                                      progress_callback=None) -> tuple:
    """Convert an image-based PPTX (from NotebookLM) to editable PPTX.

    Extracts each slide as an image, runs Claude Vision, rebuilds with text boxes.

    Args:
        pptx_path: Path to the image-based PPTX.
        output_path: Where to save the editable PPTX.
        context: Optional course context.
        progress_callback: Optional callback.

    Returns:
        (output_path, slide_count) tuple.
    """
    import tempfile
    from pathlib import Path

    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
    except ImportError as e:
        logger.error(f"Missing dependency: {e}")
        return (pptx_path, 0)

    if progress_callback:
        progress_callback("Extracting slide images from PPTX...", None)

    # Extract slide images using the existing slide image data
    source_prs = Presentation(pptx_path)
    temp_dir = tempfile.mkdtemp(prefix="pptx_slides_")
    slide_images = []

    for i, slide in enumerate(source_prs.slides):
        # Find the main image on each slide (NotebookLM puts one big image per slide)
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                image = shape.image
                img_path = os.path.join(temp_dir, f"slide_{i:03d}.png")
                with open(img_path, 'wb') as f:
                    f.write(image.blob)
                slide_images.append(img_path)
                break
        else:
            # No image found — skip this slide
            continue

    if not slide_images:
        logger.warning("No images found in PPTX — returning original")
        return (pptx_path, len(source_prs.slides))

    # Now convert using the same pipeline as PDF conversion
    return await _convert_notebooklm_to_editable_from_images(
        slide_images, output_path, context, progress_callback
    )


async def _convert_notebooklm_to_editable_from_images(
    slide_images: list, output_path: str,
    context: dict = None, progress_callback=None
) -> tuple:
    """Build editable PPTX from a list of slide images.

    Shared implementation used by both PDF and PPTX conversion paths.

    Args:
        slide_images: List of paths to slide images (PNG).
        output_path: Where to save the editable PPTX.
        context: Optional course context.
        progress_callback: Optional callback.

    Returns:
        (output_path, slide_count) tuple.
    """
    from pptx import Presentation
    from pptx.util import Inches
    import asyncio

    num_slides = len(slide_images)

    if progress_callback:
        progress_callback(f"Analyzing {num_slides} slides with Claude AI...", None)

    # Extract text using Claude Agent SDK (Read tool supports images)
    # Uses Claude Code subscription — no separate API key needed
    from courseware_agents.slides.slides_agent import extract_slides_text

    # The images are already saved as files — find their parent directory
    if slide_images:
        _img_dir = os.path.dirname(slide_images[0])
    else:
        _img_dir = ""

    if progress_callback:
        progress_callback(f"Claude AI reading {num_slides} slide images...", 25)

    extracted_data = await extract_slides_text(_img_dir) if _img_dir else []

    # Pad with empty entries if extraction returned fewer results
    empty = {"title": "", "bullets": [], "layout": "image-full",
             "is_section_header": False, "has_diagram": False}
    while len(extracted_data) < num_slides:
        extracted_data.append(dict(empty))

    if progress_callback:
        progress_callback("Building editable PPTX with images + text boxes...", 80)

    # Build editable PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for img_path, slide_data in zip(slide_images, extracted_data):
        _build_editable_slide_from_image(prs, img_path, slide_data)

    prs.save(output_path)
    slide_count = len(prs.slides)

    # Cleanup temp images
    for img_path in slide_images:
        try:
            os.unlink(img_path)
        except OSError:
            pass

    if progress_callback:
        progress_callback(f"Editable PPTX ready: {slide_count} slides!", 95)

    return (output_path, slide_count)


# =============================================================================
# NotebookLM per-topic image generation (used by sequential hybrid pipeline)
# =============================================================================

def _get_nblm_storage_paths() -> list:
    """Get available NotebookLM storage state paths for the active accounts.

    Uses the 3 designated accounts (training11, training12, training13).
    Each account can generate up to 3 decks with 21 slides each.
    Returns list of (label, path) tuples.
    """
    import pathlib as _pl
    paths = []
    accounts_dir = _pl.Path.home() / ".notebooklm" / "accounts"

    # Active accounts only — these are the ones configured for slide generation
    active_accounts = ["training11", "training12", "training13"]

    for name in active_accounts:
        ss = accounts_dir / name / "storage_state.json"
        if ss.exists():
            paths.append((name, str(ss)))

    # Fallback: if no active accounts found, check all account directories
    if not paths and accounts_dir.exists():
        for d in sorted(accounts_dir.iterdir()):
            if d.name == "default":
                continue
            ss = d / "storage_state.json"
            if ss.exists():
                paths.append((d.name, str(ss)))

    return paths


def _get_nblm_pool_images(n_images: int, topic_idx: int = 0) -> list:
    """Get images from the pre-downloaded NotebookLM image pool.

    Falls back to this pool when live NotebookLM generation is rate-limited
    or unavailable. Returns up to n_images paths, cycling through available
    decks based on topic_idx for variety.

    The pool lives at ~/.notebooklm/nblm_imgs/ with subdirectories per deck.
    Each subdirectory contains slide_NNN.png files.
    """
    import pathlib as _pl
    import random

    pool_dir = _pl.Path.home() / ".notebooklm" / "nblm_imgs"
    if not pool_dir.exists():
        return []

    # Collect all images from all decks (skip cover/summary pages)
    all_images = []
    for deck_dir in sorted(pool_dir.iterdir()):
        if not deck_dir.is_dir():
            continue
        imgs = sorted(deck_dir.glob("*.png"))
        if len(imgs) <= 2:
            all_images.extend(imgs)
        else:
            # Skip first (cover) and last (summary) slide
            all_images.extend(imgs[1:-1])

    if not all_images:
        return []

    # Deterministic shuffle based on topic_idx for variety but reproducibility
    rng = random.Random(42 + topic_idx)
    pool = list(all_images)
    rng.shuffle(pool)

    # Return up to n_images
    return [str(p) for p in pool[:n_images]]


def _run_nblm_login_subprocess() -> str:
    """Run NotebookLM login using do_auth.py as a subprocess.

    Opens Chrome for the user to log in manually.
    Returns status message.
    """
    import subprocess
    import pathlib as _pl

    auth_script = _pl.Path(__file__).parent / "do_auth.py"
    if not auth_script.exists():
        return "Auth script not found"

    try:
        proc = subprocess.Popen(
            ["python", str(auth_script)],
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
        )
        return f"Chrome opened for login (PID: {proc.pid}). Sign in and wait for confirmation."
    except Exception as e:
        return f"Failed to start login: {e}"


async def _try_generate_with_account(
    storage_path: str,
    account_label: str,
    topic_text: str,
    topic_title: str,
    topic_key: str,
    n_content_slides: int,
    cache_dir,
    progress_callback=None,
) -> tuple:
    """Try to generate NotebookLM slides with a specific account.

    Returns:
        (images_list, error_str) — images on success, error on failure.
        error_str is "rate_limited" if rate-limited, other string for other errors.
    """
    import time as _time

    try:
        from notebooklm import NotebookLMClient

        client = await NotebookLMClient.from_storage(path=storage_path)
        async with client:
            # 1. Create notebook
            nb = await client.notebooks.create(f"Slides: {topic_title[:60]}")
            nb_id = nb.id
            logger.info(f"[{topic_key}][{account_label}] Notebook created: {nb_id[:16]}...")

            # 2. Add Claude's content as source
            src = await client.sources.add_text(
                nb_id,
                f"Training Content: {topic_title}",
                topic_text[:50000]
            )
            src_id = src.id

            # 3. Wait for source processing
            try:
                await client.sources.wait_for_sources(nb_id, [src_id], timeout=30.0)
            except Exception:
                await asyncio.sleep(5)

            # 4. Generate slides — each slide gets a diagram/flowchart matching its content
            instructions = (
                f"Create EXACTLY {n_content_slides} visual content slides about: {topic_title}. "
                f"EVERY slide MUST include a diagram, flowchart, process flow, comparison chart, "
                f"or infographic that DIRECTLY illustrates that slide's specific content. "
                f"Follow the EXACT slide structure from the source — one visual per slide section. "
                f"Mix diagram types: process flows for procedures, comparison tables for categories, "
                f"cycle diagrams for iterative concepts, infographics for statistics. "
                f"STRICTLY FORBIDDEN — do NOT generate ANY of these: "
                f"TRAQOM slides, TRAQOM surveys, attendance sheets, certificates, "
                f"skill framework slides, cover slides, title slides, introduction slides, "
                f"summary slides, conclusion slides, QR code slides, survey slides, "
                f"feedback forms, evaluation forms, Q&A slides, contact slides. "
                f"ONLY generate EXACTLY {n_content_slides} topic content slides with matching visuals. "
                f"NOTHING ELSE. No extra slides of any kind."
            )
            gen = await client.artifacts.generate_slide_deck(
                nb_id,
                source_ids=[src_id],
                instructions=instructions[:3000],
            )

            # CHECK RESPONSE IMMEDIATELY for rate limit or failure
            if hasattr(gen, 'is_rate_limited') and gen.is_rate_limited:
                err = getattr(gen, 'error', 'Rate limited')
                logger.warning(f"[{topic_key}][{account_label}] RATE LIMITED: {err}")
                try:
                    await client.notebooks.delete(nb_id)
                except Exception:
                    pass
                return ([], "rate_limited")

            if hasattr(gen, 'is_failed') and gen.is_failed:
                err = getattr(gen, 'error', 'Generation failed')
                logger.warning(f"[{topic_key}][{account_label}] FAILED: {err}")
                try:
                    await client.notebooks.delete(nb_id)
                except Exception:
                    pass
                return ([], f"failed: {err}")

            # 5. Poll for completion (max 60 seconds — fallback to pool if slow)
            completed = False
            start = _time.time()
            timeout = 60
            while _time.time() - start < timeout:
                try:
                    decks = await client.artifacts.list_slide_decks(nb_id)
                    for deck in decks:
                        st = getattr(deck, 'status', None)
                        if getattr(deck, 'is_completed', False) or st == 3 or str(st) == '3':
                            completed = True
                            break
                        if getattr(deck, 'is_failed', False) or st == 4 or str(st) == '4':
                            logger.warning(f"[{topic_key}][{account_label}] Generation FAILED during poll")
                            return ([], "failed_during_generation")
                    if completed:
                        break
                except Exception:
                    pass
                await asyncio.sleep(3)
                elapsed = int(_time.time() - start)
                if elapsed % 15 == 0 and elapsed > 0:
                    logger.info(f"[{topic_key}][{account_label}] Generating... ({elapsed}s)")
                    if progress_callback:
                        progress_callback(f"[{topic_key}] NotebookLM generating images... ({elapsed}s)", None)

            if not completed:
                logger.warning(f"[{topic_key}][{account_label}] TIMEOUT after {timeout}s")
                return ([], "timeout")

            # 6. Download as PDF
            pdf_path = str(cache_dir / f'{topic_key}.pdf')
            await client.artifacts.download_slide_deck(nb_id, pdf_path)

            # 7. Extract images, skip cover + summary
            images = _extract_nblm_images(pdf_path, cache_dir, n_content_slides)
            logger.info(f"[{topic_key}][{account_label}] Extracted {len(images)} images!")
            return (images, None)

    except Exception as e:
        logger.warning(f"[{topic_key}][{account_label}] Exception: {e}")
        return ([], str(e))


async def _generate_nblm_images_for_topic(
    topic_text: str,
    topic_title: str,
    topic_key: str,
    n_content_slides: int,
    progress_callback=None,
) -> list:
    """Generate NotebookLM visual slides for one topic, return matched image paths.

    Multi-account rotation: tries each authenticated account in sequence.
    If one is rate-limited, immediately tries the next account.

    Args:
        topic_text: Claude-generated slide content formatted as text source.
        topic_title: Human-readable topic title.
        topic_key: e.g. "LU1_T1" for logging.
        n_content_slides: Number of Claude content slides (for 1:1 pairing limit).
        progress_callback: Optional progress reporter.

    Returns:
        List of image file paths (max n_content_slides). Empty if all accounts failed.
    """
    import pathlib as _pl

    cache_dir = _pl.Path.home() / 'AppData' / 'Local' / 'Temp' / 'nblm_matched_images' / topic_key
    cache_dir.mkdir(parents=True, exist_ok=True)

    # Check cache — if images already generated for this topic, reuse
    cached = sorted(cache_dir.glob('slide_*.png'))
    if cached:
        logger.info(f"[{topic_key}] Using {len(cached)} cached NotebookLM images")
        return [str(p) for p in cached[:n_content_slides]]

    # Get all available accounts
    accounts = _get_nblm_storage_paths()
    if not accounts:
        logger.warning(f"[{topic_key}] No NotebookLM accounts available")
        return []

    # Try each account until one works
    all_rate_limited = True
    for label, storage_path in accounts:
        if progress_callback:
            progress_callback(
                f"[{topic_key}] Trying account '{label}' for {n_content_slides} images...",
                None
            )

        images, error = await _try_generate_with_account(
            storage_path, label,
            topic_text, topic_title, topic_key, n_content_slides,
            cache_dir, progress_callback,
        )

        if images:
            return images  # Success!

        if error == "rate_limited":
            if progress_callback:
                progress_callback(
                    f"[{topic_key}] Account '{label}' is rate-limited, trying next...",
                    None
                )
            continue  # Try next account
        else:
            all_rate_limited = False
            # Non-rate-limit error — still try next account
            if progress_callback:
                progress_callback(
                    f"[{topic_key}] Account '{label}' failed: {error}",
                    None
                )

    # All accounts failed — skip Web UI (don't open browsers automatically)
    if all_rate_limited:
        logger.warning(f"[{topic_key}] ALL accounts rate-limited — skipping NBLM images")
    else:
        logger.warning(f"[{topic_key}] All accounts failed — skipping NBLM images")

    # No pool fallback — only fresh images for the specific course
    if progress_callback:
        progress_callback(
            f"[{topic_key}] Could not generate images — slides will use diagrams only",
            None
        )
        return [str(cache_dir / f'slide_{i:03d}.png') for i in range(len(pool_images))]
    return []


_NBLM_SKIP_KEYWORDS = [
    "traqom", "survey", "attendance", "certificate", "skill framework",
    "skills framework", "qr code", "scan the qr", "scan this qr",
    "feedback form", "evaluation form", "course evaluation",
    "thank you for attending", "end of presentation", "end of course",
    "any questions", "q & a", "q&a session", "contact us",
    "references and resources", "appendix",
]


def _is_nblm_junk_page(page) -> bool:
    """Check if a PDF page is a non-content page that should be skipped.

    Detects TRAQOM surveys, attendance sheets, certificates, QR codes,
    skill frameworks, and other non-topic pages by extracting text and
    checking against a keyword list.
    """
    try:
        text = page.get_text("text").lower()
        # Check for junk keywords
        for kw in _NBLM_SKIP_KEYWORDS:
            if kw in text:
                logger.info(f"  Skipping NBLM page (matched '{kw}')")
                return True
        # Skip pages with very little text (likely a logo-only or blank page)
        stripped = text.strip()
        if len(stripped) < 10:
            logger.info(f"  Skipping NBLM page (too little text: {len(stripped)} chars)")
            return True
    except Exception:
        pass
    return False


def _extract_nblm_images(pdf_path: str, output_dir, max_images: int) -> list:
    """Extract page images from NotebookLM PDF, skip non-content pages.

    Skips:
    - Page 0 = cover (always skip)
    - Last page = summary/outro (always skip)
    - Any page containing TRAQOM, survey, attendance, certificate, QR code,
      skill framework, or other non-content keywords
    Returns up to max_images content page images.
    """
    import fitz
    import pathlib as _pl

    output_dir = _pl.Path(output_dir)
    doc = fitz.open(pdf_path)
    total = len(doc)

    if total == 0:
        return []

    # Determine which pages to consider
    if total <= 2:
        # Too few pages to skip — use all
        start, end = 0, total
    else:
        start = 1             # Skip cover
        end = total - 1       # Skip summary

    images = []
    for idx in range(start, end):
        if len(images) >= max_images:
            break
        page = doc[idx]
        # Filter out non-content pages (TRAQOM, surveys, certificates, etc.)
        if _is_nblm_junk_page(page):
            continue
        mat = fitz.Matrix(2.0, 2.0)  # 2x resolution for quality
        pix = page.get_pixmap(matrix=mat)
        img_path = str(output_dir / f'slide_{idx:03d}.png')
        pix.save(img_path)
        images.append(img_path)

    doc.close()
    logger.info(f"  Extracted {len(images)}/{total} pages (filtered out {total - len(images) - 2} junk pages)")
    return images


# =============================================================================
# Per-Slide NotebookLM Image Generation
# =============================================================================

async def _generate_nblm_image_for_slide(
    slide_text: str,
    slide_title: str,
    slide_key: str,
    storage_path: str,
    account_label: str,
    progress_callback=None,
):
    """Generate a single NotebookLM visual image for one slide.

    Creates a notebook with the slide's content, requests a 1-slide deck,
    extracts the content page image.

    Args:
        slide_text: Formatted source text for this slide.
        slide_title: Human-readable slide title.
        slide_key: Cache key e.g. "LU1_T1_S1".
        storage_path: Path to account's storage_state.json.
        account_label: Account label for logging.
        progress_callback: Optional progress reporter.

    Returns:
        Image file path on success, "RATE_LIMITED" if rate-limited, None on failure.
    """
    import pathlib as _pl
    import time as _time

    cache_dir = _pl.Path.home() / 'AppData' / 'Local' / 'Temp' / 'nblm_slide_images' / slide_key
    cache_dir.mkdir(parents=True, exist_ok=True)

    # Check cache — reuse if already generated
    cached = sorted(cache_dir.glob('slide_*.png'))
    if cached:
        logger.info(f"[{slide_key}] Using cached image")
        return str(cached[0])

    nb_id = None
    try:
        from notebooklm import NotebookLMClient

        client = await NotebookLMClient.from_storage(path=storage_path)
        async with client:
            # 1. Create notebook for this single slide
            nb = await client.notebooks.create(f"Slide: {slide_title[:60]}")
            nb_id = nb.id

            # 2. Add slide content as source
            src = await client.sources.add_text(
                nb_id,
                f"Slide Content: {slide_title}",
                slide_text[:50000],
            )
            src_id = src.id

            # 3. Wait for source processing (shorter for single slide)
            try:
                await client.sources.wait_for_sources(nb_id, [src_id], timeout=20.0)
            except Exception:
                await asyncio.sleep(3)

            # 4. Generate single-slide deck with targeted diagram instructions
            instructions = (
                f"Create a SINGLE visual slide about: {slide_title}. "
                f"Generate EXACTLY 1 content slide with a relevant diagram, "
                f"flowchart, infographic, or illustration that explains the concepts. "
                f"Do NOT include a cover/title slide or summary slide. "
                f"Focus on creating one high-quality visual that complements the text."
            )
            gen = await client.artifacts.generate_slide_deck(
                nb_id,
                source_ids=[src_id],
                instructions=instructions[:3000],
            )

            # Check for rate limit
            if hasattr(gen, 'is_rate_limited') and gen.is_rate_limited:
                logger.warning(f"[{slide_key}][{account_label}] Rate limited")
                try:
                    await client.notebooks.delete(nb_id)
                except Exception:
                    pass
                return "RATE_LIMITED"

            if hasattr(gen, 'is_failed') and gen.is_failed:
                logger.warning(f"[{slide_key}][{account_label}] Generation failed")
                try:
                    await client.notebooks.delete(nb_id)
                except Exception:
                    pass
                return None

            # 5. Poll for completion (max 45s for single slide)
            completed = False
            start = _time.time()
            timeout = 45
            while _time.time() - start < timeout:
                try:
                    decks = await client.artifacts.list_slide_decks(nb_id)
                    for deck in decks:
                        st_val = getattr(deck, 'status', None)
                        if getattr(deck, 'is_completed', False) or st_val == 3 or str(st_val) == '3':
                            completed = True
                            break
                        if getattr(deck, 'is_failed', False) or st_val == 4 or str(st_val) == '4':
                            logger.warning(f"[{slide_key}][{account_label}] Deck generation failed (status=4)")
                            try:
                                await client.notebooks.delete(nb_id)
                            except Exception:
                                pass
                            return None
                    if completed:
                        break
                except Exception:
                    pass
                await asyncio.sleep(2)

            if not completed:
                logger.warning(f"[{slide_key}][{account_label}] Timed out after {timeout}s")
                try:
                    await client.notebooks.delete(nb_id)
                except Exception:
                    pass
                return None

            # 6. Download PDF and extract single image
            pdf_path = str(cache_dir / f'{slide_key}.pdf')
            await client.artifacts.download_slide_deck(nb_id, pdf_path)

            images = _extract_nblm_images(pdf_path, cache_dir, max_images=1)

            # 7. Cleanup notebook
            try:
                await client.notebooks.delete(nb_id)
            except Exception:
                pass

            if images:
                print(f"[NBLM] [{slide_key}][{account_label}] Image generated OK", flush=True)
            return images[0] if images else None

    except Exception as e:
        logger.warning(f"[{slide_key}][{account_label}] Exception: {e}")
        # Try cleanup
        if nb_id:
            try:
                from notebooklm import NotebookLMClient
                client = await NotebookLMClient.from_storage(path=storage_path)
                async with client:
                    await client.notebooks.delete(nb_id)
            except Exception:
                pass
        return None


async def _generate_nblm_images_for_lu_slides(
    slides_data: dict,
    lu_num: str,
    course_title: str,
    lu_title: str,
    progress_callback=None,
    total_timeout: int = 300,
):
    """Generate per-slide NBLM images for an entire LU, parallelized across accounts.

    For each individual slide, creates a dedicated NBLM notebook and generates
    a diagram/flowchart tailored to that slide's specific content.

    Args:
        slides_data: Claude-generated content with 'topics' list.
        lu_num: e.g. "LU1".
        course_title: Course title string.
        lu_title: Learning unit title string.
        progress_callback: Optional progress reporter.
        total_timeout: Max seconds for entire LU image generation (default 5 min).

    Returns:
        Dict mapping topic_idx -> list of image paths (positional, may contain None),
        or None if generation failed completely.
    """
    import time as _time

    accounts = _get_nblm_storage_paths()
    if not accounts:
        logger.warning(f"[{lu_num}] No NotebookLM accounts available for per-slide generation")
        return None

    n_accounts = len(accounts)

    # Build flat list of all slide tasks
    all_tasks = []
    for ti, topic_data in enumerate(slides_data.get('topics', [])):
        topic_title = topic_data.get('title', f'Topic {ti + 1}')
        # Strip topic number prefix (e.g. "T1: ") for cleaner titles
        clean_topic = topic_title
        if ':' in clean_topic:
            clean_topic = clean_topic.split(':', 1)[1].strip()

        for si, slide_data in enumerate(topic_data.get('slides', [])):
            slide_key = f"{lu_num}_T{ti + 1}_S{si + 1}"
            slide_title = slide_data.get('title', f'Slide {si + 1}')
            source_text = _format_slide_as_nblm_source(
                slide_data, si, clean_topic, course_title, lu_title
            )
            all_tasks.append({
                'topic_idx': ti,
                'slide_idx': si,
                'slide_key': slide_key,
                'slide_title': slide_title,
                'source_text': source_text,
            })

    total_slides = len(all_tasks)
    if total_slides == 0:
        return None

    print(
        f"[NBLM] [{lu_num}] Starting per-slide generation: "
        f"{total_slides} slides, {n_accounts} accounts",
        flush=True,
    )
    if progress_callback:
        progress_callback(
            f"[{lu_num}] Generating {total_slides} individual slide images "
            f"using {n_accounts} account(s)...",
            None,
        )

    # Results storage: topic_idx -> [None] * n_slides
    results = {}
    for ti, topic_data in enumerate(slides_data.get('topics', [])):
        n_slides = len(topic_data.get('slides', []))
        results[ti] = [None] * n_slides

    # Semaphore limits concurrency to number of accounts
    account_semaphore = asyncio.Semaphore(n_accounts)
    account_counter = [0]  # mutable counter for round-robin assignment
    completed_count = [0]
    lu_start = _time.time()

    async def _generate_one_slide(task: dict) -> None:
        """Generate image for one slide, acquiring an account from the pool."""
        async with account_semaphore:
            # Round-robin account selection
            idx = account_counter[0] % n_accounts
            account_counter[0] += 1
            label, storage_path = accounts[idx]

            ti = task['topic_idx']
            si = task['slide_idx']

            result = await _generate_nblm_image_for_slide(
                task['source_text'],
                task['slide_title'],
                task['slide_key'],
                storage_path,
                label,
                progress_callback=progress_callback,
            )

            # Handle rate limiting: try remaining accounts
            if result == "RATE_LIMITED":
                for retry_offset in range(1, n_accounts):
                    retry_idx = (idx + retry_offset) % n_accounts
                    retry_label, retry_path = accounts[retry_idx]
                    print(
                        f"[NBLM] [{task['slide_key']}] Retrying with {retry_label} "
                        f"(account {retry_idx + 1}/{n_accounts})",
                        flush=True,
                    )
                    result = await _generate_nblm_image_for_slide(
                        task['source_text'],
                        task['slide_title'],
                        task['slide_key'],
                        retry_path,
                        retry_label,
                        progress_callback=progress_callback,
                    )
                    if result != "RATE_LIMITED":
                        break

            if result and result != "RATE_LIMITED":
                results[ti][si] = result

            completed_count[0] += 1
            elapsed = _time.time() - lu_start
            if completed_count[0] % 3 == 0 or completed_count[0] == total_slides:
                print(
                    f"[NBLM] [{lu_num}] Progress: {completed_count[0]}/{total_slides} "
                    f"slides ({elapsed:.0f}s elapsed)",
                    flush=True,
                )
                if progress_callback:
                    progress_callback(
                        f"[{lu_num}] Generated {completed_count[0]}/{total_slides} "
                        f"slide images ({elapsed:.0f}s)...",
                        None,
                    )

    # Launch all slide tasks with total timeout
    try:
        await asyncio.wait_for(
            asyncio.gather(
                *[_generate_one_slide(task) for task in all_tasks],
                return_exceptions=True,
            ),
            timeout=total_timeout,
        )
    except asyncio.TimeoutError:
        logger.warning(
            f"[{lu_num}] Per-slide generation timeout ({total_timeout}s) — "
            f"completed {completed_count[0]}/{total_slides}"
        )
        print(
            f"[NBLM] [{lu_num}] TIMEOUT after {total_timeout}s — "
            f"{completed_count[0]}/{total_slides} done",
            flush=True,
        )

    # Build images_per_topic dict (same format as before for build_pptx.py)
    images_per_topic = {}
    total_generated = 0
    for ti, img_list in results.items():
        has_any = any(p is not None for p in img_list)
        if has_any:
            images_per_topic[ti] = img_list  # Preserve positional mapping with Nones
            total_generated += sum(1 for p in img_list if p is not None)

    elapsed = _time.time() - lu_start
    print(
        f"[NBLM] [{lu_num}] Per-slide generation complete: "
        f"{total_generated}/{total_slides} images in {elapsed:.0f}s",
        flush=True,
    )
    if progress_callback:
        progress_callback(
            f"[{lu_num}] Got {total_generated}/{total_slides} slide images!",
            None,
        )

    return images_per_topic if images_per_topic else None


def _format_topic_as_nblm_source(topic_data: dict, course_title: str, lu_title: str) -> str:
    """Format Claude-generated topic content as a text source for NotebookLM.

    ONLY includes slide titles and bullet points — NO activities, surveys, or
    assessment references that could cause NBLM to generate unrelated slides.
    Each slide section is clearly labeled so NBLM generates matching visuals.
    """
    title = topic_data.get('title', 'Topic')
    slides = topic_data.get('slides', [])

    lines = [
        f"Training Course: {course_title}",
        f"Learning Unit: {lu_title}",
        f"Topic: {title}",
        "",
        "IMPORTANT INSTRUCTIONS:",
        f"Generate EXACTLY {len(slides)} visual content slides — one per section below.",
        "Each slide MUST contain a diagram, flowchart, or infographic matching its content.",
        "",
        "STRICTLY FORBIDDEN — do NOT generate any of the following:",
        "- TRAQOM slides or TRAQOM survey slides",
        "- Attendance sheets or attendance tracking slides",
        "- Certificate slides or certificate templates",
        "- Skill framework or skills framework slides",
        "- Cover slides, title slides, or introduction slides",
        "- Summary slides, conclusion slides, or end slides",
        "- QR code slides or scan-to-feedback slides",
        "- Survey slides, feedback forms, or evaluation forms",
        "- Q&A slides, 'Any Questions?' slides, or contact slides",
        "",
        "ONLY generate content slides with matching visuals as described below:",
        "",
        "=" * 60,
        "",
    ]
    for i, slide in enumerate(slides):
        s_title = slide.get('title', '')
        s_bullets = slide.get('bullets', [])
        lines.append(f"=== SLIDE {i + 1}: {s_title} ===")
        lines.append(f"Create a visual diagram or infographic for this slide:")
        lines.append("")
        for bullet in s_bullets:
            lines.append(f"  - {bullet}")
        lines.append("")

    # Explicitly tell NBLM to NOT add extra content
    lines.append("=" * 60)
    lines.append("")
    lines.append(f"TOTAL: Generate EXACTLY {len(slides)} content slides with visuals. NOTHING ELSE.")

    return "\n".join(lines)


def _format_slide_as_nblm_source(
    slide_data: dict,
    slide_idx: int,
    topic_title: str,
    course_title: str,
    lu_title: str,
) -> str:
    """Format a SINGLE slide's content as text source for NotebookLM.

    Provides topic/course context so NBLM understands the domain,
    but focuses on one slide's specific bullets for targeted diagram generation.
    """
    s_title = slide_data.get('title', f'Slide {slide_idx + 1}')
    s_bullets = slide_data.get('bullets', [])

    lines = [
        f"Course: {course_title}",
        f"Learning Unit: {lu_title}",
        f"Topic: {topic_title}",
        "",
        "=" * 60,
        "",
        f"--- Slide: {s_title} ---",
        "",
    ]
    for bullet in s_bullets:
        lines.append(f"  - {bullet}")
    lines.append("")
    lines.append(
        "Create a single visual slide with a relevant diagram, flowchart, "
        "or infographic that illustrates the concepts above. "
        "The visual should complement the text content and help learners "
        "understand the key points at a glance."
    )
    return "\n".join(lines)


def _build_fallback_slides_from_cp(lu: dict, slides_per_topic: int = 5) -> dict:
    """Build slide content directly from Course Proposal data when Claude AI times out.

    Creates structured slides from the CP's topic titles and bullet points,
    ensuring we always produce output even if the agent SDK is unavailable.
    """
    topics = lu.get('Topics', [])
    topic_data_list = []
    for i, topic in enumerate(topics):
        t_title = topic.get('Topic_Title', f'Topic {i + 1}')
        bullet_points = topic.get('Bullet_Points', []) or topic.get('Key_Points', [])
        t_desc = topic.get('Topic_Description', '')

        # Build slides from bullet points (chunk into groups of 5)
        fallback_slides = []
        if bullet_points:
            chunk_size = 5
            for j in range(0, len(bullet_points), chunk_size):
                chunk = bullet_points[j:j + chunk_size]
                if chunk:
                    s_title = t_title if j == 0 else f"{t_title} (Cont'd)"
                    fallback_slides.append({"title": s_title, "bullets": chunk})
        if not fallback_slides:
            fallback_slides = [{"title": t_title, "bullets": [t_desc or f"Content for {t_title}"]}]

        # Pad to slides_per_topic if needed
        while len(fallback_slides) < max(2, slides_per_topic // 2):
            fallback_slides.append({
                "title": f"{t_title} — Details",
                "bullets": [f"Detailed content for {t_title} to be added"],
            })

        topic_data_list.append({
            "title": f"T{i + 1}: {t_title}",
            "ref": "",
            "slides": fallback_slides,
            "activity": [
                f"Activity: {t_title} Practice",
                f"Scenario: Apply {t_title} concepts to a workplace situation",
                "Step 1: Review the key concepts covered",
                "Step 2: Apply to a real-world scenario",
                "Step 3: Discuss findings with your group",
                "Expected Output: Brief presentation of findings",
                "Duration: 15-20 minutes",
            ],
        })

    return {"topics": topic_data_list}


# =============================================================================
# Editable PPTX generation (Claude AI + python-pptx, optional NotebookLM images)
# =============================================================================

async def _generate_editable_pptx(context: dict, course_title: str,
                                    config: Dict[str, Any],
                                    progress_callback=None,
                                    skip_lu_indices: set = None,
                                    use_nblm: bool = False) -> Dict[str, Any]:
    """Generate editable PPTX slide decks using Claude AI + python-pptx.

    For each Learning Unit:
    1. Claude AI generates structured slide content (titles, bullets, activities)
    2. (Optional) NotebookLM generates matching images per topic
    3. python-pptx builds real editable PPTX with text + matched images

    When use_nblm=True:
    - After Claude generates content, each topic's content is sent to NotebookLM
    - NotebookLM generates visual slides matching the exact content
    - Cover/summary pages are filtered out
    - Images are paired 1:1 with Claude content slides
    - If NotebookLM fails for a topic, that topic uses shape diagrams only

    Slide count scales with course duration:
    - 1-day course (8h): 60-100 slides
    - 2-day course (16h): 120-160 slides
    """
    from courseware_agents.slides.slides_agent import generate_slide_content
    from generate_slides.build_pptx import (
        build_lu_deck, Presentation, TEMPLATE_PATH, SLIDE_W, SLIDE_H,
        _remove_all_slides, _strip_template_footers,
    )

    lus = context.get('Learning_Units', [])
    num_lus = len(lus)
    if num_lus == 0:
        return {"success": False, "message": "No Learning Units found in course info."}

    lu_indices_with_topics = [i for i in range(num_lus) if len(lus[i].get('Topics', [])) > 0]
    if not lu_indices_with_topics:
        return {"success": False, "message": "No topics found in any Learning Unit."}

    if skip_lu_indices is None:
        skip_lu_indices = set()

    # ── Calculate dynamic slides_per_topic based on course duration ──
    total_hours_raw = (
        context.get('Total_Course_Duration_Hours', '')
        or context.get('Total_Training_Hours', '')
        or ''
    )
    _hrs_str = str(total_hours_raw).lower()
    _hrs_str = _hrs_str.replace('hours', '').replace('hrs', '').replace('hr', '').replace('h', '').strip()
    _hrs_match = re.search(r'[\d.]+', _hrs_str)
    try:
        course_hours = float(_hrs_match.group()) if _hrs_match else 8
    except (ValueError, TypeError):
        course_hours = 8
    course_days = max(1, round(course_hours / 8))

    # Total topics across ALL LUs (for proportional calculation)
    total_topics = sum(len(lus[i].get('Topics', [])) for i in lu_indices_with_topics)
    total_topics = max(1, total_topics)

    # Target: 1-day = 80 slides (midpoint of 60-100), 2-day = 140 total (midpoint of 120-160)
    target_total = 80 if course_days == 1 else 140 if course_days == 2 else course_days * 70
    # Overhead: ~6 fixed slides (cover, objectives, summary, cert, assessment, thank you)
    # + 3 per topic (section header, diagram, activity)
    overhead = 6 + total_topics * 3
    content_needed = target_total - overhead
    slides_per_topic = max(5, min(15, round(content_needed / total_topics)))

    logger.info(
        f"Dynamic slide scaling: {course_hours}h / {course_days} day(s), "
        f"{total_topics} topics, target ~{target_total} slides, "
        f"{slides_per_topic} content slides per topic"
    )
    if progress_callback:
        progress_callback(
            f"{course_days}-day course ({total_topics} topics) → ~{slides_per_topic} slides/topic, "
            f"target {target_total} total slides", 2
        )

    lu_results = []
    total = len(lu_indices_with_topics)
    generated = 0

    # Create ONE Presentation object — all LUs share it (no merge needed)
    if os.path.exists(TEMPLATE_PATH):
        prs = Presentation(TEMPLATE_PATH)
        _remove_all_slides(prs)
        _strip_template_footers(prs)
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

    for order, lu_idx in enumerate(lu_indices_with_topics):
        if order in skip_lu_indices:
            continue

        lu = lus[lu_idx]
        lu_num = lu.get('LU_Number', f'LU{lu_idx + 1}')
        lu_title = lu.get('LU_Title', 'Learning Unit')
        is_first = (order == 0)
        is_last = (order == len(lu_indices_with_topics) - 1)

        num_t = len(lu.get('Topics', []))
        if progress_callback:
            progress_callback(
                f"[{lu_num}] Researching & generating slides for {num_t} topics ({order + 1}/{total})...",
                int((order / total) * 80)
            )

        try:
            # Step 1: Claude AI generates structured content (parallel per topic)
            # Timeout: 5 minutes per LU (agent SDK can hang; fallback uses CP data)
            CONTENT_TIMEOUT = 300
            print(f"[SLIDES] [{lu_num}] Starting generate_slide_content ({num_t} topics, {slides_per_topic} slides/topic, timeout={CONTENT_TIMEOUT}s)...", flush=True)
            try:
                slides_data = await asyncio.wait_for(
                    generate_slide_content(
                        context, lu, lu_idx, num_lus,
                        slides_per_topic=slides_per_topic
                    ),
                    timeout=CONTENT_TIMEOUT,
                )
            except asyncio.TimeoutError:
                print(f"[SLIDES] [{lu_num}] generate_slide_content TIMED OUT after {CONTENT_TIMEOUT}s — using CP fallback", flush=True)
                if progress_callback:
                    progress_callback(f"[{lu_num}] Content generation timed out — using course proposal data", None)
                # Build fallback from CP data
                slides_data = _build_fallback_slides_from_cp(lu, slides_per_topic)
            print(f"[SLIDES] [{lu_num}] generate_slide_content returned: {len(slides_data.get('topics', []))} topics", flush=True)

            topic_count = len(slides_data.get('topics', []))
            slide_count_est = sum(
                len(t.get('slides', [])) + 2  # +2 for section header + activity
                for t in slides_data.get('topics', [])
            )
            if progress_callback:
                progress_callback(
                    f"[{lu_num}] Content ready ({topic_count} topics, ~{slide_count_est} slides).",
                    None
                )

            # Step 2 (optional): NotebookLM images per topic
            # Each topic gets a NBLM deck with images matching its slides (1:1)
            # Wrapped in try/except so NBLM failures never crash the PPTX build
            images_per_topic = None
            if use_nblm:
                try:
                    if progress_callback:
                        progress_callback(
                            f"[{lu_num}] Getting NotebookLM images for {topic_count} topics...",
                            None
                        )
                    images_per_topic = {}
                    nblm_lu_start = time.time()
                    NBLM_LU_TIMEOUT = 180  # 3 minutes per LU

                    for ti, topic_data in enumerate(slides_data.get('topics', [])):
                        if time.time() - nblm_lu_start > NBLM_LU_TIMEOUT:
                            print(f"[SLIDES] [{lu_num}] NBLM timeout ({NBLM_LU_TIMEOUT}s) — building PPTX with available images", flush=True)
                            if progress_callback:
                                progress_callback(f"[{lu_num}] NBLM timeout — building PPTX with available images", None)
                            break

                        topic_key = f'{lu_num}_T{ti+1}'
                        t_title = topic_data.get('title', f'Topic {ti+1}')
                        n_slides = len(topic_data.get('slides', []))

                        source_text = _format_topic_as_nblm_source(
                            topic_data, course_title, lu_title
                        )
                        if progress_callback:
                            progress_callback(
                                f"[{topic_key}] Getting {n_slides} images for '{t_title[:40]}'...",
                                None
                            )

                        try:
                            topic_images = await asyncio.wait_for(
                                _generate_nblm_images_for_topic(
                                    source_text, t_title, topic_key, n_slides,
                                    progress_callback=progress_callback,
                                ),
                                timeout=90,  # 90s per topic (enough for 1 deck generation)
                            )
                        except asyncio.TimeoutError:
                            print(f"[SLIDES] [{topic_key}] NBLM topic timeout (90s)", flush=True)
                            topic_images = []
                        except Exception as te:
                            print(f"[SLIDES] [{topic_key}] NBLM topic error: {te}", flush=True)
                            topic_images = []

                        if topic_images:
                            images_per_topic[ti] = topic_images
                            if progress_callback:
                                progress_callback(
                                    f"[{topic_key}] Got {len(topic_images)} images!",
                                    None
                                )

                        # Small delay between topics to avoid rate limits
                        if ti < topic_count - 1:
                            await asyncio.sleep(1)

                    # Clear empty dict
                    if not any(images_per_topic.values()):
                        images_per_topic = None
                except Exception as nblm_err:
                    import traceback as _tb
                    print(f"[SLIDES] [{lu_num}] NBLM image step crashed (non-fatal): {nblm_err}\n{_tb.format_exc()}", flush=True)
                    images_per_topic = None

            # Step 3: Build editable PPTX (with matched images if available)
            n_img = sum(len(v) for v in (images_per_topic or {}).values()) if images_per_topic else 0
            print(f"[SLIDES] [{lu_num}] Building PPTX ({topic_count} topics, {n_img} images)...", flush=True)
            if progress_callback:
                img_msg = f" + {n_img} images" if n_img else ""
                progress_callback(f"[{lu_num}] Building PPTX{img_msg}...", None)

            _, slide_count = build_lu_deck(
                context, lu_idx, slides_data,
                is_first=is_first, is_last=is_last,
                images_per_topic=images_per_topic,
                prs=prs,  # Shared Presentation — no merge needed
            )
            print(f"[SLIDES] [{lu_num}] Added {slide_count} slides", flush=True)

            if progress_callback:
                progress_callback(f"[{lu_num}] {slide_count} editable slides added.", None)

            lu_results.append({
                "lu_num": lu_num,
                "lu_title": lu_title,
                "notebook_id": "",
                "notebook_title": f"{course_title} - {lu_num}",
                "task_id": "",
                "generation_status": "completed",
                "research_sources_count": 0,
                "total_sources": 0,
                "chunk_idx": order,
                "slide_count": slide_count,
                "slides_data": slides_data,
                "lu_idx": lu_idx,
                "is_first": is_first,
                "is_last": is_last,
            })
            generated += 1

        except Exception as e:
            import traceback
            tb = traceback.format_exc()
            logger.error(f"[{lu_num}] Editable PPTX generation failed: {e}\n{tb}")
            print(f"[SLIDES] [{lu_num}] FAILED: {e}\n{tb}")
            if progress_callback:
                progress_callback(f"[{lu_num}] Failed: {e}", None)
            lu_results.append({
                "lu_num": lu_num,
                "lu_title": lu_title,
                "notebook_id": "",
                "notebook_title": "",
                "task_id": "",
                "generation_status": f"failed: {e}",
                "chunk_idx": order,
                "pptx_path": None,
            })

    total_slides = len(prs.slides)

    # Save the single PPTX (all LUs already built into one Presentation — no merge needed)
    import tempfile
    safe_title = course_title.replace(":", "").replace("/", "-").replace(" ", "_")[:40]
    merged_path = tempfile.mktemp(suffix=f"_{safe_title}_ALL_LUs.pptx")
    prs.save(merged_path)
    logger.info(f"Built single PPTX: {merged_path} ({total_slides} slides)")

    if progress_callback:
        progress_callback(f"PPTX ready: {total_slides} slides!", 100)

    msg = f"Generated {generated}/{total} LUs in 1 file. {total_slides} total slides."

    return {
        "success": generated > 0,
        "message": msg,
        "lu_results": lu_results,
        "num_lus": num_lus,
        "total_chunks": total,
        "is_resume": bool(skip_lu_indices),
        "pptx_paths": [merged_path],
        "merged_pptx_path": merged_path,
    }


# =============================================================================

def app():
    """Streamlit page for Slides Generation."""
    st.title("Generate Slides")
    st.write("NotebookLM Version 2 Generate Slides")

    # Prompt Templates (editable, collapsed)
    from utils.prompt_template_editor import render_prompt_templates
    render_prompt_templates("slides", "Prompt Templates (Slides)")

    # =================================================================
    # Section 1: Generate Slides per LU (requires extracted course info)
    # =================================================================
    extracted_info = st.session_state.get('extracted_course_info')
    if not extracted_info:
        st.warning("Please extract course info first on the **Extract Course Info** page to generate slides.")
    else:
        course_title = extracted_info.get('Course_Title', 'Course')
        lus = extracted_info.get('Learning_Units', [])
        num_lus = len(lus)
        num_topics = sum(len(lu.get('Topics', [])) for lu in lus)
        # Dynamic target decks based on course duration from CP
        total_hours_raw = (
            extracted_info.get('Total_Course_Duration_Hours', '')
            or extracted_info.get('Total_Training_Hours', '')
            or ''
        )
        # Parse hours robustly — handle "18 hrs", "16", "8 hours", "24h", etc.
        import re as _re
        _hours_str = str(total_hours_raw).lower()
        _hours_str = _hours_str.replace('hours', '').replace('hrs', '').replace('hr', '').replace('h', '').strip()
        if _hours_str in ('n/a', 'na', 'nil', 'none', '-', ''):
            _hours_str = ''
        _hours_match = _re.search(r'[\d.]+', _hours_str)
        try:
            _hours = float(_hours_match.group()) if _hours_match else 0
        except (ValueError, TypeError):
            _hours = 0
        # Fallback: estimate from topic count if hours not in CP
        if _hours < 1.0 and num_topics > 0:
            _hours = max(8.0, num_topics * 2.0)
        if _hours < 8.0:
            _hours = 8.0
        _days = max(1, round(_hours / 8))
        # 1 deck per LU — single large deck with ALL topics per Learning Unit
        _lu_with_topics = [i for i in range(num_lus) if len(lus[i].get('Topics', [])) > 0]
        total_decks = len(_lu_with_topics)
        from generate_slides.multi_agent_config import SLIDE_TARGETS, SLIDES_PER_DAY_DEFAULT
        if _days in SLIDE_TARGETS:
            _min_slides, _max_slides = SLIDE_TARGETS[_days]
        else:
            _min_slides = 140 + (_days - 2) * SLIDES_PER_DAY_DEFAULT
            _max_slides = _min_slides + 20
        st.caption(
            f"**{course_title}** | {_hours:.0f} hours | **{_days} day{'s' if _days > 1 else ''}** | "
            f"{num_lus} LUs | {num_topics} topics | **{total_decks} deck(s)** (1 per LU) | "
            f"**{_min_slides}-{_max_slides} slides** target"
        )

        from utils.agent_runner import submit_agent_job
        from utils.agent_status import render_page_job_status

        # ── Build deck list for selection (1 deck per LU) ──
        deck_options = []
        for _ci, lu_idx in enumerate(_lu_with_topics):
            lu = lus[lu_idx]
            lu_num = lu.get('LU_Number', f'LU{lu_idx + 1}')
            lu_title = lu.get('LU_Title', 'Learning Unit')
            num_t = len(lu.get('Topics', []))
            label = f"{lu_num}: {lu_title} ({num_t} topics)"
            deck_options.append((_ci, label))

        # Deck selector — default all selected
        all_labels = [d[1] for d in deck_options]
        selected_labels = st.multiselect(
            "Select decks to generate",
            options=all_labels,
            default=all_labels,
            help="Uncheck decks you already have. Only selected decks will be generated."
        )
        # Map selected labels back to chunk indices to skip
        selected_indices = {d[0] for d in deck_options if d[1] in selected_labels}
        all_indices = {d[0] for d in deck_options}
        skip_indices = all_indices - selected_indices

        # Generation mode — V2 pipeline only
        use_multi_agent = True
        use_hybrid = False

        # V2 configuration panel
        _multi_agent_config = {}
        if use_multi_agent:
            with st.expander("Slide Generation Configuration", expanded=False):
                _ma_col1, _ma_col2 = st.columns(2)
                with _ma_col1:
                    _research_depth = st.select_slider(
                        "Research depth (sources per topic)",
                        options=[10, 15, 20, 25, 30],
                        value=20,
                        help="Light=10, Normal=20, Deep=30. More sources = richer content but slower.",
                    )
                    _content_model = st.selectbox(
                        "Content model",
                        ["claude-sonnet-4-20250514", "claude-opus-4-5-20251101"],
                        index=0,
                        help="Sonnet 4 is fast and balanced. Opus 4.5 produces richer content but costs more.",
                    )
                with _ma_col2:
                    _blocks_options = ["Auto"] + list(range(4, 15))
                    _blocks_per_topic_raw = st.select_slider(
                        "Infographic slides per topic",
                        options=_blocks_options,
                        value="Auto",
                        help="Auto = computed from course duration (recommended). "
                             "1-day → ~6/topic, 2-day → ~12/topic. Override manually if needed.",
                    )
                    _blocks_per_topic = None if _blocks_per_topic_raw == "Auto" else int(_blocks_per_topic_raw)
                    _skip_infographics = st.checkbox(
                        "Skip infographic generation",
                        value=False,
                        help="If checked, uses text fallback slides instead of infographic images.",
                    )
                _multi_agent_config = {
                    "research_depth": _research_depth,
                    "model": _content_model,
                    "skip_infographics": _skip_infographics,
                }
                if _blocks_per_topic is not None:
                    _multi_agent_config["num_blocks_per_topic"] = _blocks_per_topic

            # Status check
            if _multi_agent_config.get("skip_infographics"):
                st.success("**Infographics**: Skipped — using text fallback slides")

        # Generate button
        button_label = "Generate Slides"
        if st.button(button_label, type="primary"):
            _info = extracted_info
            _title = course_title
            _config = {'_context': extracted_info}
            _skip = skip_indices.copy() if skip_indices else None

            _progress_list = []

            def _progress(msg, pct):
                _progress_list.append((msg, pct))
                logger.info(f"[PROGRESS] {msg}")

            if use_multi_agent:
                from generate_slides.multi_agent_orchestrator import orchestrate_multi_agent_slides
                _ma_cfg = _multi_agent_config.copy()
                # Pass selected company info for branding
                _selected_co = st.session_state.get('selected_company', {})
                # Fallback: auto-match from CP's Registered Training Provider
                if not _selected_co and _info.get('Name_of_Organisation'):
                    _selected_co = _auto_match_company(_info['Name_of_Organisation'])
                if _selected_co:
                    _ma_cfg["company"] = _selected_co

                async def _generate_slides():
                    return await orchestrate_multi_agent_slides(
                        context=_info,
                        config=_ma_cfg,
                        progress_callback=_progress,
                    )
            else:
                async def _generate_slides():
                    return await _generate_editable_pptx(
                        _info, _title, _config,
                        progress_callback=_progress,
                        skip_lu_indices=_skip,
                        use_nblm=False,
                    )

            # Clear previous results for fresh generation
            st.session_state.pop('slides_result', None)
            job = submit_agent_job(
                key="generate_slides",
                label="Generate Slides",
                async_fn=_generate_slides,
            )
            if job:
                # Attach the shared progress list to the job dict
                job["progress_messages"] = _progress_list

            if job is None:
                st.warning("Slide generation is already running.")
            else:
                st.rerun()

        # Agent Status
        def _on_slides_complete(job):
            result = job.get("result")
            if result:
                st.session_state['slides_result'] = result

        job_status = render_page_job_status(
            "generate_slides",
            on_complete=_on_slides_complete,
            running_message=(
                f"Generating {len(selected_indices)} slide deck(s)... "
                + ("Research → Content → Editor → Infographic → PPTX"
                   if use_multi_agent else "Claude AI is building editable PPTX")
            ),
        )

        if job_status == "running":
            # Show live progress messages from background thread
            from utils.agent_runner import get_job
            _running_job = get_job("generate_slides")
            if _running_job:
                msgs = _running_job.get("progress_messages", [])
                if msgs:
                    with st.expander("Live Progress", expanded=True):
                        # Show last 10 messages to keep output concise
                        for msg, _pct in msgs[-10:]:
                            st.text(msg)
            st.stop()

        # Display Results
        result = st.session_state.get('slides_result')
        if result:
            lu_results = result.get('lu_results', [])

            if result.get("success"):
                total_slides = sum(lr.get("slide_count", 0) for lr in lu_results)
                st.success(f"Generated **{total_slides} editable slides** in a single PPTX!")

                # Removed verbose stats and info messages — keep UI clean

                course_name = st.session_state.get('extracted_course_info', {}).get('Course_Title', 'Slides')
                safe_name = re.sub(r'[^\w\s\-]', '', course_name).strip().replace(' ', '_')[:50]

                # Primary download: single merged PPTX (all LUs in one file)
                merged_path = result.get("merged_pptx_path")
                if merged_path and os.path.exists(merged_path):
                    try:
                        with open(merged_path, "rb") as f:
                            merged_bytes = f.read()
                        st.download_button(
                            label=f"Download Complete PPTX ({total_slides} slides)",
                            data=merged_bytes,
                            file_name=f"{safe_name}_Slides.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            type="primary",
                            key="dl_merged_pptx",
                        )
                    except Exception as e:
                        st.warning(f"Could not read merged PPTX: {e}")

                # Also offer individual LU downloads in an expander
                individual_available = [lr for lr in lu_results if lr.get("pptx_path") and os.path.exists(lr.get("pptx_path", ""))]
                if len(individual_available) > 1:
                    with st.expander("Download Individual LU Files"):
                        for lr in individual_available:
                            pptx_path = lr.get("pptx_path")
                            slide_count = lr.get("slide_count", 0)
                            try:
                                with open(pptx_path, "rb") as f:
                                    pptx_bytes = f.read()
                                st.download_button(
                                    label=f"Download {lr.get('lu_number', lr.get('lu_num', 'LU'))} ({slide_count} slides)",
                                    data=pptx_bytes,
                                    file_name=f"{safe_name}_{lr.get('lu_number', lr.get('lu_num', 'LU'))}.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    key=f"dl_pptx_{lr.get('lu_number', lr.get('lu_num', 'LU'))}",
                                )
                            except Exception as e:
                                st.warning(f"Could not read {lr.get('lu_number', 'LU')}: {e}")

                if not merged_path and not individual_available:
                    st.warning("PPTX files not available for download.")

                with st.expander("Generation Details"):
                    for lr in lu_results:
                        slides = lr.get('slide_count', 0)
                        status = lr.get('generation_status', 'N/A')
                        st.markdown(
                            f"- **{lr.get('lu_number', lr.get('lu_num', 'LU'))}: {lr.get('lu_title', '')}** | "
                            f"{slides} slides | Status: {status}"
                        )
            else:
                st.error("Slide generation failed.")
                st.markdown(result.get("message", "Unknown error occurred."))
                if lu_results:
                    st.warning(
                        f"Partially completed: {len(lu_results)} deck(s). "
                        f"Uncheck completed decks and click **Generate Editable Slides** to retry."
                    )


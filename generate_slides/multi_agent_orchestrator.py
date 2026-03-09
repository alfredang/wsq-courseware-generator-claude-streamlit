"""
Multi-Agent Slide Generation Orchestrator (5-Phase Linear Chain)

Coordinates specialized agents to produce infographic-based training slides:
  Phase 1: Research Agent → Web research per topic (parallel)
  Phase 2: Content Generator (1st pass) → Structured content blocks
  Phase 3: Editor Agent → Skeleton with infographic assignments
  Phase 4: Infographic Agent → AntV infographic PNGs (parallel per topic)
  Phase 5: Assembly → Map PNGs to slide positions → Build PPTX

Content slides are ALL infographic images. Standard WSQ slides remain text.
"""

import asyncio
import logging
import os
import tempfile
from typing import Callable, Optional

from courseware_agents.slides.research_agent import research_all_topics
from courseware_agents.slides.content_generator_agent import (
    generate_all_content_blocks,
    assemble_final_slides,
)
from courseware_agents.slides.editor_agent import generate_skeleton
from courseware_agents.slides.infographic_agent import generate_all_infographics
from generate_slides.multi_agent_config import (
    DEFAULT_MODEL,
    DEFAULT_RESEARCH_DEPTH,
    DEFAULT_BLOCKS_PER_TOPIC,
    compute_slides_per_topic,
    compute_per_topic_distribution,
    compute_total_target,
    compute_standard_slide_count,
)

logger = logging.getLogger(__name__)


async def orchestrate_multi_agent_slides(
    context: dict,
    config: dict = None,
    progress_callback: Optional[Callable] = None,
) -> dict:
    """Main orchestrator for multi-agent slide generation (5-phase linear chain).

    Args:
        context: Parsed Course Proposal context dict (from interpret_cp).
        config: Optional config dict with:
            - research_depth: int (10/20/30 sources per topic)
            - model: str (model ID for agents)
            - infographic_model: str (model for infographic agent)
            - skip_infographics: bool (skip image generation)
            - num_blocks_per_topic: int (content blocks per topic, 4-8)
        progress_callback: Optional callback(message, percent) for UI updates.

    Returns:
        Dict with success, message, pptx paths, stats, etc.
    """
    if config is None:
        config = {}

    model = config.get("model", DEFAULT_MODEL)
    research_depth = config.get("research_depth", DEFAULT_RESEARCH_DEPTH)
    skip_infographics = config.get("skip_infographics", False)

    course_title = context.get("Course_Title", "Course")
    lus = context.get("Learning_Units", [])

    # Compute slide count dynamically from CP training hours
    total_hours_raw = (
        context.get("Total_Course_Duration_Hours", "")
        or context.get("Total_Training_Hours", "")
        or ""
    )
    import re as _re
    _hrs_str = str(total_hours_raw).lower().replace("hours", "").replace("hrs", "").replace("hr", "").replace("h", "").strip()
    # Remove 'N/A', 'n/a', 'na', etc.
    if _hrs_str in ("n/a", "na", "nil", "none", "-", ""):
        _hrs_str = ""
    # Extract first number from string
    _hrs_match = _re.search(r'[\d.]+', _hrs_str)
    try:
        total_hours = float(_hrs_match.group()) if _hrs_match else 0.0
    except (ValueError, TypeError):
        total_hours = 0.0

    # Count total topics across all LUs
    total_topics = sum(len(lu.get("Topics", [])) for lu in lus)

    # FALLBACK: If hours not found in CP, estimate from topic count
    # A typical WSQ course has ~4-5 topics per day (8 hrs)
    if total_hours < 1.0 and total_topics > 0:
        estimated_hours = max(8.0, total_topics * 2.0)  # ~2 hrs per topic
        logger.warning(
            f"CP training hours not found (raw='{total_hours_raw}'). "
            f"Estimating {estimated_hours}h from {total_topics} topics."
        )
        total_hours = estimated_hours

    # Safety floor
    if total_hours < 8.0:
        total_hours = 8.0

    course_days = max(1, round(total_hours / 8))
    logger.info(
        f"[HOURS DEBUG] Raw='{total_hours_raw}' | Parsed={total_hours}h | "
        f"Days={course_days} | Topics={total_topics}"
    )

    # Dynamic blocks per topic based on course days and topic count
    total_target = compute_total_target(total_hours)
    standard_slides = compute_standard_slide_count(total_topics)
    num_blocks = config.get("num_blocks_per_topic")
    if num_blocks is None:
        num_blocks = compute_slides_per_topic(total_hours, total_topics)

    # Per-topic distribution to hit exact target (e.g. exactly 120 for 2-day)
    per_topic_blocks = compute_per_topic_distribution(total_hours, total_topics)
    expected_content = sum(per_topic_blocks)
    expected_total = standard_slides + expected_content

    logger.info(
        f"Course: {course_title} | {total_hours}h / {course_days} day(s) | "
        f"{total_topics} topics | {num_blocks} blocks/topic (uniform) | "
        f"Target: {total_target} slides | Standard: {standard_slides} | "
        f"Content: {expected_content} | Expected: {expected_total}"
    )

    def _progress(msg: str, pct: int):
        logger.info(f"[{pct}%] {msg}")
        if progress_callback:
            try:
                progress_callback(msg, pct)
            except Exception:
                pass

    # ===================================================================
    # PHASE 1: Research Agent — Research all topics (parallel)
    # ===================================================================
    _progress(f"Phase 1/5: Researching {total_topics} topics (WebSearch)...", 5)

    # Flatten all topics from context
    all_topics = []
    for lu in lus:
        lo_desc = lu.get("LO", "")
        lu_title = lu.get("LU_Title", "")
        for t in lu.get("Topics", []):
            all_topics.append({
                "topic_title": t.get("Topic_Title", "Topic"),
                "bullet_points": t.get("Bullet_Points", []),
                "lo_description": lo_desc,
                "lu_title": lu_title,
            })

    research_map = {}
    if all_topics:
        try:
            research_map = await research_all_topics(
                topics=all_topics,
                course_title=course_title,
                research_depth=research_depth,
                model=model,
            )
        except Exception as e:
            logger.warning(f"Research phase failed, continuing without research: {e}")

    researched_count = sum(
        1 for v in research_map.values()
        if v.get("sources") and len(v["sources"]) > 0
    )
    total_sources = sum(len(v.get("sources", [])) for v in research_map.values())
    _progress(
        f"Phase 1/5: Research complete — "
        f"{researched_count}/{len(all_topics)} topics, {total_sources} total sources",
        20,
    )

    # ===================================================================
    # PHASE 2: Content Generator (1st pass) — Structured content blocks
    # ===================================================================
    _progress("Phase 2/5: Generating content blocks for infographics...", 25)

    content_map = {}
    if all_topics:
        try:
            content_map = await generate_all_content_blocks(
                topics=all_topics,
                research_map=research_map,
                course_title=course_title,
                num_blocks_per_topic=num_blocks,
                per_topic_blocks=per_topic_blocks,
                model=model,
            )
        except Exception as e:
            logger.warning(f"Content generation failed, continuing with fallback: {e}")

    total_blocks = sum(len(v.get("content_blocks", [])) for v in content_map.values())
    # Log per-topic block counts for debugging
    for i, (t_title, t_data) in enumerate(content_map.items()):
        block_count = len(t_data.get("content_blocks", []))
        target_for_topic = per_topic_blocks[i] if i < len(per_topic_blocks) else num_blocks
        logger.info(f"  Content '{t_title}': {block_count}/{target_for_topic} blocks")
    _progress(
        f"Phase 2/5: Content blocks complete — "
        f"{total_blocks} blocks across {len(content_map)} topics "
        f"(target: {expected_content} content + {standard_slides} standard = {expected_total} total)",
        40,
    )

    # ===================================================================
    # PHASE 3: Editor Agent — Skeleton with infographic assignments
    # ===================================================================
    _progress("Phase 3/5: Creating slide skeleton with infographic assignments...", 45)

    try:
        skeleton = await generate_skeleton(
            context=context,
            content_map=content_map,
            model=model,
        )
    except Exception as e:
        logger.error(f"Skeleton generation failed: {e}")
        return {"success": False, "message": f"Skeleton generation failed: {e}"}

    total_assignments = 0
    for lo in skeleton.get("learning_outcomes", []):
        for lu in lo.get("learning_units", []):
            for topic in lu.get("topics", []):
                total_assignments += len(topic.get("infographic_assignments", []))

    _progress(
        f"Phase 3/5: Skeleton created — "
        f"{total_assignments} infographic assignments",
        50,
    )

    # ===================================================================
    # PHASE 4: Infographic Agent — Generate AntV infographic PNGs
    # ===================================================================
    infographic_map = {}
    infographic_dir = None

    if not skip_infographics:
        _progress("Phase 4/5: Generating infographic images (AntV DSL → HTML → PNG)...", 55)

        infographic_dir = tempfile.mkdtemp(prefix="multi_agent_infographics_")
        try:
            infographic_map = await generate_all_infographics(
                skeleton=skeleton,
                content_map=content_map,
                output_dir=infographic_dir,
                model=config.get("infographic_model", model),
            )
        except Exception as e:
            logger.error(f"Infographic phase failed: {e}", exc_info=True)

        total_infographics = sum(len(v) for v in infographic_map.values())
        generated_count = sum(
            sum(1 for r in v if r.get("generated"))
            for v in infographic_map.values()
        )

        # Log per-topic infographic status
        for topic_title, results in infographic_map.items():
            gen = sum(1 for r in results if r.get("generated"))
            errors = [r.get("error", "") for r in results if not r.get("generated") and r.get("error")]
            logger.info(f"  Infographic '{topic_title}': {gen}/{len(results)} generated")
            for err in errors[:2]:
                logger.warning(f"    Error: {err[:100]}")

        _progress(
            f"Phase 4/5: Infographics — {generated_count}/{total_infographics} generated",
            75,
        )
    else:
        _progress("Phase 4/5: Infographics skipped", 75)

    # ===================================================================
    # PHASE 5: Assembly + PPTX Build
    # ===================================================================
    _progress("Phase 5/5: Assembling slides and building PPTX...", 80)

    # 5a: Assembly — map infographic PNGs to slide positions (pure Python)
    lu_data_map = assemble_final_slides(
        skeleton=skeleton,
        infographic_results=infographic_map,
        content_map=content_map,
    )

    # 5b: Build PPTX files
    try:
        pptx_result = _build_infographic_pptx(
            context=context,
            skeleton=skeleton,
            lu_data_map=lu_data_map,
            content_map=content_map,
            infographic_results=infographic_map,
            progress_callback=_progress,
            company=config.get("company"),
        )
    except Exception as e:
        logger.error(f"PPTX build failed: {e}")
        return {"success": False, "message": f"PPTX build failed: {e}"}

    _progress(f"Phase 5/5: PPTX built! {pptx_result.get('message', '')}", 100)

    # Compute stats
    total_infographics = sum(len(v) for v in infographic_map.values())
    generated_count = sum(
        sum(1 for r in v if r.get("generated"))
        for v in infographic_map.values()
    )

    return {
        "success": True,
        "message": pptx_result.get("message", "Multi-agent slides generated"),
        "merged_pptx_path": pptx_result.get("merged_pptx_path"),
        "pptx_paths": pptx_result.get("pptx_paths", []),
        "lu_results": pptx_result.get("lu_results", []),
        "skeleton": skeleton,
        "research_stats": {
            "topics_researched": researched_count,
            "total_sources": total_sources,
        },
        "content_stats": {
            "total_blocks": total_blocks,
            "topics_with_blocks": len(content_map),
        },
        "infographic_stats": {
            "generated": generated_count,
            "total": total_infographics,
            "output_dir": infographic_dir,
        },
    }


def _recover_infographic_images(infographic_data: dict, infographic_results: dict):
    """Fill in missing image_path on infographic_slides from infographic_results.

    This is the critical safety net: even if assembly or fallback lost image paths,
    this function recovers them by fuzzy-matching topic titles and slide positions
    against the infographic agent's actual output.
    """
    if not infographic_results:
        return

    def _norm(s):
        return s.lower().replace(" ", "").replace("_", "").replace("-", "").strip()

    def _find_results(topic_title):
        """Find infographic results for a topic using fuzzy matching."""
        if topic_title in infographic_results:
            return infographic_results[topic_title]
        norm_key = _norm(topic_title)
        for k, v in infographic_results.items():
            if _norm(k) == norm_key:
                return v
        key_lower = topic_title.lower().strip()
        for k, v in infographic_results.items():
            k_lower = k.lower().strip()
            if key_lower in k_lower or k_lower in key_lower:
                return v
        return []

    for topic in infographic_data.get("topics", []):
        t_title = topic.get("title", "")
        infographic_slides = topic.get("infographic_slides", [])
        missing_count = sum(1 for s in infographic_slides if not s.get("image_path") or not os.path.exists(str(s.get("image_path", ""))))
        if missing_count == 0:
            continue

        # Find matching infographic results
        results = _find_results(t_title)
        if not results:
            logger.warning(f"[IMAGE RECOVERY] No infographic results found for '{t_title}'")
            continue

        recovered = 0
        for slide in infographic_slides:
            img_path = slide.get("image_path")
            if img_path and os.path.exists(img_path):
                continue  # Already has valid image

            pos = slide.get("position", -1)
            slide_title = slide.get("title", "")

            # Try matching by slide_position
            match = None
            for r in results:
                if r.get("generated") and r.get("image_path") and r.get("slide_position") == pos:
                    if os.path.exists(r["image_path"]):
                        match = r
                        break

            # Fallback: match by sub_title
            if not match:
                for r in results:
                    if r.get("generated") and r.get("image_path") and os.path.exists(r["image_path"]):
                        r_sub = _norm(r.get("sub_title", ""))
                        s_title = _norm(slide_title)
                        if r_sub and s_title and (r_sub == s_title or r_sub in s_title or s_title in r_sub):
                            match = r
                            break

            # Fallback: match by index
            if not match:
                idx = infographic_slides.index(slide)
                if idx < len(results) and results[idx].get("generated") and results[idx].get("image_path"):
                    if os.path.exists(results[idx]["image_path"]):
                        match = results[idx]

            if match:
                slide["image_path"] = match["image_path"]
                recovered += 1

        if recovered > 0:
            logger.info(
                f"[IMAGE RECOVERY] '{t_title}': recovered {recovered}/{missing_count} images"
            )
        else:
            logger.warning(
                f"[IMAGE RECOVERY] '{t_title}': could NOT recover any of {missing_count} missing images"
            )


def _build_infographic_pptx(
    context: dict,
    skeleton: dict,
    lu_data_map: dict,
    content_map: dict = None,
    infographic_results: dict = None,
    progress_callback: Optional[Callable] = None,
    company: dict = None,
) -> dict:
    """Build the final PPTX from assembled infographic slide data.

    Builds ALL LUs into a single Presentation object — no merge step needed.
    If lu_data_map has no topics for an LU, falls back to building slides
    directly from content_map (bypasses skeleton/assembly failures).
    """
    from generate_slides.build_pptx import (
        build_lu_deck, Presentation, TEMPLATE_PATH, SLIDE_W, SLIDE_H,
        _remove_all_slides, _strip_template_footers,
    )

    if content_map is None:
        content_map = {}

    lus = context.get("Learning_Units", [])
    num_lus = len(lus)
    lu_results = []
    total_slides = 0

    # Create ONE Presentation object from template — all LUs share it
    if os.path.exists(TEMPLATE_PATH):
        prs = Presentation(TEMPLATE_PATH)
        _remove_all_slides(prs)
        _strip_template_footers(prs)
    else:
        from pptx.util import Emu
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

    # Build a normalized lookup for lu_data_map keys (strip spaces, case-insensitive)
    def _normalize_lu(key):
        return key.replace(" ", "").replace("_", "").upper()

    lu_data_normalized = {_normalize_lu(k): v for k, v in lu_data_map.items()}
    logger.info(f"lu_data_map keys: {list(lu_data_map.keys())}")
    logger.info(f"context LU_Numbers: {[lu.get('LU_Number') for lu in lus]}")
    logger.info(f"content_map keys: {list(content_map.keys())}")

    for lu_idx, lu in enumerate(lus):
        lu_num = lu.get("LU_Number", f"LU{lu_idx + 1}")
        is_first = (lu_idx == 0)
        is_last = (lu_idx == num_lus - 1)

        # ALWAYS build from CONTEXT topics — this guarantees correct topic-to-LU
        # mapping regardless of skeleton quality. The skeleton often misassigns
        # topics to wrong LUs, causing missing or duplicated topics.
        infographic_data = _build_lu_data_from_content_map(
            lu, content_map, lu_idx,
        )
        logger.info(
            f"Built {lu_num} from context: "
            f"{len(infographic_data.get('topics', []))} topics "
            f"({[t.get('title','')[:40] for t in infographic_data.get('topics', [])]})"
        )

        # ── IMAGE RECOVERY ──
        # Fill in image_path from infographic_results for EVERY topic.
        if infographic_results:
            _recover_infographic_images(infographic_data, infographic_results)

        try:
            # NEVER pass is_last=True here — closing slides are added AFTER
            # padding to ensure correct slide ordering.
            _, slides_added = build_lu_deck(
                context=context,
                lu_idx=lu_idx,
                slides_data=infographic_data,
                is_first=is_first,
                is_last=False,
                infographic_mode=True,
                prs=prs,
                company=company,
            )
            total_slides += slides_added
            lu_results.append({
                "lu_number": lu_num,
                "slide_count": slides_added,
                "topic_count": len(infographic_data.get("topics", [])),
            })
            logger.info(f"Built {lu_num}: {slides_added} slides")

        except Exception as e:
            logger.error(f"Failed to build PPTX for {lu_num}: {e}", exc_info=True)
            lu_results.append({
                "lu_number": lu_num,
                "error": str(e),
            })

    # ── TARGET ENFORCEMENT ──
    # RULE: Slide count must ALWAYS hit the target for the course duration.
    # If pipeline produced fewer slides than target, pad with additional content slides.
    total_hours_raw = (
        context.get("Total_Course_Duration_Hours", "")
        or context.get("Total_Training_Hours", "")
        or ""
    )
    import re as _re
    _hrs_str = str(total_hours_raw).lower().replace("hours", "").replace("hrs", "").replace("hr", "").replace("h", "").strip()
    if _hrs_str in ("n/a", "na", "nil", "none", "-", ""):
        _hrs_str = ""
    _hrs_match = _re.search(r'[\d.]+', _hrs_str)
    try:
        _total_hours = float(_hrs_match.group()) if _hrs_match else 8.0
    except (ValueError, TypeError):
        _total_hours = 8.0
    if _total_hours < 8.0:
        _total_hours = 8.0

    slide_target = compute_total_target(_total_hours)
    current_count = len(prs.slides)
    # Account for ~7 closing slides that will be added after padding
    CLOSING_SLIDES_COUNT = 7
    shortfall = slide_target - current_count - CLOSING_SLIDES_COUNT

    if shortfall > 0:
        logger.warning(
            f"SLIDE TARGET ENFORCEMENT: {current_count} slides < {slide_target} target. "
            f"Adding {shortfall} padding slides."
        )
        from generate_slides.build_pptx import add_content_slide, add_infographic_slide

        # Gather all topic titles for padding content
        all_topic_titles = []
        for lu in lus:
            for t in lu.get("Topics", []):
                all_topic_titles.append(t.get("Topic_Title", "Topic"))

        # Build a flat list of ALL available infographic images (with generated=True)
        # for use in padding slides
        _all_infographic_images = []
        if infographic_results:
            for _t_title, _results in infographic_results.items():
                for r in _results:
                    if r.get("generated") and r.get("image_path") and os.path.exists(r["image_path"]):
                        _all_infographic_images.append(r)

        # Add padding slides distributed across topics — USE IMAGES when available
        for pad_i in range(shortfall):
            topic_title = all_topic_titles[pad_i % len(all_topic_titles)] if all_topic_titles else "Course Content"
            # Get content blocks for this topic if available
            topic_content = content_map.get(topic_title, {}) if content_map else {}
            blocks = topic_content.get("content_blocks", [])

            if blocks:
                block = blocks[pad_i % len(blocks)]
                items = block.get("data", {}).get("items", [])
                bullets = [
                    f"{it.get('label', '')}: {it.get('desc', '')}"
                    for it in items[:6]
                ] or [f"Key concepts in {topic_title}"]
                slide_title = block.get("sub_title", topic_title)
            else:
                bullets = [f"Key concepts in {topic_title}"]
                slide_title = topic_title

            # Try to use infographic image for padding slide
            if _all_infographic_images:
                img_info = _all_infographic_images[pad_i % len(_all_infographic_images)]
                add_infographic_slide(prs, slide_title, img_info["image_path"], img_info.get("caption", ""))
            else:
                add_content_slide(prs, slide_title, bullets)

        logger.info(
            f"SLIDE TARGET ENFORCEMENT: Padded {shortfall} slides. "
            f"Final count: {len(prs.slides)} (target: {slide_target})"
        )

    # ── CLOSING SLIDES ── (added AFTER padding so they appear at the end)
    from generate_slides.build_pptx import add_closing_slides
    add_closing_slides(prs, context)
    logger.info(f"Added closing slides. Total: {len(prs.slides)}")

    # Save the single PPTX
    course_title = context.get("Course_Title", "Course")
    safe_title = course_title.replace(":", "").replace("/", "-").replace(" ", "_")[:40]
    pptx_path = tempfile.mktemp(suffix=f"_{safe_title}_multi_agent.pptx")
    prs.save(pptx_path)
    total_slide_count = len(prs.slides)
    logger.info(f"Built single PPTX: {pptx_path} ({total_slide_count} slides)")

    return {
        "message": f"{total_slide_count} slides across {num_lus} LUs",
        "merged_pptx_path": pptx_path,
        "pptx_paths": [pptx_path],
        "lu_results": lu_results,
    }


def _build_lu_data_from_content_map(lu: dict, content_map: dict, lu_idx: int) -> dict:
    """Build infographic slide data directly from content_map when assembly fails.

    This is the critical fallback — ensures topics always have slides even when
    the skeleton/assembly pipeline produces empty data.
    """
    topics = lu.get("Topics", [])
    lo_num = lu.get("LO_Number", f"LO{lu_idx + 1}")
    lu_num = lu.get("LU_Number", f"LU{lu_idx + 1}")
    lu_title = lu.get("LU_Title", "")
    lo_title = lu.get("LO", "")
    topics_data = []

    def _fuzzy_content(title):
        """Fuzzy lookup in content_map."""
        if not content_map or not title:
            return {}
        if title in content_map:
            return content_map[title]
        norm = lambda s: s.lower().replace(" ", "").replace("_", "").replace("-", "").strip()
        nk = norm(title)
        for k, v in content_map.items():
            if norm(k) == nk:
                return v
        tl = title.lower().strip()
        for k, v in content_map.items():
            kl = k.lower().strip()
            if tl in kl or kl in tl:
                return v
        return {}

    for ti, t in enumerate(topics):
        t_title = t.get("Topic_Title", f"Topic {ti + 1}")
        t_content = _fuzzy_content(t_title)
        blocks = t_content.get("content_blocks", [])
        activity_data = t_content.get("activity", {})

        # Build infographic_slides from content blocks
        infographic_slides = []
        if blocks:
            for bi, block in enumerate(blocks):
                items = block.get("data", {}).get("items", [])
                fallback_bullets = [
                    f"{it.get('label', '')}: {it.get('desc', '')}"
                    for it in items[:6]
                ]
                infographic_slides.append({
                    "position": bi,
                    "title": block.get("sub_title", t_title),
                    "image_path": None,
                    "caption": block.get("caption", ""),
                    "fallback_bullets": fallback_bullets,
                })
        else:
            # No content blocks either — build from CP bullet points
            bps = t.get("Bullet_Points", [])
            if bps:
                chunk_size = 4
                for si in range(0, max(1, len(bps)), chunk_size):
                    chunk = bps[si:si + chunk_size]
                    infographic_slides.append({
                        "position": si // chunk_size,
                        "title": chunk[0][:40] if chunk else t_title,
                        "image_path": None,
                        "caption": "",
                        "fallback_bullets": chunk,
                    })
            else:
                infographic_slides.append({
                    "position": 0,
                    "title": t_title,
                    "image_path": None,
                    "caption": "",
                    "fallback_bullets": [f"Content for {t_title}"],
                })

        # Format activity
        activity_lines = []
        if isinstance(activity_data, dict):
            if activity_data.get("title"):
                activity_lines.append(f"Activity: {activity_data['title']}")
            for step in activity_data.get("steps", []):
                activity_lines.append(step)
        elif isinstance(activity_data, list):
            activity_lines = activity_data

        topics_data.append({
            "title": t_title,
            "topic_number": f"T{ti + 1}",
            "lo_number": lo_num,
            "lo_title": lo_title,
            "lu_number": lu_num,
            "lu_title": lu_title,
            "ref": "",
            "infographic_slides": infographic_slides,
            "activity": activity_lines,
        })

    return {
        "topics": topics_data,
        "lo_number": lo_num,
        "lo_title": lo_title,
        "lu_number": lu_num,
        "lu_title": lu_title,
    }



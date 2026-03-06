"""
build_pptx.py - Generate professional editable WSQ training PPTX slide decks.

Accepts dynamic course content (from Claude AI or extracted CP data) and builds
real editable PPTX files with text boxes, formatted bullets, logos, and copyright.

Uses slide_template.pptx with predefined layouts for consistent styling.
"""

import os
import logging
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE_PATH = os.path.join(BASE_DIR, ".claude", "skills", "generate_slides",
                              "templates", "slide_template.pptx")
LOGO_DIR = os.path.join(BASE_DIR, "assets", "slide_logos")
WSQ_LOGO = os.path.join(LOGO_DIR, "wsq_logo.png")
TERTIARY_LOGO = os.path.join(LOGO_DIR, "tertiary_infotech_logo.png")
CERT_TEMPLATE = os.path.join(LOGO_DIR, "certificate_template.png")
LETS_KNOW_IMG = os.path.join(LOGO_DIR, "lets_know_each_other.png")

# Colours (matching supervisor template)
DARK_NAVY = RGBColor(0x1B, 0x2A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x34, 0x98, 0xDB)
ACCENT_TEAL = RGBColor(0x1A, 0xBC, 0x9C)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)
RED = RGBColor(0xCC, 0x00, 0x00)

# Template uses Google Slides dimensions (10" x 5.625")
SLIDE_W = Emu(9144000)
SLIDE_H = Emu(5143500)

# Layout indices (from slide_template.pptx)
LY_BLANK = 0
LY_TITLE_BODY = 1
LY_COVER = 4
LY_SECTION = 5
LY_TWO_COL = 6
LY_TITLE_ONLY = 9

COPYRIGHT = "This material belongs to Tertiary Infotech Pte Ltd (UEN: 20120096W). All Rights Reserved"


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------

def _strip_template_footers(prs):
    """Remove ALL inherited copyright and page-number shapes from master + layouts.

    Must be called once after loading the template and before creating any slides.
    This ensures _add_copyright() is the ONLY source of copyright text.
    """
    def _should_remove(shape):
        if not shape.has_text_frame:
            return False
        txt = shape.text_frame.text.strip()
        txt_lower = txt.lower()
        # Copyright text
        if "all rights reserved" in txt_lower:
            return True
        # Page number placeholder (various unicode quote styles around #)
        stripped = ''.join(c for c in txt if c.isalnum() or c == '#')
        if stripped == '#' and len(txt) <= 5:
            return True
        return False

    for master in prs.slide_masters:
        for shape in list(master.shapes):
            if _should_remove(shape):
                shape._element.getparent().remove(shape._element)
    for layout in prs.slide_layouts:
        for shape in list(layout.shapes):
            if _should_remove(shape):
                shape._element.getparent().remove(shape._element)


def _add_copyright(slide):
    """Add single clean copyright footer to bottom of slide.

    A thin line separator + centered copyright text at the very bottom.
    Master/layout copyright shapes are removed by _strip_template_footers().
    """
    # Remove any slide-level copyright text boxes (from previous calls)
    to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip().lower()
            if "all rights reserved" in txt and "tertiary" in txt:
                to_remove.append(shape)
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # Thin separator line above copyright
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(300000), Emu(4920000), Emu(8544000), Emu(10000)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_GRAY
    line.line.fill.background()

    # Copyright text — centered, clean
    txBox = slide.shapes.add_textbox(Emu(100000), Emu(4940000), Emu(8944000), Emu(200000))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = COPYRIGHT
    p.font.size = Pt(7)
    p.font.bold = False
    p.font.color.rgb = GRAY
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER


def _set_ph_text(ph, text, size=None, bold=None, color=None):
    """Set placeholder text with formatting."""
    ph.text = text
    for para in ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Arial"
            if size:
                run.font.size = size
            if bold is not None:
                run.font.bold = bold
            if color:
                run.font.color.rgb = color


def _fill_body(ph, lines, size=Pt(14)):
    """Fill a body placeholder with bullet lines, well-spaced for readability."""
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        is_sub = line.startswith("  ")
        p.text = line.strip()
        p.font.name = "Arial"
        p.font.size = Pt(12) if is_sub else size
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        p.level = 1 if is_sub else 0


def _remove_all_slides(prs):
    """Remove all slides from presentation."""
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_cover(prs, course_title, tgs_code=""):
    """Cover slide with course title centered in the middle, logos and version info at bottom."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK])

    # --- Course title (vertically centered in the middle of the slide) ---
    title_h = Emu(1600000)
    title_top = (SLIDE_H - title_h) // 2  # Vertically centered
    title_box = slide.shapes.add_textbox(
        Emu(400000), title_top, SLIDE_W - Emu(800000), title_h
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = course_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    # --- WSQ logo (bottom-left) ---
    if os.path.exists(WSQ_LOGO):
        slide.shapes.add_picture(WSQ_LOGO, Emu(300000), Emu(3450000), height=Emu(800000))

    # --- Tertiary Infotech logo (below WSQ, fits within slide) ---
    if os.path.exists(TERTIARY_LOGO):
        slide.shapes.add_picture(TERTIARY_LOGO, Emu(300000), Emu(4300000), height=Emu(700000))

    # --- Version info (bottom-right, black text) ---
    info_text = "Version: 1.0"
    if tgs_code:
        info_text = f"Version: 1.0\nCourse Code: {tgs_code}"
    info_text += "\nWebsite: www.tertiarycourses.com.sg"
    info_box = slide.shapes.add_textbox(
        Emu(5500000), Emu(4100000), Emu(3500000), Emu(700000)
    )
    inf_tf = info_box.text_frame
    inf_tf.word_wrap = True
    for li, line in enumerate(info_text.split("\n")):
        if li == 0:
            ip = inf_tf.paragraphs[0]
        else:
            ip = inf_tf.add_paragraph()
        ip.text = line
        ip.font.size = Pt(12)
        ip.font.color.rgb = RGBColor(0, 0, 0)  # Black text
        ip.alignment = PP_ALIGN.RIGHT

    _add_copyright(slide)


def _add_title_only_slide(prs, title):
    """Slide with only a bold centered title — blank body for manual editing."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK])
    title_box = slide.shapes.add_textbox(Emu(151275), Emu(100000), Emu(8800000), Emu(600000))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    _add_copyright(slide)


def _add_title_image_slide(prs, title, image_path):
    """Slide with bold title at top and a full-width image below.

    Falls back to a text-only slide if the image file doesn't exist.
    """
    if not os.path.exists(image_path):
        add_tb_slide(prs, title, ["[Image placeholder — place image at: " + image_path + "]"])
        return
    slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK])
    # Title at top — consistent with other slide types
    title_box = slide.shapes.add_textbox(Emu(151275), Emu(100000), Emu(8800000), Emu(600000))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    # Image — fill remaining space below title
    img_top = Emu(700000)
    img_left = Emu(300000)
    img_width = Emu(8500000)
    img_height = Emu(4100000)
    slide.shapes.add_picture(image_path, img_left, img_top, img_width, img_height)
    _add_copyright(slide)


def add_tb_slide(prs, title, lines, font_size=Pt(14)):
    """Title + body slide with good spacing and teal accent bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_BODY])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            _set_ph_text(ph, title, size=Pt(20), bold=True, color=DARK_NAVY)
        elif ph.placeholder_format.idx == 1:
            _fill_body(ph, lines, font_size)
    # Accent bar under title
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(151275), Emu(572700), Emu(8000000), Emu(30000)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_TEAL
    bar.line.fill.background()
    _add_copyright(slide)


def add_section(prs, text):
    """Section header — centered title with dark navy background and accent bar."""
    import re as _re
    clean_text = _re.sub(r'\s*[\(\[]\s*[KA]\d+.*?[\)\]]', '', text).strip()
    clean_text = _re.sub(r'\s*[-–—]\s*[KA]\d+[\s,KA\d]*$', '', clean_text).strip()
    clean_text = _re.sub(r'\s*[KA]\d+\s*[,&]\s*[KA]\d+[\s,&KA\d]*$', '', clean_text).strip()

    slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK])

    # Dark navy background fill
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_NAVY

    # Accent bar at top
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, Emu(80000)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_TEAL
    bar.line.fill.background()

    # Centered title text
    title_box = slide.shapes.add_textbox(
        Emu(500000), Emu(1800000), Emu(8144000), Emu(1500000)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clean_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    # Bottom accent bar
    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(5063500), SLIDE_W, Emu(80000)
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = ACCENT_TEAL
    bar2.line.fill.background()


def add_content_slide(prs, title, bullets, ref_tag="", image_path=None):
    """Content slide — full-width text OR two-column (text + image).

    When image_path is provided: two-column layout (editable text LEFT, image RIGHT).
    When no image: full-width text layout (no empty right column).
    """
    # Clean title: remove ALL K/A references in any format
    import re as _re
    clean_title = _re.sub(r'\s*[\(\[]\s*[KA]\d+.*?[\)\]]', '', title).strip()
    clean_title = _re.sub(r'\s*[-–—]\s*[KA]\d+[\s,KA\d]*$', '', clean_title).strip()
    clean_title = _re.sub(r'\s*[KA]\d+\s*[,&]\s*[KA]\d+[\s,&KA\d]*$', '', clean_title).strip()
    full_title = clean_title

    if image_path and os.path.exists(image_path):
        # --- Two-column: editable text LEFT + real image RIGHT ---
        slide = prs.slides.add_slide(prs.slide_layouts[LY_TWO_COL])
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                _set_ph_text(ph, full_title, size=Pt(18), bold=True, color=DARK_NAVY)
            elif ph.placeholder_format.idx == 1:
                _fill_body(ph, bullets, Pt(12))
            elif ph.placeholder_format.idx == 2:
                col_left = ph.left
                col_top = ph.top
                col_w = ph.width
                col_h = ph.height
                sp = ph._element
                sp.getparent().remove(sp)

                try:
                    from PIL import Image as _PILImg
                    with _PILImg.open(image_path) as im:
                        iw, ih = im.size
                    img_ratio = iw / ih
                    col_ratio = col_w / col_h
                    if img_ratio > col_ratio:
                        fit_w = col_w
                        fit_h = int(col_w / img_ratio)
                    else:
                        fit_h = col_h
                        fit_w = int(col_h * img_ratio)
                    x = col_left + (col_w - fit_w) // 2
                    y = col_top + (col_h - fit_h) // 2
                    slide.shapes.add_picture(image_path, Emu(x), Emu(y), Emu(fit_w), Emu(fit_h))
                except ImportError:
                    slide.shapes.add_picture(image_path, col_left, col_top, col_w, col_h)
        # Accent bar under title — consistent with other slide types
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(151275), Emu(572700), Emu(8000000), Emu(30000)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT_TEAL
        bar.line.fill.background()
    else:
        # --- Full-width: professional text content with accent bar ---
        slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_BODY])
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                _set_ph_text(ph, full_title, size=Pt(18), bold=True, color=DARK_NAVY)
            elif ph.placeholder_format.idx == 1:
                _fill_body(ph, bullets, Pt(14))

        # Accent bar under title — professional visual element
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(151275), Emu(572700), Emu(8000000), Emu(30000)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT_TEAL
        bar.line.fill.background()

    _add_copyright(slide)


def add_activity(prs, topic_short, steps):
    """Activity slide with teal accent bar and structured format.

    Steps should include scenario/objective/steps/output/duration.
    Title is never truncated.
    """
    # Ensure full title (no truncation)
    full_title = f"Activity: {topic_short}"
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_BODY])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            # Use smaller font if title is long to avoid truncation
            title_size = Pt(15) if len(full_title) > 60 else Pt(18)
            _set_ph_text(ph, full_title, size=title_size, bold=True, color=DARK_NAVY)
        elif ph.placeholder_format.idx == 1:
            _fill_body(ph, steps, Pt(14))
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(151275), Emu(572700), Emu(8000000), Emu(30000)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_TEAL
    bar.line.fill.background()
    _add_copyright(slide)


# ---------------------------------------------------------------------------
# Diagram / Visual slide builders (editable shapes)
# ---------------------------------------------------------------------------

# Diagram colour palette (professional, matches existing theme)
_DIAGRAM_COLORS = [
    RGBColor(0x34, 0x98, 0xDB),  # Blue
    RGBColor(0x1A, 0xBC, 0x9C),  # Teal
    RGBColor(0xE7, 0x4C, 0x3C),  # Red
    RGBColor(0xF3, 0x9C, 0x12),  # Orange
    RGBColor(0x9B, 0x59, 0xB6),  # Purple
    RGBColor(0x27, 0xAE, 0x60),  # Green
    RGBColor(0x2C, 0x3E, 0x50),  # Dark blue
    RGBColor(0xE6, 0x7E, 0x22),  # Dark orange
]


def _add_diagram_title(slide, title, ref_tag=""):
    """Add a prominent title to a diagram slide, aligned within slide bounds."""
    import re as _re
    clean_title = _re.sub(r'\s*[\(\[]\s*[KA]\d+.*?[\)\]]', '', title).strip()
    # Also strip standalone K/A refs like "K3, A2" at end of title
    clean_title = _re.sub(r'\s*[-–—]\s*[KA]\d+[\s,KA\d]*$', '', clean_title).strip()
    ttl = slide.shapes.add_textbox(Emu(151275), Emu(60000), Emu(8800000), Emu(500000))
    tf = ttl.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clean_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY
    p.font.name = "Arial"


def add_process_flow_slide(prs, title, steps, ref_tag=""):
    """Process flow diagram — large horizontal boxes connected by arrows.

    Use for: How It Works, Step-by-step processes, Workflows.
    steps: list of strings (3-6 items work best).
    All elements fit within slide bounds (10" x 5.625").
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK])
    _add_diagram_title(slide, title, ref_tag)

    n = min(len(steps), 6)
    if n == 0:
        _add_copyright(slide)
        return

    # Content area: x=0.2" to 9.8", y=0.6" to 5.1" — bigger boxes, more visible
    margin_x = 183000   # 0.2"
    usable_w = 8780000  # 9.6" in EMU
    arrow_w = 180000
    gap = 60000
    box_w = (usable_w - (n - 1) * (gap + arrow_w + gap)) // n
    box_h = 900000      # Much taller (0.98" vs 0.7")
    # Vertically center in content area
    content_mid_y = (650000 + 4800000) // 2
    center_y = content_mid_y - box_h // 2

    for i, step in enumerate(steps[:n]):
        x = margin_x + i * (box_w + gap + arrow_w + gap)
        color = _DIAGRAM_COLORS[i % len(_DIAGRAM_COLORS)]

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(center_y), Emu(box_w), Emu(box_h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()  # No border — clean look

        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.margin_top = Emu(80000)
        tf.margin_bottom = Emu(80000)
        tf.margin_left = Emu(50000)
        tf.margin_right = Emu(50000)
        from pptx.enum.text import MSO_ANCHOR
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(11) if len(step) > 35 else (Pt(13) if len(step) > 20 else Pt(14))
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER

        # Step number below box
        num_box = slide.shapes.add_textbox(
            Emu(x), Emu(center_y + box_h + 50000), Emu(box_w), Emu(220000)
        )
        num_p = num_box.text_frame.paragraphs[0]
        num_p.text = f"Step {i + 1}"
        num_p.font.size = Pt(9)
        num_p.font.color.rgb = GRAY
        num_p.font.name = "Arial"
        num_p.alignment = PP_ALIGN.CENTER

        # Small chevron arrow between boxes
        if i < n - 1:
            arrow_x = x + box_w + gap
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON, Emu(arrow_x),
                Emu(center_y + box_h // 2 - 80000),
                Emu(arrow_w), Emu(160000)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()

    _add_copyright(slide)


def add_comparison_slide(prs, title, items, ref_tag=""):
    """Grid of large colored boxes — for Key Components, Categories, Comparisons.

    items: list of dicts [{"label": "Title", "desc": "Description"}, ...]
    All elements fit within slide bounds (10" x 5.625").
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK])
    _add_diagram_title(slide, title, ref_tag)

    n = min(len(items), 6)
    if n == 0:
        _add_copyright(slide)
        return

    # Content area: x=0.2" to 9.8", y=0.65" to 5.1" — maximized
    margin_x = 183000
    usable_w = 8780000
    content_top = 650000
    content_bot = 4750000
    usable_h = content_bot - content_top

    cols = 3 if n > 4 else (2 if n > 1 else 1)
    rows = (n + cols - 1) // cols
    gap_x = 100000
    gap_y = 100000
    box_w = (usable_w - (cols - 1) * gap_x) // cols
    box_h = min(1300000, (usable_h - (rows - 1) * gap_y) // rows)
    grid_w = cols * box_w + (cols - 1) * gap_x
    grid_h = rows * box_h + (rows - 1) * gap_y
    start_x = margin_x + (usable_w - grid_w) // 2
    start_y = content_top + (usable_h - grid_h) // 2

    for i, item in enumerate(items[:n]):
        row = i // cols
        col = i % cols
        x = start_x + col * (box_w + gap_x)
        y = start_y + row * (box_h + gap_y)
        color = _DIAGRAM_COLORS[i % len(_DIAGRAM_COLORS)]

        label = item.get("label", item) if isinstance(item, dict) else str(item)
        desc = item.get("desc", "") if isinstance(item, dict) else ""

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(y), Emu(box_w), Emu(box_h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.margin_top = Emu(80000)
        tf.margin_left = Emu(60000)
        tf.margin_right = Emu(60000)

        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER

        if desc:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(6)
            p2.text = desc
            p2.font.size = Pt(10)
            p2.font.color.rgb = WHITE
            p2.font.name = "Arial"
            p2.alignment = PP_ALIGN.CENTER

    _add_copyright(slide)


def add_cycle_slide(prs, title, stages, center_text="", ref_tag=""):
    """Circular cycle diagram — for frameworks, lifecycles, iterative processes.

    stages: list of strings (3-6 items).
    center_text: optional text in the center circle.
    All elements fit within slide bounds (10" x 5.625").
    """
    import math

    slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK])
    _add_diagram_title(slide, title, ref_tag)

    n = min(len(stages), 6)
    if n == 0:
        _add_copyright(slide)
        return

    # Content area center — bigger radius and nodes for visibility
    cx = SLIDE_W // 2                    # 4572000 = 5.0"
    cy = (650000 + 4750000) // 2         # ~2700000 = ~2.95"
    radius = 1500000                     # 1.64" — bigger circle
    node_w, node_h = 1700000, 560000     # bigger nodes for readability

    if center_text:
        center = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Emu(cx - 650000), Emu(cy - 400000),
            Emu(1300000), Emu(800000)
        )
        center.fill.solid()
        center.fill.fore_color.rgb = DARK_NAVY
        center.line.fill.background()
        ctf = center.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = center_text
        cp.font.size = Pt(12)
        cp.font.bold = True
        cp.font.color.rgb = WHITE
        cp.font.name = "Arial"
        cp.alignment = PP_ALIGN.CENTER

    for i, stage in enumerate(stages[:n]):
        angle = -math.pi / 2 + (2 * math.pi * i / n)
        nx = cx + int(radius * math.cos(angle)) - node_w // 2
        ny = cy + int(radius * math.sin(angle)) - node_h // 2

        # Clamp within slide bounds (with 0.15" margin)
        nx = max(140000, min(nx, SLIDE_W - node_w - 140000))
        ny = max(580000, min(ny, SLIDE_H - node_h - 280000))

        color = _DIAGRAM_COLORS[i % len(_DIAGRAM_COLORS)]

        node = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(nx), Emu(ny), Emu(node_w), Emu(node_h)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = color
        node.line.fill.background()

        ntf = node.text_frame
        ntf.word_wrap = True
        ntf.auto_size = None
        ntf.margin_top = Emu(40000)
        ntf.margin_left = Emu(40000)
        ntf.margin_right = Emu(40000)
        np_ = ntf.paragraphs[0]
        np_.text = stage
        np_.font.size = Pt(10) if len(stage) > 25 else Pt(11)
        np_.font.bold = True
        np_.font.color.rgb = WHITE
        np_.font.name = "Arial"
        np_.alignment = PP_ALIGN.CENTER

        # Arrow connector between nodes
        if n > 1:
            next_angle = -math.pi / 2 + (2 * math.pi * ((i + 1) % n) / n)
            mid_angle = (angle + next_angle) / 2
            if i == n - 1:
                mid_angle = angle + (2 * math.pi / n) / 2
            ax = cx + int((radius - 200000) * math.cos(mid_angle)) - 80000
            ay = cy + int((radius - 200000) * math.sin(mid_angle)) - 80000
            ax = max(140000, min(ax, SLIDE_W - 300000))
            ay = max(580000, min(ay, SLIDE_H - 360000))
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Emu(ax), Emu(ay), Emu(160000), Emu(160000)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = ACCENT_TEAL
            dot.line.fill.background()

    _add_copyright(slide)


def add_infographic_slide(prs, title, image_path, caption=""):
    """Infographic image slide matching the approved reference PPTX layout.

    Layout (10" x 5.62" slide):
      - Title textbox at top (0.17", 0.11") — 9.62" wide
      - Thin separator line at (0.17", 0.57")
      - Centered infographic image at (1.15", 0.68") — 7.66" x 4.37"
      - Source caption at (0.17", 5.14")
      - Copyright footer at bottom
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LY_BLANK])

    # ---- Title textbox ----
    title_left = Emu(155000)    # 0.17"
    title_top = Emu(100000)     # 0.11"
    title_w = Emu(8796000)      # 9.62"
    title_h = Emu(400000)       # 0.44"
    title_box = slide.shapes.add_textbox(title_left, title_top, title_w, title_h)
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.text = title or ""
    title_p.font.size = Pt(18)
    title_p.font.bold = True
    title_p.font.color.rgb = DARK_NAVY

    # ---- Thin separator line ----
    sep_top = Emu(520000)       # 0.57"
    sep_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, title_left, sep_top, title_w, Emu(18000)
    )
    sep_shape.fill.solid()
    sep_shape.fill.fore_color.rgb = ACCENT_TEAL
    sep_shape.line.fill.background()

    # ---- Infographic image (centered) ----
    img_left = Emu(1051000)     # 1.15"
    img_top = Emu(621000)       # 0.68"
    img_area_w = Emu(7004000)   # 7.66"
    img_area_h = Emu(3996000)   # 4.37"

    if image_path and os.path.exists(image_path):
        try:
            from PIL import Image as _PILImg
            with _PILImg.open(image_path) as im:
                iw, ih = im.size
            img_ratio = iw / ih
            area_ratio = img_area_w / img_area_h
            if img_ratio > area_ratio:
                fit_w = img_area_w
                fit_h = int(img_area_w / img_ratio)
            else:
                fit_h = img_area_h
                fit_w = int(img_area_h * img_ratio)
            x = img_left + (img_area_w - fit_w) // 2
            y = img_top + (img_area_h - fit_h) // 2
            slide.shapes.add_picture(image_path, x, y, fit_w, fit_h)
        except ImportError:
            slide.shapes.add_picture(
                image_path, img_left, img_top, img_area_w, img_area_h
            )

    # ---- Source caption ----
    if caption:
        cap_top = Emu(4700000)  # 5.14"
        cap_box = slide.shapes.add_textbox(
            title_left, cap_top, title_w, Emu(247000)
        )
        cap_tf = cap_box.text_frame
        cap_tf.word_wrap = True
        cap_p = cap_tf.paragraphs[0]
        cap_p.text = caption
        cap_p.font.size = Pt(9)
        cap_p.font.italic = True
        cap_p.font.color.rgb = GRAY
        cap_p.alignment = PP_ALIGN.LEFT

    _add_copyright(slide)


def add_diagram_slide(prs, title, diagram_type, items, ref_tag="", center_text=""):
    """Dispatch to the appropriate diagram builder.

    diagram_type: "process" | "comparison" | "cycle"
    items: depends on type — list of strings or list of dicts
    """
    if diagram_type == "process":
        add_process_flow_slide(prs, title, items, ref_tag)
    elif diagram_type == "comparison":
        add_comparison_slide(prs, title, items, ref_tag)
    elif diagram_type == "cycle":
        add_cycle_slide(prs, title, items, center_text, ref_tag)
    else:
        # Default to comparison grid
        add_comparison_slide(prs, title, items, ref_tag)


def add_certificate(prs, tsc_code="", course_title=""):
    """Certificate of Accomplishment slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_BODY])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            _set_ph_text(ph, "Certificate of Accomplishment", size=Pt(20), bold=True)
        elif ph.placeholder_format.idx == 1:
            lines = [
                "Upon successful completion you will receive:",
                "",
            ]
            if tsc_code:
                lines.append(f"  WSQ Statement of Attainment (SOA) - {tsc_code}")
            else:
                lines.append("  WSQ Statement of Attainment (SOA)")
            lines += [
                "  Certificate from Tertiary Infotech Academy Pte Ltd",
                "",
                "Requirements:",
                "  Minimum 75% attendance",
                "  Achieve 'Competent' in assessment",
                "  Complete TRAQOM survey",
            ]
            _fill_body(ph, lines, Pt(12))
    if os.path.exists(CERT_TEMPLATE):
        slide.shapes.add_picture(CERT_TEMPLATE, Emu(6000000), Emu(700000), height=Emu(3800000))
    _add_copyright(slide)


# ---------------------------------------------------------------------------
# Intro & closing slide sets (dynamic from context)
# ---------------------------------------------------------------------------

def add_intro_slides(prs, context):
    """Add standard intro slides using course context data.

    Reads from the CP interpreter's output schema:
    - TGS_Ref_No (not TGS_Number)
    - Learning_Units[].LO, LO_Description
    - Learning_Units[].K_numbering_description, A_numbering_description
    - TSC_Description / Proficiency_Description
    """
    course_title = context.get('Course_Title', 'Course')
    tgs_code = context.get('TGS_Ref_No', '') or context.get('TGS_Number', '') or context.get('Course_Code', '')
    tsc_code = context.get('TSC_Code', '')
    tsc_title = context.get('TSC_Title', '')
    lus = context.get('Learning_Units', [])

    add_cover(prs, course_title, tgs_code)

    # Digital Attendance (Mandatory)
    add_tb_slide(prs, "Digital Attendance (Mandatory)", [
        "It is mandatory for you to take both AM, PM and Assessment digital attendance for WSQ funded courses.",
        "",
        "The trainer or administrator will show you the digital attendance QR code generated from SSG portal.",
        "",
        "Please scan the QR code from your mobile phone camera and submit your attendance.",
    ])

    # About the Trainer — plain slide, trainer fills in manually
    _add_title_only_slide(prs, "About the Trainer")

    # Let's Know Each Other — title + icebreaker image
    _add_title_image_slide(prs, "Let's Know Each Other...", LETS_KNOW_IMG)

    # Ground Rules
    add_tb_slide(prs, "Ground Rules", [
        "Set your mobile phone to silent mode",
        "Actively participate in the class. No question is stupid.",
        "Respect each other views. Agree to disagree.",
        "Each person should only speak one at a time.",
        "Be punctual. Back from breaks on time.",
        "Exit the class silently if you need to step out for phone call, toilet break etc.",
        "75% attendance is required for WSQ funding eligibility.",
    ])

    # Skills Framework — TSC info + LOs per LU in structured format
    sf_lines = []
    if tsc_title:
        sf_lines.append(f"TSC Title: {tsc_title}")
    if tsc_code:
        sf_lines.append(f"TSC Code: {tsc_code}")
    tsc_desc = (context.get('TSC_Description', '')
                or context.get('Proficiency_Description', ''))
    if tsc_desc:
        sf_lines += ["", f"TSC Description: {tsc_desc[:300]}"]
    prof_level = context.get('Proficiency_Level', '')
    if prof_level:
        sf_lines.append(f"Proficiency Level: {prof_level}")

    # Add Learning Outcomes per LU directly into Skills Framework
    sf_lines.append("")
    sf_lines.append("Learning Outcomes:")
    for lu_i, lu in enumerate(lus):
        lu_num = lu.get('LU_Number', '') or f'LU{lu_i + 1}'
        lu_title = lu.get('LU_Title', '')
        lo = lu.get('LO', '')
        lo_desc = lu.get('LO_Description', '')
        # Truncate long LO descriptions to prevent slide overflow
        if lo_desc and len(lo_desc) > 120:
            lo_desc = lo_desc[:117] + "..."
        lo_text = f"{lo}: {lo_desc}" if lo_desc else lo
        sf_lines.append(f"  {lu_num} ({lu_title})")
        sf_lines.append(f"    {lo_text}")
    if sf_lines:
        add_tb_slide(prs, "Skills Framework", sf_lines, Pt(10))

    # K&A Statements — extracted from Learning Units (CP interpreter schema)
    ka_lines = []
    for lu_i, lu in enumerate(lus):
        lu_num = lu.get('LU_Number', '') or f'LU{lu_i + 1}'
        lu_title = lu.get('LU_Title', '')
        k_list = lu.get('K_numbering_description', [])
        a_list = lu.get('A_numbering_description', [])
        if k_list or a_list:
            ka_lines.append(f"{lu_num}: {lu_title}")
            for k in k_list:
                if isinstance(k, dict):
                    desc = k.get('Description', '')
                    if len(desc) > 150:
                        desc = desc[:147] + "..."
                    ka_lines.append(f"  {k.get('K_number', '')}: {desc}")
                elif isinstance(k, str):
                    ka_lines.append(f"  {k[:150]}")
            for a in a_list:
                if isinstance(a, dict):
                    desc = a.get('Description', '')
                    if len(desc) > 150:
                        desc = desc[:147] + "..."
                    ka_lines.append(f"  {a.get('A_number', '')}: {desc}")
                elif isinstance(a, str):
                    ka_lines.append(f"  {a[:150]}")
            ka_lines.append("")
    if ka_lines:
        # Always fit on ONE slide — shrink font for long K&A lists
        non_empty_ka = [l for l in ka_lines if l.strip()]
        if len(non_empty_ka) > 20:
            ka_font = Pt(7)
        elif len(non_empty_ka) > 14:
            ka_font = Pt(8)
        else:
            ka_font = Pt(10)
        add_tb_slide(prs, "Knowledge & Ability Statements", ka_lines, ka_font)

    # Course Outline (NO K/A references — clean titles only)
    outline_lines = []
    for lu_idx_o, lu in enumerate(lus):
        lu_num = lu.get('LU_Number', '') or f"LU{lu_idx_o + 1}"
        lu_title = lu.get('LU_Title', '')
        topics = lu.get('Topics', [])
        outline_lines.append(f"{lu_num}: {lu_title}")
        for ti, t in enumerate(topics):
            outline_lines.append(f"  T{ti+1}: {t.get('Topic_Title', '')}")
        outline_lines.append("")
    if outline_lines:
        # Always one page — shrink font to fit all topics
        non_empty = [l for l in outline_lines if l.strip()]
        if len(non_empty) > 16:
            font_sz = Pt(7)
        elif len(non_empty) > 12:
            font_sz = Pt(8)
        elif len(non_empty) > 8:
            font_sz = Pt(9)
        else:
            font_sz = Pt(10)
        add_tb_slide(prs, "Course Outline", outline_lines, font_sz)

    # Assessment info — always shown (default text if not in CP)
    assessment_method = (context.get('Assessment_Method', '')
                         or context.get('Mode_of_Assessment', '')
                         or 'Written Assessment')
    add_tb_slide(prs, "Assessment Methods & Briefing", [
        f"Assessment: {assessment_method}",
        "",
        "Assessment format: Open Book",
        "Open book assessment allows you to reference your learning materials.",
        "",
        "Duration: As per course schedule",
        "Grading: Competent / Not Yet Competent",
        "",
        "Assessment Rules:",
        "  Place phones & other materials under the table",
        "  No photos or recording of assessment scripts",
        "  No discussion with other learners during assessment",
        "  Raise your hand if you have any questions",
    ])

    add_tb_slide(prs, "Criteria for Funding", [
        "Minimum attendance rate of 75% based on SSG Digital Attendance record.",
        "Complete the assessment and be assessed as 'Competent'.",
        "Complete the TRAQOM survey.",
        "",
        "For more information on WSQ funding:",
        "Visit SkillsFuture portal: www.skillsfuture.gov.sg",
        "",
        "Eligible individuals may use SkillsFuture Credit to offset course fees.",
    ])


def add_closing_slides(prs, context):
    """Add standard closing slides."""
    tsc_code = context.get('TSC_Code', '')

    add_section(prs, "Summary & Q&A")

    add_tb_slide(prs, "TRAQOM Survey", [
        "Access the survey here",
        "Key in your last four NRIC/FIN characters and the six-digit course run ID",
        "",
        "Your feedback helps us improve training quality",
        "The survey is mandatory for WSQ-funded courses",
        "Takes approximately 5-10 minutes",
        "All responses are confidential",
    ])

    add_certificate(prs, tsc_code, context.get('Course_Title', ''))

    add_tb_slide(prs, "Digital Attendance", [
        "It is mandatory for you to take both AM, PM and Assessment digital attendance.",
        "",
        "Please scan the QR code from your mobile phone camera and submit your attendance.",
        "",
        "Ensure your attendance is recorded for all course days.",
        "This is required for funding and certification purposes.",
    ])

    add_section(prs, "Final Assessment")

    add_tb_slide(prs, "Support", [
        "If you have any enquiries during and after the class, you can contact us below",
        "",
        "  Email: enquiry@tertiaryinfotech.com",
        "  Tel: +65 6318 4588",
        "  Website: www.tertiarycourses.com.sg",
    ])

    add_section(prs, "Thank You!")


# ---------------------------------------------------------------------------
# Build topic slides from Claude-generated content
# ---------------------------------------------------------------------------

def build_infographic_topic_slides(prs, topic_data, topic_idx=0, lu_label=""):
    """Build slides for a topic using infographic images ONLY (no text bullets).

    Used in multi-agent infographic_mode. Each topic's content slides are
    full-width infographic images. Activity slide remains as text.

    topic_data = {
        "title": "Topic Title",
        "ref": "",
        "infographic_slides": [
            {
                "position": 0,
                "title": "Slide Sub-Title",
                "image_path": "/path/to/infographic.png" or None,
                "caption": "Source: ...",
                "fallback_bullets": ["Point 1", "Point 2", ...]
            },
            ...
        ],
        "activity": ["Step 1", "Step 2", ...]
    }
    """
    import re as _re
    title = topic_data.get("title", "Topic")
    title = _re.sub(r'\s*[\(\[]\s*[KA]\d+.*?[\)\]]', '', title).strip()
    title = _re.sub(r'\s*[-–—]\s*[KA]\d+[\s,KA\d]*$', '', title).strip()
    title = _re.sub(r'\s*[KA]\d+\s*[,&]\s*[KA]\d+[\s,&KA\d]*$', '', title).strip()

    # Section header: LO | LU | T format (matching approved PPTX)
    t_num = topic_data.get("topic_number", f"T{topic_idx + 1}")
    lo_num = topic_data.get("lo_number", "")
    lu_num = topic_data.get("lu_number", lu_label)
    if lo_num and lu_num:
        section_title = f"{lo_num} | {lu_num} | {t_num}: {title}"
    elif lu_num:
        section_title = f"{lu_num} | {t_num}: {title}"
    else:
        section_title = f"{t_num}: {title}"
    add_section(prs, section_title)

    # Infographic content slides
    infographic_slides = topic_data.get("infographic_slides", [])
    for info_slide in infographic_slides:
        slide_title = info_slide.get("title", title)
        image_path = info_slide.get("image_path")
        caption = info_slide.get("caption", "")
        fallback_bullets = info_slide.get("fallback_bullets", [])

        if image_path and os.path.exists(image_path):
            # Full-width infographic image slide
            add_infographic_slide(prs, slide_title, image_path, caption)
        elif fallback_bullets:
            # Text fallback when infographic generation failed
            add_content_slide(prs, slide_title, fallback_bullets)
        else:
            # Minimal fallback
            add_content_slide(prs, slide_title, [f"Content for: {slide_title}"])

    # Activity slide
    activity = topic_data.get("activity", [])
    if activity:
        add_activity(prs, title, activity)


def build_topic_slides(prs, topic_data, image_paths=None, topic_idx=0, lu_label=""):
    """Build slides for a single topic from structured data.

    topic_data = {
        "title": "Topic Title",
        "ref": "K3, A2",
        "slides": [
            {"title": "Slide Title", "bullets": ["point 1", "point 2", ...]},
            ...
        ],
        "diagram": {
            "type": "process" | "comparison" | "cycle",
            "title": "Diagram Title",
            "items": ["Step 1", "Step 2", ...] or [{"label": "X", "desc": "Y"}, ...],
            "center_text": ""  (optional, for cycle)
        },
        "activity": ["Step 1", "Step 2", ...]
    }

    image_paths: When provided, used as slide images (two-column layout).
                 When None, full-width text layout with editable diagrams.
    topic_idx: 0-based index of this topic within the LU.
               Activity slides are only added for every 2nd topic (odd indices: 1, 3, 5...).
    """
    import re as _re
    title = topic_data.get("title", "Topic")
    # Clean ALL K/A references from topic title (any format)
    title = _re.sub(r'\s*[\(\[]\s*[KA]\d+.*?[\)\]]', '', title).strip()
    title = _re.sub(r'\s*[-–—]\s*[KA]\d+[\s,KA\d]*$', '', title).strip()
    title = _re.sub(r'\s*[KA]\d+\s*[,&]\s*[KA]\d+[\s,&KA\d]*$', '', title).strip()
    ref = ""  # Never show K/A refs on slides
    slides = topic_data.get("slides", [])
    activity = topic_data.get("activity", [])
    diagram = topic_data.get("diagram")

    # Section header for topic — includes LU and topic number
    topic_num = f"T{topic_idx + 1}"
    if lu_label:
        section_title = f"{lu_label} | {topic_num}: {title}"
    else:
        section_title = f"{topic_num}: {title}"
    add_section(prs, section_title)

    # Content slides — with images (two-column) or full-width text
    # When image_paths are provided (from NotebookLM), pair them with content slides
    img_list = image_paths if image_paths else []
    diagram_inserted = False
    for i, slide_data in enumerate(slides):
        s_title = slide_data.get("title", "")
        s_bullets = slide_data.get("bullets", [])
        img_path = img_list[i] if i < len(img_list) else None
        if s_title and s_bullets:
            add_content_slide(prs, s_title, s_bullets, image_path=img_path)

        # Insert diagram after the 3rd content slide (middle of topic)
        if i == 2 and not diagram_inserted:
            if diagram and isinstance(diagram, dict):
                d_type = diagram.get("type", "comparison")
                d_title = diagram.get("title", f"{title} -- Overview")
                d_items = diagram.get("items", [])
                d_center = diagram.get("center_text", "")
                if d_items:
                    add_diagram_slide(prs, d_title, d_type, d_items, center_text=d_center)
                    diagram_inserted = True
            elif slides:
                # Auto-generate a process flow from slide titles
                slide_titles = [s.get("title", "") for s in slides if s.get("title")]
                if len(slide_titles) >= 3:
                    steps = slide_titles[:5]
                    add_process_flow_slide(prs, f"{title} -- Overview", steps)
                    diagram_inserted = True

    # If diagram wasn't inserted (fewer than 3 content slides), add it now
    if not diagram_inserted:
        if diagram and isinstance(diagram, dict):
            d_type = diagram.get("type", "comparison")
            d_title = diagram.get("title", f"{title} — Overview")
            d_items = diagram.get("items", [])
            d_center = diagram.get("center_text", "")
            if d_items:
                add_diagram_slide(prs, d_title, d_type, d_items, center_text=d_center)
                diagram_inserted = True

        # ALWAYS generate a diagram if Claude didn't provide one — auto-generate from slide titles
        if not diagram_inserted and slides:
            slide_titles = [s.get("title", "") for s in slides if s.get("title")]
            if len(slide_titles) >= 3:
                steps = slide_titles[:5]
                add_process_flow_slide(prs, f"{title} — Key Concepts", steps)
            elif len(slide_titles) >= 1:
                # Even with few slides, create a comparison grid from bullet points
                items = []
                for s in slides[:4]:
                    t = s.get("title", "")
                    b = s.get("bullets", [])
                    if t:
                        items.append({"label": t, "desc": b[0] if b else ""})
                if items:
                    add_comparison_slide(prs, f"{title} — Summary", items)

    # Activity slide for each topic
    if activity:
        add_activity(prs, title, activity)


# ---------------------------------------------------------------------------
# Main: Build LU deck from dynamic content
# ---------------------------------------------------------------------------

def build_lu_deck(context, lu_idx, slides_data, is_first=False, is_last=False,
                  image_paths=None, images_per_topic=None, infographic_mode=False,
                  prs=None):
    """Build an editable PPTX deck for a Learning Unit.

    Args:
        context: Extracted course info dict.
        lu_idx: Index of the LU in context['Learning_Units'].
        slides_data: Claude-generated content dict with 'topics' list.
        is_first: Whether this is the first LU (adds intro slides).
        is_last: Whether this is the last LU (adds closing slides).
        image_paths: Optional list of image paths for the whole LU.
        images_per_topic: Optional dict mapping topic index -> list of image paths.
                          When provided, each topic gets ONLY its own images.
        infographic_mode: When True, use build_infographic_topic_slides() for topics
                          (full-width infographic images instead of text bullets).
        prs: Optional existing Presentation object. When provided, slides are added
             to this presentation instead of creating a new one. This allows building
             all LUs into a single PPTX without merging.

    Returns:
        Tuple of (pptx_path, slide_count) when prs is None (creates new file).
        Tuple of (None, slides_added) when prs is provided (caller manages saving).
    """
    lus = context.get('Learning_Units', [])
    lu = lus[lu_idx] if lu_idx < len(lus) else {}
    lu_num = lu.get('LU_Number', f'LU{lu_idx + 1}')
    lu_title = lu.get('LU_Title', 'Learning Unit')

    owns_prs = prs is None
    if owns_prs:
        # Use template if available, else blank
        if os.path.exists(TEMPLATE_PATH):
            prs = Presentation(TEMPLATE_PATH)
            _remove_all_slides(prs)
            _strip_template_footers(prs)  # Remove master/layout copyright & page numbers
        else:
            prs = Presentation()
            prs.slide_width = SLIDE_W
            prs.slide_height = SLIDE_H

    slides_before = len(prs.slides)

    # Intro slides (first LU only)
    if is_first:
        add_intro_slides(prs, context)

    # LO/LU tracking info (used by topic section headers)
    _lu_lo = slides_data.get("lo_number", "") or lu.get('LO', '')
    _lu_lo_title = slides_data.get("lo_title", "")

    # Skip LU section header and LU Overview slides — topics have their own LO|LU|T headers
    topics = slides_data.get("topics", [])

    if infographic_mode:
        # Infographic mode: full-width image slides for topic content
        for ti, topic_data in enumerate(topics):
            build_infographic_topic_slides(prs, topic_data, topic_idx=ti, lu_label=lu_num)
    else:
        # Standard mode: text bullets with optional images
        topic_image_lists = [None] * len(topics)

        if images_per_topic:
            for ti in range(len(topics)):
                topic_imgs = images_per_topic.get(ti, [])
                if topic_imgs:
                    topic_image_lists[ti] = topic_imgs
        elif image_paths:
            img_idx = 0
            for ti, topic_data in enumerate(topics):
                n_slides = len([s for s in topic_data.get("slides", [])
                               if s.get("title") and s.get("bullets")])
                if n_slides > 0:
                    topic_imgs = []
                    for _ in range(n_slides):
                        topic_imgs.append(image_paths[img_idx % len(image_paths)])
                        img_idx += 1
                    topic_image_lists[ti] = topic_imgs

        for ti, topic_data in enumerate(topics):
            build_topic_slides(prs, topic_data, image_paths=topic_image_lists[ti],
                               topic_idx=ti, lu_label=lu_num)

    # Closing slides (last LU only)
    if is_last:
        add_closing_slides(prs, context)

    slides_added = len(prs.slides) - slides_before

    if owns_prs:
        # Save to temp file (original behavior when no prs passed)
        course_title = context.get('Course_Title', 'Course')
        safe_course = course_title.replace(':', '').replace('/', '-').replace(' ', '_')[:40]
        safe_lu = lu_num.replace(':', '').replace('/', '-').replace(' ', '_')

        pptx_path = tempfile.mktemp(suffix=f"_{safe_lu}.pptx")
        prs.save(pptx_path)
        slide_count = len(prs.slides)
        logger.info(f"Built editable PPTX: {pptx_path} ({slide_count} slides)")
        return pptx_path, slide_count
    else:
        # Caller owns the Presentation — just return slide count added
        logger.info(f"Added {slides_added} slides for {lu_num}")
        return None, slides_added

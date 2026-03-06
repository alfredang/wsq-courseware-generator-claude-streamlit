"""Add activity slides after each topic in the existing PPTX."""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

INPUT = r"C:\Users\Afzaana Jaffar\Downloads\TGS-2025056191.pptx"
OUTPUT = r"C:\Users\Afzaana Jaffar\Downloads\TGS-2025056191_with_activities.pptx"

ACCENT_TEAL = RGBColor(0x00, 0x96, 0x88)
DARK_NAVY = RGBColor(0x1A, 0x23, 0x2F)

# Topic boundaries: (topic_name, insert_after_slide_0indexed, activity_steps)
TOPICS = [
    ("Understanding Customers & Customer Service", 16, [
        "Activity: Customer Service Role Play",
        "Scenario: A new customer walks into your store looking confused and unsure about what they need.",
        "Step 1: Greet the customer warmly and introduce yourself",
        "Step 2: Ask open-ended questions to understand their needs",
        "Step 3: Listen actively and paraphrase their requirements",
        "Step 4: Suggest appropriate products/services based on their needs",
        "Step 5: Check if the customer is satisfied with your suggestions",
        "Expected Output: A completed Customer Interaction Form",
        "Duration: 15 minutes",
    ]),
    ("Establishing Your Attitude", 23, [
        "Activity: Attitude Self-Assessment & Role Play",
        "Scenario: You have had a tough morning, but a customer needs your full attention and positive attitude.",
        "Step 1: Complete the Service Attitude Self-Assessment checklist",
        "Step 2: Identify your top 3 attitude strengths and 2 areas for improvement",
        "Step 3: Pair up and role-play: one is the customer, the other the service provider",
        "Step 4: Practice maintaining positive attitude even under pressure",
        "Step 5: Switch roles and repeat the exercise",
        "Expected Output: Completed self-assessment with personal action plan",
        "Duration: 20 minutes",
    ]),
    ("Identifying and Addressing Customer Needs", 30, [
        "Activity: Customer Needs Analysis Workshop",
        "Scenario: A long-time customer calls to complain that your service no longer meets their growing business needs.",
        "Step 1: Read the customer profile card provided",
        "Step 2: List the customer's spoken and unspoken needs",
        "Step 3: Map needs to the four levels: Basic, Expected, Desired, Unanticipated",
        "Step 4: Propose a tailored solution addressing each level",
        "Step 5: Present your solution to the class (2 min per group)",
        "Expected Output: Customer Needs Analysis template with proposed solutions",
        "Duration: 20 minutes",
    ]),
    ("In-Person Customer Service", 37, [
        "Activity: In-Person Service Simulation",
        "Scenario: A walk-in customer arrives at your service counter visibly frustrated about a delayed order.",
        "Step 1: Practice appropriate body language (eye contact, open posture, smile)",
        "Step 2: Greet the customer and acknowledge their frustration",
        "Step 3: Use active listening to understand the full issue",
        "Step 4: Provide a clear resolution timeline and follow-up plan",
        "Step 5: Close the interaction professionally and thank the customer",
        "Expected Output: Completed In-Person Service Evaluation rubric",
        "Duration: 15 minutes",
    ]),
    ("Giving Customer Service over the Phone", 44, [
        "Activity: Professional Phone Call Practice",
        "Scenario: A customer calls to enquire about a billing discrepancy. They are polite but concerned.",
        "Step 1: Answer using the 3-ring rule with proper phone greeting",
        "Step 2: Verify customer identity following company protocol",
        "Step 3: Listen to the concern and take notes",
        "Step 4: Explain the billing details clearly and resolve the issue",
        "Step 5: Summarise the resolution and confirm customer satisfaction",
        "Expected Output: Completed Phone Call Log Sheet",
        "Duration: 15 minutes",
    ]),
    ("Providing Customer Service via Email", 51, [
        "Activity: Professional Email Drafting Exercise",
        "Scenario: A customer emails complaining about receiving the wrong product and requests an immediate replacement.",
        "Step 1: Analyse the customer email for key concerns and tone",
        "Step 2: Draft a professional response using AIDA (Acknowledge, Inform, Do, Assure)",
        "Step 3: Apply proper netiquette - check tone, grammar, and formatting",
        "Step 4: Peer review - exchange emails with a partner and provide feedback",
        "Step 5: Revise and submit final email draft",
        "Expected Output: Professional customer service email (printed/shared)",
        "Duration: 20 minutes",
    ]),
    ("Generating Return Business", 57, [
        "Activity: Customer Retention Strategy Workshop",
        "Scenario: Your company's repeat customer rate has dropped 15% this quarter. Management wants a recovery plan.",
        "Step 1: Analyse the provided customer feedback data (satisfaction surveys)",
        "Step 2: Identify the top 3 reasons customers are not returning",
        "Step 3: Develop a Customer Retention Action Plan with specific strategies",
        "Step 4: Design one loyalty initiative suitable for Singapore's market",
        "Step 5: Present your plan to the class (3 min per group)",
        "Expected Output: Customer Retention Action Plan (1-page template)",
        "Duration: 20 minutes",
    ]),
    ("Service Recovery & Handling Difficult Customers", 63, [
        "Activity: HEARD De-Escalation Role Play",
        "Scenario: An angry customer threatens to post a negative review online because of repeated service failures.",
        "Step 1: Review the HEARD framework (Hear, Empathise, Apologise, Resolve, Diagnose)",
        "Step 2: Pair up - one person plays the angry customer, the other the service provider",
        "Step 3: Apply each step of HEARD to de-escalate the situation",
        "Step 4: Switch roles and handle a different difficult customer scenario",
        "Step 5: Debrief - discuss what worked and what could improve",
        "Expected Output: Completed De-Escalation Reflection Sheet",
        "Duration: 20 minutes",
    ]),
    ("Customer Service Escalation", 69, [
        "Activity: Escalation Decision-Making Exercise",
        "Scenario: You are handling a customer who has become verbally abusive and refuses to calm down.",
        "Step 1: Assess the situation using the Escalation Decision Framework",
        "Step 2: Identify which escalation triggers are present (vulgarity, threats, etc.)",
        "Step 3: Decide - can you handle it, or must you escalate? Justify your decision",
        "Step 4: If escalating, draft a brief handover note for your supervisor",
        "Step 5: Discuss as a class - share decisions and reasoning",
        "Expected Output: Completed Escalation Assessment Form",
        "Duration: 15 minutes",
    ]),
]

prs = Presentation(INPUT)
orig_count = len(prs.slides)
print(f"Original slides: {orig_count}")

# Find layout
layout = None
for sl in prs.slide_layouts:
    if "Title" in (sl.name or "") and "Content" in (sl.name or ""):
        layout = sl
        break
if not layout:
    layout = prs.slide_layouts[1]

# Add activity slides (appended at end first)
activity_positions = []
for topic_name, insert_after_idx, steps in TOPICS:
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ""
            tf = ph.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"Activity: {topic_name}"
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = DARK_NAVY
        elif ph.placeholder_format.idx == 1:
            ph.text = ""
            tf = ph.text_frame
            tf.word_wrap = True
            for si, step in enumerate(steps):
                p = tf.paragraphs[0] if si == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = step
                if step.startswith("Activity:") or step.startswith("Scenario:"):
                    run.font.bold = True
                    run.font.size = Pt(13)
                elif step.startswith("Expected Output:") or step.startswith("Duration:"):
                    run.font.bold = True
                    run.font.size = Pt(12)
                    run.font.color.rgb = ACCENT_TEAL
                else:
                    run.font.size = Pt(12)
                run.font.color.rgb = DARK_NAVY
                p.space_after = Pt(4)

    # Add teal accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(530000), Emu(9144000), Emu(40000)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_TEAL
    bar.line.fill.background()

    activity_positions.append(insert_after_idx)
    print(f"  Added activity for: {topic_name}")

# Reorder: move activity slides from end to correct positions
sldIdLst = prs.slides._sldIdLst
sldIds = list(sldIdLst)

activity_map = {}
for ai, pos in enumerate(activity_positions):
    activity_map[pos] = sldIds[orig_count + ai]

new_order = []
for i in range(orig_count):
    new_order.append(sldIds[i])
    if i in activity_map:
        new_order.append(activity_map[i])

# Rebuild slide order
for child in list(sldIdLst):
    sldIdLst.remove(child)
for item in new_order:
    sldIdLst.append(item)

print(f"\nFinal slides: {len(prs.slides)} ({len(TOPICS)} activities added)")
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
